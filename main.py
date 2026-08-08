import io
import json
import os
import time
from abc import ABC, abstractmethod
from concurrent.futures import ThreadPoolExecutor, as_completed
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

import msal
import requests

# Document & Spreadsheet Libraries
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

# Colab browser file download support
try:
    from google.colab import files as colab_files
except ImportError:
    colab_files = None


# ==========================================
# 1. CONFIGURATION & GEMINI SETUP
# ==========================================
CLIENT_ID = os.getenv("CLIENT_ID")
TENANT_ID = "common"

# Your Gemini API key from Google AI Studio (aistudio.google.com)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Gemini Model Target
MODEL_NAME = os.getenv("MODEL_NAME")

SCOPES = ["https://graph.microsoft.com/Files.Read"]

client = OpenAI(
    api_key=GEMINI_API_KEY,
    base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
)


# ==========================================
# 2. PLUGGABLE EXTRACTOR REGISTRY
# ==========================================
class BaseParser(ABC):
    @abstractmethod
    def parse(self, file_bytes: bytes) -> str:
        pass


class ParserRegistry:
    _parsers = {}

    @classmethod
    def register(cls, *extensions):
        def decorator(parser_cls):
            instance = parser_cls()
            for ext in extensions:
                cls._parsers[ext.lower()] = instance
            return parser_cls
        return decorator

    @classmethod
    def parse(cls, file_bytes: bytes, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        parser = cls._parsers.get(ext, DefaultTextParser())
        try:
            parsed_content = parser.parse(file_bytes)
            return parsed_content.strip() if parsed_content else "[Document is empty or contains no readable text]"
        except Exception as e:
            return f"[Error parsing {filename}: {str(e)}]"


# --- CONCRETE PARSER STRATEGIES ---

@ParserRegistry.register(".docx", ".doc")
class DocxParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        doc = Document(io.BytesIO(file_bytes))
        output = []

        for p in doc.paragraphs:
            if p.text.strip():
                output.append(p.text.strip())

        for table_idx, table in enumerate(doc.tables, start=1):
            output.append(f"\n[Table {table_idx}]")
            for row in table.rows:
                row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                output.append(" | ".join(row_cells))

        return "\n".join(output)


@ParserRegistry.register(".pptx", ".ppt")
class PptxParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines = [f"--- Slide {idx} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    slide_lines.append(shape.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_str = " | ".join([cell.text.strip().replace("\n", " ") for cell in row.cells])
                        slide_lines.append(row_str)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"[Speaker Notes]: {notes}")

            slides_text.append("\n".join(slide_lines))

        return "\n\n".join(slides_text)


@ParserRegistry.register(".xlsx", ".xls", ".csv")
class ExcelParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        buffer = io.BytesIO(file_bytes)
        output = []

        try:
            excel_file = pd.ExcelFile(buffer)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                if not df.empty:
                    md_table = df.to_markdown(index=False)
                    output.append(f"--- Sheet: {sheet_name} ---\n{md_table}")
        except Exception:
            buffer.seek(0)
            df = pd.read_csv(buffer)
            if not df.empty:
                output.append(df.to_markdown(index=False))

        return "\n\n".join(output) if output else "[Empty Sheet/Data]"


@ParserRegistry.register(".pdf")
class PdfParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        pdf_pages = []

        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_lines = [f"--- Page {page_num} ---"]

                text = page.extract_text()
                if text:
                    page_lines.append(text.strip())

                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        clean_row = [str(cell).strip().replace("\n", " ") if cell else "" for cell in row]
                        page_lines.append(" | ".join(clean_row))

                pdf_pages.append("\n".join(page_lines))

        return "\n\n".join(pdf_pages)


@ParserRegistry.register(".png", ".jpg", ".jpeg", ".tiff", ".bmp")
class ImageOcrParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        image = Image.open(io.BytesIO(file_bytes))
        extracted_text = pytesseract.image_to_string(image)
        return extracted_text.strip() if extracted_text.strip() else "[Image contains no readable text]"


class DefaultTextParser(BaseParser):
    def parse(self, file_bytes: bytes) -> str:
        return file_bytes.decode("utf-8", errors="ignore")


# ==========================================
# 3. ONEDRIVE MSAL AUTHENTICATION
# ==========================================
def authenticate_onedrive():
    authority = f"https://login.microsoftonline.com/{TENANT_ID}"
    app = msal.PublicClientApplication(CLIENT_ID, authority=authority)

    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result:
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    if "user_code" not in flow:
        raise Exception(f"Failed to create device flow: {flow.get('error_description')}")

    print("=" * 60)
    print(f"1. Open URL in browser: {flow['verification_uri']}")
    print(f"2. Enter this code:     {flow['user_code']}")
    print("=" * 60)

    result = app.acquire_token_by_device_flow(flow)
    if "access_token" in result:
        print("\n[SUCCESS] OneDrive Authentication Successful!\n")
        return result["access_token"]
    else:
        raise Exception(f"Authentication failed: {result.get('error_description')}")


# ==========================================
# 4. SEARCH & CRAWLER ENGINES
# ==========================================
FILE_TYPE_MAP = {
    "word": [".docx", ".doc"],
    "docx": [".docx"],
    "doc": [".doc"],
    "powerpoint": [".pptx", ".ppt"],
    "ppt": [".pptx", ".ppt"],
    "pptx": [".pptx"],
    "excel": [".xlsx", ".xls", ".csv"],
    "spreadsheet": [".xlsx", ".xls", ".csv"],
    "xls": [".xls"],
    "xlsx": [".xlsx"],
    "csv": [".csv"],
    "pdf": [".pdf"],
    "image": [".png", ".jpg", ".jpeg", ".tiff", ".bmp"],
    "text": [".txt"]
}


class OneDriveSearchCache:
    def __init__(self, ttl_seconds=180):
        self.cache = {}
        self.ttl = ttl_seconds

    def get(self, key):
        if key in self.cache:
            data, timestamp = self.cache[key]
            if time.time() - timestamp < self.ttl:
                return data
            else:
                del self.cache[key]
        return None

    def set(self, key, value):
        self.cache[key] = (value, time.time())


SEARCH_CACHE = OneDriveSearchCache(ttl_seconds=180)


def native_onedrive_search(access_token, keyword):
    if not keyword:
        return []

    cached_res = SEARCH_CACHE.get(keyword.lower())
    if cached_res is not None:
        return cached_res

    headers = {"Authorization": f"Bearer {access_token}"}
    url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{keyword}')"

    all_matches = []
    while url:
        res = requests.get(url, headers=headers)
        if res.status_code != 200:
            break
        data = res.json()
        items = data.get("value", [])

        for item in items:
            parent_path = item.get("parentReference", {}).get("path", "")
            clean_parent = parent_path.split("root:")[-1] if "root:" in parent_path else ""
            item["calculated_path"] = f"Root{clean_parent}/{item.get('name')}"
            all_matches.append(item)

        url = data.get("@odata.nextLink")

    SEARCH_CACHE.set(keyword.lower(), all_matches)
    return all_matches


def get_recent_files(access_token):
    headers = {"Authorization": f"Bearer {access_token}"}
    url = "https://graph.microsoft.com/v1.0/me/drive/recent"
    res = requests.get(url, headers=headers)

    if res.status_code == 200:
        items = res.json().get("value", [])
        for item in items:
            parent_path = item.get("parentReference", {}).get("path", "")
            clean_parent = parent_path.split("root:")[-1] if "root:" in parent_path else ""
            item["calculated_path"] = f"Root{clean_parent}/{item.get('name')}"
        return items
    return []


def get_folder_files(access_token, folder_name):
    headers = {"Authorization": f"Bearer {access_token}"}
    search_res = native_onedrive_search(access_token, folder_name)
    folders = [item for item in search_res if "folder" in item]

    if not folders:
        all_items = crawl_onedrive_tree_parallel(access_token)
        folders = [item for item in all_items if "folder" in item and folder_name.lower() in item.get("name", "").lower()]

    if not folders:
        return []

    target_folder = folders[0]
    folder_id = target_folder["id"]

    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{folder_id}/children"
    res = requests.get(url, headers=headers)
    if res.status_code == 200:
        items = res.json().get("value", [])
        folder_path = target_folder.get("calculated_path", folder_name)
        for item in items:
            item["calculated_path"] = f"{folder_path}/{item.get('name')}"
        return items
    return []


def parallel_fetch_folder_children(access_token, folder_item):
    headers = {"Authorization": f"Bearer {access_token}"}
    item_id = folder_item["id"]
    path = folder_item["calculated_path"]
    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/children"

    res = requests.get(url, headers=headers)
    if res.status_code != 200:
        return [], []

    items = res.json().get("value", [])
    subfolders = []
    files = []

    for item in items:
        item["calculated_path"] = f"{path}/{item['name']}"
        if "folder" in item:
            subfolders.append(item)
        else:
            files.append(item)

    return subfolders, files


def crawl_onedrive_tree_parallel(access_token, max_workers=8):
    root_folder = {"id": "root", "calculated_path": "Root"}
    queue = [root_folder]
    all_folders = []
    all_files = []

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        while queue:
            futures = {executor.submit(parallel_fetch_folder_children, access_token, folder): folder for folder in queue}
            queue = []

            for future in as_completed(futures):
                subfolders, files = future.result()
                all_folders.extend(subfolders)
                all_files.extend(files)
                queue.extend(subfolders)

    return all_folders + all_files


def list_all_folders(access_token):
    print("\n[FAST PARALLEL SCAN] Fetching folder directory structure...")
    start_time = time.perf_counter()

    all_items = crawl_onedrive_tree_parallel(access_token)
    folders = [item for item in all_items if "folder" in item]

    elapsed = time.perf_counter() - start_time
    print(f"[TIMING] Directory crawl completed in {elapsed:.3f} seconds.")

    if not folders:
        return "No subfolders found in OneDrive (only Root directory exists)."

    folder_list = [f"- {f['calculated_path']} (Created: {f.get('createdDateTime', 'N/A')})" for f in folders]
    return f"Found {len(folders)} folder(s) in OneDrive ({elapsed:.2f}s):\n" + "\n".join(folder_list)


# ==========================================
# 5. CORE ACTION FUNCTIONS
# ==========================================
def find_and_download_latest_file(access_token, keyword=None, file_type=None, folder_name=None):
    headers = {"Authorization": f"Bearer {access_token}"}

    # Clean generic terms passed as keywords
    if keyword:
        kw_clean = keyword.lower().strip()
        if kw_clean in FILE_TYPE_MAP:
            if not file_type:
                file_type = kw_clean
            keyword = None
        elif kw_clean in ["file", "latest", "document", "file contents", "latest file", "contents", "show me"]:
            keyword = None

    target_extensions = []
    if file_type:
        clean_type = file_type.lower().strip(".")
        target_extensions = FILE_TYPE_MAP.get(clean_type, [f".{clean_type}"])

    print(f"\n[SEARCH] Finding latest file (Keyword: '{keyword}', Type: '{file_type}', Folder: '{folder_name}')...")
    search_start = time.perf_counter()

    candidate_items = []

    if folder_name:
        candidate_items = get_folder_files(access_token, folder_name)
    elif keyword:
        candidate_items = native_onedrive_search(access_token, keyword)
    elif file_type and target_extensions:
        # --- FIX: Search Graph API natively by file extension first ---
        for ext in target_extensions:
            # Query Graph search for the extension (e.g. search(q='.xlsx'))
            ext_matches = native_onedrive_search(access_token, ext)
            candidate_items.extend(ext_matches)
        
        # Deduplicate results by file ID
        seen_ids = set()
        unique_candidates = []
        for item in candidate_items:
            if item["id"] not in seen_ids:
                seen_ids.add(item["id"])
                unique_candidates.append(item)
        candidate_items = unique_candidates
    else:
        candidate_items = get_recent_files(access_token)

    # Filter out folders and verify extensions
    filtered_items = []
    for item in candidate_items:
        if "folder" in item:
            continue

        file_name = item.get("name", "")
        if target_extensions:
            if any(file_name.lower().endswith(ext) for ext in target_extensions):
                filtered_items.append(item)
        else:
            filtered_items.append(item)

    # --- FALLBACK: If recent API or native search returned nothing, crawl the full drive ---
    if not filtered_items and not folder_name:
        print("[FALLBACK] Native search returned no items. Crawling OneDrive tree...")
        all_items = crawl_onedrive_tree_parallel(access_token)
        for item in all_items:
            if "folder" in item:
                continue
            file_name = item.get("name", "")
            if target_extensions:
                if any(file_name.lower().endswith(ext) for ext in target_extensions):
                    filtered_items.append(item)
            else:
                filtered_items.append(item)

    search_end = time.perf_counter()
    print(f"[TIMING] Search completed in {search_end - search_start:.3f} seconds.")

    if not filtered_items:
        folder_msg = f" in folder '{folder_name}'" if folder_name else ""
        type_msg = f" of type '{file_type}'" if file_type else ""
        kw_msg = f" matching '{keyword}'" if keyword else ""
        return f"No matching files found{kw_msg}{type_msg}{folder_msg} in OneDrive."

    # Sort descending by creation/modification date
    filtered_items.sort(
        key=lambda x: x.get("createdDateTime") or x.get("lastModifiedDateTime") or "",
        reverse=True
    )

    winning_file = filtered_items[0]
    file_id = winning_file["id"]
    file_name = winning_file["name"]
    created_time = winning_file.get("createdDateTime") or winning_file.get("lastModifiedDateTime", "Unknown")
    folder_path = winning_file.get("calculated_path", "Unknown Path")

    print(f"\n[FOUND MATCH] File Name: '{file_name}' | Created/Modified: {created_time} | Path: {folder_path}")

    download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
    dl_res = requests.get(download_url, headers=headers)

    if dl_res.status_code == 200:
        with open(file_name, "wb") as f:
            f.write(dl_res.content)
        print(f"[DOWNLOAD SUCCESS] Saved locally to: {os.path.abspath(file_name)}")

        if colab_files:
            try:
                colab_files.download(file_name)
            except Exception:
                pass

        extracted_content = ParserRegistry.parse(dl_res.content, file_name)

        return (
            f"=== FILE FOUND & DOWNLOADED ===\n"
            f"• File Name: {file_name}\n"
            f"• Created/Modified Date: {created_time}\n"
            f"• OneDrive Path: {folder_path}\n"
            f"• Local Path: {os.path.abspath(file_name)}\n\n"
            f"=== EXTRACTED FILE CONTENTS ===\n"
            f"{extracted_content}"
        )
    else:
        return f"Failed to download file '{file_name}': {dl_res.status_code} - {dl_res.text}"

def summarize_recent_documents(access_token, keyword, count=7):
    headers = {"Authorization": f"Bearer {access_token}"}
    print(f"\n[MULTI-DOC SEARCH] Querying OneDrive for recent '{keyword}' documents (Target: last {count})...")

    start_time = time.perf_counter()
    matching_files = [
        item for item in native_onedrive_search(access_token, keyword)
        if "folder" not in item
    ]
    print(f"[TIMING] Indexed search completed in {time.perf_counter() - start_time:.3f} seconds.")

    if not matching_files:
        return f"No documents found matching keyword '{keyword}' in OneDrive to summarize."

    matching_files.sort(
        key=lambda x: x.get("createdDateTime") or x.get("lastModifiedDateTime") or "",
        reverse=True
    )
    selected_files = matching_files[:count]
    print(f"[FOUND] Summarizing the newest {len(selected_files)} document(s)...")

    corpus = []
    for idx, item in enumerate(selected_files, start=1):
        file_id = item["id"]
        file_name = item["name"]
        created_date = item.get("createdDateTime", "Unknown Date")
        path = item.get("calculated_path", "Unknown Path")

        print(f" -> Reading [{idx}/{len(selected_files)}]: '{file_name}' ({created_date})")

        dl_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
        dl_res = requests.get(dl_url, headers=headers)

        if dl_res.status_code == 200:
            extracted_text = ParserRegistry.parse(dl_res.content, file_name)
            if extracted_text.strip():
                corpus.append(
                    f"=== DOCUMENT {idx}: {file_name} ===\n"
                    f"Created Date: {created_date}\n"
                    f"Path: {path}\n"
                    f"Content:\n{extracted_text}\n"
                )

    if not corpus:
        return f"Found documents matching '{keyword}', but could not extract readable text from them."

    full_text_context = "\n\n".join(corpus)[:40000]

    summary_prompt = [
        {
            "role": "system",
            "content": (
                "You are an executive assistant specializing in document synthesis. "
                "Analyze the provided set of documents chronologically. Produce a structured report with two sections:\n\n"
                "1. **Individual Document Gists**: Brief 2-3 sentence overview for each document in chronological order.\n"
                "2. **Integrated Executive Summary**: A rephrased overall summary connecting key themes, decisions made over time, and outstanding action items."
            )
        },
        {
            "role": "user",
            "content": f"Please summarize these recent {len(selected_files)} '{keyword}' documents:\n\n{full_text_context}"
        }
    ]

    print("[SYNTHESIZING] Sending document corpus to Gemini for multi-doc summary...")
    synth_res = client.chat.completions.create(
        model=MODEL_NAME,
        messages=summary_prompt
    )

    return synth_res.choices[0].message.content


# ==========================================
# 6. TOOL SCHEMAS & DISPATCHER
# ==========================================
tools = [
    {
        "type": "function",
        "function": {
            "name": "find_and_download_latest_file",
            "description": (
                "Searches, downloads, and reads the contents of the latest file matching keywords, file formats "
                "(excel, ppt, word, pdf, image), or folder names. ALWAYS use this tool whenever a user asks to find, "
                "download, show contents of, or open a file."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "keyword": {"type": "string", "description": "Topic or keyword in filename/content (e.g. 'kmm', 'budget'). Optional."},
                    "file_type": {"type": "string", "description": "File format or type (e.g. 'excel', 'ppt', 'word', 'pdf', 'csv'). Optional."},
                    "folder_name": {"type": "string", "description": "Specific folder name to look inside. Optional."}
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "summarize_recent_documents",
            "description": "Reads and synthesizes the last N created documents matching a keyword across OneDrive into individual gists and an integrated summary.",
            "parameters": {
                "type": "object",
                "properties": {
                    "keyword": {"type": "string", "description": "Target document topic or keyword."},
                    "count": {"type": "integer", "description": "Number of recent documents to include (default: 7)."}
                },
                "required": ["keyword"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_all_folders",
            "description": "Lists all directories and folders present across OneDrive. Use ONLY when user explicitly asks to list folders.",
            "parameters": {"type": "object", "properties": {}}
        }
    }
]

TOOL_DISPATCHER = {
    "find_and_download_latest_file": lambda token, args: find_and_download_latest_file(
        access_token=token,
        keyword=args.get("keyword"),
        file_type=args.get("file_type"),
        folder_name=args.get("folder_name")
    ),
    "summarize_recent_documents": lambda token, args: summarize_recent_documents(
        access_token=token,
        keyword=args.get("keyword"),
        count=args.get("count", 7)
    ),
    "list_all_folders": lambda token, args: list_all_folders(
        access_token=token
    )
}


# ==========================================
# 7. INTERACTIVE CHATBOT LOOP
# ==========================================
def run_chatbot():
    access_token = authenticate_onedrive()

    print("--- Gemini OneDrive Agent Activated! ---")
    print("Supported queries:")
    print(" - 'find me latest excel file'")
    print(" - 'show me contents of latest excel file'")
    print(" - 'find me latest kmm ppt'")
    print(" - 'find me latest file in Marketing folder'")
    print(" - 'Give summary of last 7 MOMs'")
    print(" - 'List all folders in my OneDrive'")
    print("Type 'exit' to quit.\n")

    messages = [
        {
            "role": "system",
            "content": (
                "You are an executive OneDrive assistant. When users ask to locate, search, open, or view file contents "
                "(e.g., 'find me latest excel file', 'show me contents of latest excel file', 'find me latest kmm ppt'), "
                "you MUST call `find_and_download_latest_file`. "
                "ALWAYS include the File Name, Created Date, OneDrive Path, and Full Extracted Contents in your final response."
            )
        }
    ]

    while True:
        user_input = input("\nYou: ").strip()
        if user_input.lower() in ["exit", "quit"]:
            print("Goodbye!")
            break

        if not user_input:
            continue

        messages.append({"role": "user", "content": user_input})

        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=messages,
            tools=tools,
            tool_choice="auto"
        )

        response_message = response.choices[0].message
        messages.append(response_message)

        if response_message.tool_calls:
            for tool_call in response_message.tool_calls:
                func_name = tool_call.function.name
                args = json.loads(tool_call.function.arguments or "{}")

                print(f"\n[Agent executing tool: {func_name}({args})...]")

                handler = TOOL_DISPATCHER.get(func_name)
                if handler:
                    result = handler(access_token, args)
                else:
                    result = f"Error: Tool handler for '{func_name}' is not registered."

                messages.append({
                    "role": "tool",
                    "tool_call_id": tool_call.id,
                    "content": result
                })

            final_response = client.chat.completions.create(
                model=MODEL_NAME,
                messages=messages
            )

            output_text = final_response.choices[0].message.content
            if output_text and output_text.strip():
                print(f"\nAgent:\n{output_text}")
                messages.append(final_response.choices[0].message)
            else:
                print(f"\nAgent:\n{result}")
        else:
            print(f"\nAgent: {response_message.content}")


if __name__ == "__main__":
    run_chatbot()