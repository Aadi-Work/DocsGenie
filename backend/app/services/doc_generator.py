from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Optional

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches, Pt

from app.config import get_settings
from app.models.schemas import OutputFormat, TemplateMeta
from app.services.catalog import CatalogService


class DocumentGenerator:
    def __init__(self, catalog: CatalogService):
        self.catalog = catalog
        settings = get_settings()
        self.output_dir = Path(settings.storage_path) / "generated"
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        template: TemplateMeta,
        filled: dict[str, str],
        output_format: Optional[OutputFormat] = None,
    ) -> tuple[str, Path]:
        fmt = output_format or template.output_format
        version = self.catalog.latest_version(template).version
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = template.id.replace("tmpl-", "")
        filename = f"{safe_name}_v{version}_{stamp}.{fmt.value}"
        path = self.output_dir / filename

        if fmt == OutputFormat.docx:
            self._write_docx(template, filled, version, path)
        elif fmt == OutputFormat.xlsx:
            self._write_xlsx(template, filled, version, path)
        else:
            self._write_pptx(template, filled, version, path)

        return filename, path

    def _write_docx(
        self,
        template: TemplateMeta,
        filled: dict[str, str],
        version: str,
        path: Path,
    ) -> None:
        doc = Document()
        doc.add_heading(template.name, level=0)
        doc.add_paragraph(f"YMSLI Template Hub · Approved version {version}")
        doc.add_paragraph(template.description)
        doc.add_heading("Document Context", level=1)
        for key, value in filled.items():
            p = doc.add_paragraph()
            p.add_run(f"{key}: ").bold = True
            p.add_run(value)
        doc.add_heading("Content", level=1)
        for section in template.content_outline:
            doc.add_heading(section, level=2)
            # Replace placeholders inside section guidance
            body = filled.get(section) or filled.get(section.replace(" ", "_"), "")
            if not body:
                body = (
                    f"This section covers {section.lower()}. "
                    f"Generated from the approved {template.name} template."
                )
            doc.add_paragraph(body)
        doc.add_paragraph("")
        doc.add_paragraph("— End of document —")
        doc.save(path)

    def _write_xlsx(
        self,
        template: TemplateMeta,
        filled: dict[str, str],
        version: str,
        path: Path,
    ) -> None:
        wb = Workbook()
        ws = wb.active
        ws.title = "Summary"
        ws["A1"] = template.name
        ws["A1"].font = Font(bold=True, size=14)
        ws["A2"] = f"Version {version}"
        ws["A3"] = template.description
        ws["A5"] = "Field"
        ws["B5"] = "Value"
        ws["A5"].font = Font(bold=True)
        ws["B5"].font = Font(bold=True)
        row = 6
        for key, value in filled.items():
            ws[f"A{row}"] = key
            ws[f"B{row}"] = value
            row += 1

        plan = wb.create_sheet("Outline")
        plan["A1"] = "Section"
        plan["B1"] = "Notes"
        plan["A1"].font = Font(bold=True)
        plan["B1"].font = Font(bold=True)
        for idx, section in enumerate(template.content_outline, start=2):
            plan[f"A{idx}"] = section
            plan[f"B{idx}"] = filled.get(section, "Populate during project kickoff")
        wb.save(path)

    def _write_pptx(
        self,
        template: TemplateMeta,
        filled: dict[str, str],
        version: str,
        path: Path,
    ) -> None:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = template.name
        title_slide.placeholders[1].text = (
            f"YMSLI · Version {version}\n{filled.get('Project Name', filled.get('Title', ''))}"
        )

        context = prs.slides.add_slide(prs.slide_layouts[1])
        context.shapes.title.text = "Context"
        body = context.shapes.placeholders[1].text_frame
        body.clear()
        for key, value in filled.items():
            p = body.add_paragraph()
            p.text = f"{key}: {value}"
            p.level = 0
            p.font.size = Pt(18)

        for section in template.content_outline:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = filled.get(section, f"Details for {section}")
        prs.save(path)
