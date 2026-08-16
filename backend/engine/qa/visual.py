"""
Visual QA (optional).

    Office file -> LibreOffice -> PDF -> page images -> checks

Two levels:
  * geometric - runs without any model: clipped text, overflowing columns,
    blank pages, page-count blow-ups. Cheap, deterministic, catches most of it.
  * vision - if a vision-capable LLM is configured, the page image is described
    and asked about layout breakage.

Everything here degrades gracefully: no LibreOffice, no PDF tooling or no
model means a report with `skipped` set, never a crash.
"""

from __future__ import annotations

import base64
import glob
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

from .structural import QAIssue, QAReport

SOFFICE_CANDIDATES = ["soffice", "libreoffice",
                      "/usr/bin/soffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"]


def _soffice() -> Optional[str]:
    for c in SOFFICE_CANDIDATES:
        p = shutil.which(c) if not os.path.isabs(c) else (c if os.path.exists(c) else None)
        if p:
            return p
    return None


def to_pdf(path: str, out_dir: Optional[str] = None, timeout: int = 180) -> Optional[str]:
    """Convert any Office file to PDF. Returns None if LibreOffice is absent."""
    exe = _soffice()
    if exe is None:
        return None
    out_dir = out_dir or tempfile.mkdtemp(prefix="ote_pdf_")
    os.makedirs(out_dir, exist_ok=True)
    profile = tempfile.mkdtemp(prefix="ote_lo_profile_")
    cmd = [exe, "--headless", "--norestore", f"-env:UserInstallation=file://{profile}",
           "--convert-to", "pdf", "--outdir", out_dir, path]
    try:
        subprocess.run(cmd, capture_output=True, timeout=timeout, check=False)
    except subprocess.TimeoutExpired:
        return None
    pdf = os.path.join(out_dir, os.path.splitext(os.path.basename(path))[0] + ".pdf")
    return pdf if os.path.exists(pdf) else None


def pdf_to_images(pdf_path: str, out_dir: Optional[str] = None, dpi: int = 110) -> List[str]:
    out_dir = out_dir or tempfile.mkdtemp(prefix="ote_png_")
    os.makedirs(out_dir, exist_ok=True)
    prefix = os.path.join(out_dir, "page")
    try:
        import fitz                                     # PyMuPDF, if present
        paths = []
        with fitz.open(pdf_path) as doc:
            for i, page in enumerate(doc, 1):
                p = f"{prefix}-{i}.png"
                page.get_pixmap(dpi=dpi).save(p)
                paths.append(p)
        return paths
    except Exception:
        pass
    if shutil.which("pdftoppm"):
        subprocess.run(["pdftoppm", "-png", "-r", str(dpi), pdf_path, prefix],
                       capture_output=True, timeout=180, check=False)
        return sorted(glob.glob(prefix + "*.png"))
    return []


# --------------------------------------------------------------------------
def geometric_qa(output_path: str) -> QAReport:
    """
    Model-free layout checks. For xlsx this is where truncated text actually
    gets caught: a value wider than its column, with no wrap and no empty
    neighbour, will render clipped no matter how correct the value is.
    """
    rep = QAReport()
    ext = os.path.splitext(output_path)[1].lower().lstrip(".")

    if ext in ("xlsx", "xlsm"):
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        wb = load_workbook(output_path, data_only=False)
        for ws in wb.worksheets:
            merged = {}
            for rng in ws.merged_cells.ranges:
                merged[(rng.min_row, rng.min_col)] = rng
            for row in ws.iter_rows():
                for cell in row:
                    if not isinstance(cell.value, str) or not cell.value.strip():
                        continue
                    text = cell.value
                    letter = get_column_letter(cell.column)
                    width = (ws.column_dimensions[letter].width
                             if letter in ws.column_dimensions and ws.column_dimensions[letter].width
                             else 8.43)
                    rng = merged.get((cell.row, cell.column))
                    if rng is not None:
                        width = sum((ws.column_dimensions[get_column_letter(c)].width
                                     if get_column_letter(c) in ws.column_dimensions
                                     and ws.column_dimensions[get_column_letter(c)].width else 8.43)
                                    for c in range(rng.min_col, rng.max_col + 1))
                    wraps = bool(cell.alignment and cell.alignment.wrap_text)
                    neighbour = ws.cell(row=cell.row, column=cell.column + 1).value
                    if len(text) > width * 1.15 and not wraps and neighbour not in (None, ""):
                        rep.add("overflow",
                                f"text ({len(text)} chars) wider than column ({width:.0f}) "
                                f"with a filled neighbour - will render clipped",
                                "warning", where=f"{ws.title}!{cell.coordinate}")
                    if wraps and rng is not None and rng.min_row == rng.max_row:
                        est_lines = max(1, len(text) / max(1.0, width))
                        height = ws.row_dimensions[cell.row].height or 15
                        if est_lines * 13 > height * 1.3:
                            rep.add("row_height",
                                    f"wrapped text needs ~{est_lines:.0f} lines but row height "
                                    f"is {height:.0f}pt", "warning",
                                    where=f"{ws.title}!{cell.coordinate}")
        wb.close()

    elif ext in ("pptx", "potx"):
        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation(output_path)
        sw, sh = prs.slide_width, prs.slide_height
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # A shape - table or otherwise - that extends past the slide
                # edge will be clipped regardless of what its own cells say.
                # A table's declared column widths can also exceed the width
                # its graphicFrame was created with (python-pptx does not
                # rebalance the other columns when one is resized), so this
                # check catches authoring problems that per-cell text length
                # never would.
                if shape.left is not None and shape.width is not None:
                    right = shape.left + shape.width
                    if right > sw * 1.005:
                        over_in = (right - sw) / Emu(1) / 914400
                        rep.add("off_slide",
                                f"'{shape.name}' extends {over_in:.2f}in past the right edge "
                                f"of the slide", where=f"slide {i}")
                if shape.top is not None and shape.height is not None:
                    bottom = shape.top + shape.height
                    if bottom > sh * 1.005:
                        over_in = (bottom - sh) / Emu(1) / 914400
                        rep.add("off_slide",
                                f"'{shape.name}' extends {over_in:.2f}in past the bottom edge "
                                f"of the slide", where=f"slide {i}")
                if getattr(shape, "has_table", False) and shape.has_table:
                    total_w = sum(c.width or 0 for c in shape.table.columns)
                    if shape.width and total_w > shape.width * 1.02:
                        over_in = (total_w - shape.width) / Emu(1) / 914400
                        rep.add("off_slide",
                                f"table '{shape.name}' column widths sum to {over_in:.2f}in "
                                f"more than the table's own declared width - it will render "
                                f"wider than intended", "warning", where=f"slide {i}")
                    _pptx_table_overflow(shape, i, rep)
                    continue
                if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                    continue
                text = shape.text_frame.text or ""
                if not text.strip() or shape.width is None or shape.height is None:
                    continue
                size_pt = 18.0
                try:
                    r = shape.text_frame.paragraphs[0].runs
                    if r and r[0].font.size:
                        size_pt = r[0].font.size.pt
                except Exception:
                    pass
                chars_per_line = max(1, (shape.width / Emu(1) / 914400 * 96) / (size_pt * 0.55))
                lines = sum(max(1, len(l) / chars_per_line) for l in text.split("\n"))
                capacity = (shape.height / Emu(1) / 914400 * 72) / (size_pt * 1.25)
                if lines > capacity * 1.1:
                    rep.add("overflow",
                            f"~{lines:.0f} lines of text in a box that fits ~{capacity:.0f}",
                            "warning", where=f"slide {i} / {shape.name}")
        rep.stats["slides"] = len(prs.slides)

    return rep


def _pptx_table_overflow(shape, slide_no: int, rep: QAReport) -> None:
    """
    A table cell has no wrap-to-fit like a text box does - text that doesn't
    fit is simply clipped at the cell boundary. Estimated the same way as a
    free-text shape, but per column width instead of per shape width.
    """
    from pptx.util import Emu
    table = shape.table
    col_widths = [c.width for c in table.columns]
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            text = (cell.text or "").strip()
            if not text or ci >= len(col_widths) or col_widths[ci] is None:
                continue
            size_pt = 14.0
            try:
                for para in cell.text_frame.paragraphs:
                    if para.runs and para.runs[0].font.size:
                        size_pt = para.runs[0].font.size.pt
                        break
            except Exception:
                pass
            width_in = col_widths[ci] / Emu(1) / 914400
            margin_in = 0.2       # pptx default cell inset (~0.1" each side)
            chars_per_line = max(1, (max(width_in - margin_in, 0.1) * 96) / (size_pt * 0.55))
            if len(text) > chars_per_line and "\n" not in text and len(text.split()) <= 2:
                # a short, unbroken token (like a status word) that is wider
                # than the column has nowhere to wrap to - it will clip
                rep.add("overflow",
                        f"cell r{ri}c{ci} text {text!r} ({len(text)} chars) wider than its "
                        f"column (~{chars_per_line:.0f} chars) - will render clipped",
                        "warning", where=f"slide {slide_no} / {shape.name}")


def visual_qa(output_path: str, llm=None, max_pages: int = 5,
              keep_artifacts: bool = False) -> QAReport:
    """Geometric checks + optional vision pass over rendered pages."""
    rep = geometric_qa(output_path)

    pdf = to_pdf(output_path)
    if pdf is None:
        rep.add("render", "LibreOffice not available - skipped page rendering", "info")
        rep.stats["rendered"] = False
        return rep
    rep.stats["rendered"] = True
    rep.stats["pdf"] = pdf if keep_artifacts else os.path.basename(pdf)

    images = pdf_to_images(pdf)[:max_pages]
    rep.stats["pages"] = len(images)
    if not images:
        rep.add("render", "no page rasteriser available (install pymupdf or poppler-utils)", "info")
        return rep

    for i, img in enumerate(images, 1):
        if os.path.getsize(img) < 3000:
            rep.add("blank_page", f"page {i} looks blank", "warning", where=f"page {i}")

    if llm is not None and getattr(llm, "available", lambda: False)():
        _vision_pass(images, llm, rep)
    return rep


def _vision_pass(images: List[str], llm, rep: QAReport) -> None:
    """Only Anthropic/OpenAI-style multimodal calls; anything else is skipped."""
    import json
    system = ("You are inspecting a rendered page of a generated business document. "
              "Report ONLY layout defects you can see: clipped or truncated text, "
              "text overflowing its cell or box, broken table borders, misaligned "
              "columns, overlapping elements, obviously missing content. "
              'Return JSON: {"issues":[{"type":"...","detail":"...","severity":"error|warning"}]}')
    for i, img in enumerate(images, 1):
        try:
            with open(img, "rb") as fh:
                b64 = base64.b64encode(fh.read()).decode()
            payload = getattr(llm, "complete_vision_json", None)
            if payload is None:
                rep.add("vision", f"{llm.name} has no vision entrypoint - skipped", "info")
                return
            data = payload(system, "Inspect this page.", b64) or {}
            for issue in data.get("issues", []):
                rep.add("vision", f"{issue.get('type')}: {issue.get('detail')}",
                        issue.get("severity", "warning"), where=f"page {i}")
        except Exception as exc:
            rep.add("vision", f"page {i}: vision check failed ({exc})", "info")
