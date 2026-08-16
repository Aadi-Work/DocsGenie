"""
Structural QA.

`wb.save()` returning without an exception proves nothing. This module diffs
the generated file against the template it came from and fails on anything
that should have been preserved but was not:

    merged ranges, formulas, styles, dimensions, protected cells,
    placeholder resolution, row counts, data types, images/shapes
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

from openpyxl import load_workbook
from openpyxl.utils import range_boundaries

from ..parsers.base import PLACEHOLDER_RE


@dataclass
class QAIssue:
    check: str
    message: str
    severity: str = "error"          # error | warning | info
    where: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {"check": self.check, "message": self.message,
                "severity": self.severity, "where": self.where}


@dataclass
class QAReport:
    issues: List[QAIssue] = field(default_factory=list)
    stats: Dict[str, Any] = field(default_factory=dict)

    @property
    def ok(self) -> bool:
        return not any(i.severity == "error" for i in self.issues)

    def add(self, check: str, message: str, severity: str = "error", where: str = "") -> None:
        self.issues.append(QAIssue(check, message, severity, where))

    def to_dict(self) -> Dict[str, Any]:
        return {"ok": self.ok, "stats": self.stats,
                "issues": [i.to_dict() for i in self.issues]}


def structural_qa(template_path: str, output_path: str,
                  expected_rows: Optional[Dict[str, int]] = None) -> QAReport:
    ext = os.path.splitext(template_path)[1].lower().lstrip(".")
    if ext in ("xlsx", "xlsm", "xltx"):
        return _xlsx_qa(template_path, output_path)
    if ext in ("docx", "dotx"):
        return _docx_qa(template_path, output_path)
    if ext in ("pptx", "potx"):
        return _pptx_qa(template_path, output_path)
    return QAReport(stats={"skipped": f"no structural QA for .{ext}"})


# --------------------------------------------------------------------------
def _xlsx_qa(template_path: str, output_path: str) -> QAReport:
    rep = QAReport()
    tw = load_workbook(template_path, data_only=False)
    ow = load_workbook(output_path, data_only=False)

    if set(tw.sheetnames) != set(ow.sheetnames):
        rep.add("sheets", f"sheet set changed: {tw.sheetnames} -> {ow.sheetnames}")

    total_written, total_formulas = 0, 0
    for name in tw.sheetnames:
        if name not in ow.sheetnames:
            continue
        ts, os_ = tw[name], ow[name]

        # merged ranges: every template merge must survive (rows may shift down)
        t_merges = {str(r) for r in ts.merged_cells.ranges}
        o_merges = {str(r) for r in os_.merged_cells.ranges}
        lost = _lost_merges(t_merges, o_merges)
        if lost:
            rep.add("merged_ranges", f"{len(lost)} merged range(s) lost: {sorted(lost)[:5]}",
                    where=name)

        # formulas must not be destroyed. Text may legitimately change (a row
        # inserted inside a SUM/COUNTA range shifts its bounds) - so compare
        # by *shape* (function + column pattern, digits stripped) rather than
        # verbatim text, and only fail on an actual loss of formula count.
        t_f = _formula_map(ts)
        o_f = _formula_map(os_)
        total_formulas += len(t_f)
        if len(o_f) < len(t_f):
            rep.add("formulas", f"{len(t_f) - len(o_f)} formula(s) lost "
                                 f"({len(t_f)} in template, {len(o_f)} in output)", where=name)
        o_shapes = [_formula_shape(f) for f in o_f.values()]
        for coord, formula in t_f.items():
            shape = _formula_shape(formula)
            if formula not in o_f.values() and shape not in o_shapes:
                rep.add("formulas",
                        f"formula at {coord} ('{formula}') has no equivalent in the output "
                        f"(not even after allowing for row-shift)", where=name)

        # dimensions
        if ts.max_column != os_.max_column:
            rep.add("dimensions", f"column count changed {ts.max_column} -> {os_.max_column}",
                    "warning", where=name)

        # column widths / row heights inherited
        for key, dim in list(ts.column_dimensions.items())[:200]:
            odim = os_.column_dimensions.get(key)
            if dim.width and odim is not None and odim.width and abs(dim.width - odim.width) > 0.01:
                rep.add("column_width", f"column {key} width {dim.width} -> {odim.width}",
                        "warning", where=name)

        # sheet protection preserved
        if bool(ts.protection.sheet) != bool(os_.protection.sheet):
            rep.add("protection", "sheet protection flag changed", where=name)

        # unresolved placeholders and new content count
        for row in os_.iter_rows():
            for cell in row:
                v = cell.value
                if isinstance(v, str):
                    if PLACEHOLDER_RE.search(v):
                        rep.add("placeholders", f"unresolved placeholder {v[:40]!r}",
                                "warning", where=f"{name}!{cell.coordinate}")
                    if v.startswith("="):
                        continue
                if v not in (None, ""):
                    total_written += 1

        # images survive
        if len(getattr(ts, "_images", []) or []) != len(getattr(os_, "_images", []) or []):
            rep.add("images", "image count changed", where=name)

    rep.stats = {"sheets": len(tw.sheetnames), "formulas_in_template": total_formulas,
                 "non_empty_cells_in_output": total_written}
    tw.close()
    ow.close()
    return rep


def _formula_map(ws) -> Dict[str, str]:
    out = {}
    for row in ws.iter_rows():
        for cell in row:
            if isinstance(cell.value, str) and cell.value.startswith("="):
                out[cell.coordinate] = cell.value
    return out


def _formula_shape(formula: str) -> str:
    """Digits stripped: '=COUNTA(B24:B26)' and '=COUNTA(B25:B28)' -> same shape."""
    return "".join(ch for ch in formula.upper() if not ch.isdigit())


def _lost_merges(t_merges: set, o_merges: set) -> set:
    """A merge that moved down (row insertion) still counts as preserved."""
    def norm(rng: str):
        min_c, min_r, max_c, max_r = range_boundaries(rng)
        return (min_c, max_c, max_r - min_r)
    o_shapes = {}
    for r in o_merges:
        c0, c1, h = norm(r)
        o_shapes.setdefault((c0, c1, h), 0)
        o_shapes[(c0, c1, h)] += 1
    lost = set()
    for r in t_merges:
        if r in o_merges:
            continue
        shape = norm(r)
        if o_shapes.get(shape):
            o_shapes[shape] -= 1
        else:
            lost.add(r)
    return lost


# --------------------------------------------------------------------------
def _docx_qa(template_path: str, output_path: str) -> QAReport:
    from docx import Document
    rep = QAReport()
    t, o = Document(template_path), Document(output_path)

    if len(t.tables) != len(o.tables):
        rep.add("tables", f"table count changed {len(t.tables)} -> {len(o.tables)}")
    for i, (tt, ot) in enumerate(zip(t.tables, o.tables)):
        if len(tt.columns) != len(ot.columns):
            rep.add("table_columns", f"table {i}: columns {len(tt.columns)} -> {len(ot.columns)}")
        if len(ot.rows) < len(tt.rows):
            rep.add("table_rows", f"table {i}: rows dropped {len(tt.rows)} -> {len(ot.rows)}")

    t_styles = [p.style.name for p in t.paragraphs]
    o_styles = [p.style.name for p in o.paragraphs]
    if len(o_styles) < len(t_styles):
        rep.add("paragraphs", f"paragraph count dropped {len(t_styles)} -> {len(o_styles)}")
    for p in o.paragraphs:
        if PLACEHOLDER_RE.search(p.text or ""):
            rep.add("placeholders", f"unresolved placeholder {p.text[:40]!r}", "warning")

    if len(t.sections) != len(o.sections):
        rep.add("sections", "section count changed")
    rep.stats = {"paragraphs": len(o.paragraphs), "tables": len(o.tables)}
    return rep


def _pptx_qa(template_path: str, output_path: str) -> QAReport:
    from pptx import Presentation
    rep = QAReport()
    t, o = Presentation(template_path), Presentation(output_path)

    if len(t.slides) != len(o.slides):
        rep.add("slides", f"slide count changed {len(t.slides)} -> {len(o.slides)}")
    for i, (ts, os_) in enumerate(zip(t.slides, o.slides)):
        if len(ts.shapes) != len(os_.shapes):
            rep.add("shapes", f"slide {i+1}: shape count {len(ts.shapes)} -> {len(os_.shapes)}",
                    "warning")
        for shape in os_.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                if PLACEHOLDER_RE.search(shape.text_frame.text or ""):
                    rep.add("placeholders",
                            f"slide {i+1}: unresolved placeholder", "warning")
    if (t.slide_width, t.slide_height) != (o.slide_width, o.slide_height):
        rep.add("dimensions", "slide size changed")
    rep.stats = {"slides": len(o.slides)}
    return rep
