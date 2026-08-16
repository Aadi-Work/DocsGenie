"""Source readers: PDF / DOCX / PPTX / XLSX / TXT / MD / CSV / JSON / VTT -> text."""

from __future__ import annotations

import csv
import io
import json
import os
import re
from typing import List

VTT_TS_RE = re.compile(r"^\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->")


def read_source(path: str) -> str:
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    fn = {
        "txt": _txt, "md": _txt, "log": _txt, "eml": _txt,
        "json": _json, "csv": _csv, "tsv": _csv,
        "docx": _docx, "pptx": _pptx, "xlsx": _xlsx, "xlsm": _xlsx,
        "pdf": _pdf, "vtt": _vtt, "srt": _vtt,
    }.get(ext)
    if fn is None:
        return _txt(path)
    return fn(path)


def _txt(path: str) -> str:
    with open(path, "rb") as fh:
        raw = fh.read()
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _json(path: str) -> str:
    with open(path, encoding="utf-8") as fh:
        data = json.load(fh)

    def flatten(obj, prefix="") -> List[str]:
        out = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                out += flatten(v, f"{prefix}{k}: " if not isinstance(v, (dict, list)) else f"{prefix}{k}.")
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                out += flatten(v, prefix)
        else:
            out.append(f"{prefix}{obj}")
        return out

    return "\n".join(flatten(data))


def _csv(path: str) -> str:
    delim = "\t" if path.lower().endswith(".tsv") else ","
    lines = []
    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        for row in csv.reader(fh, delimiter=delim):
            lines.append(" | ".join(str(c).strip() for c in row))
    return "\n".join(lines)


def _docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if (p.text or "").strip()]
    for t in doc.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                if (shape.text_frame.text or "").strip():
                    parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("Notes: " + slide.notes_slide.notes_text_frame.text.strip())
    return "\n".join(parts)


def _xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _vtt(path: str) -> str:
    out = []
    for line in _txt(path).splitlines():
        s = line.strip()
        if not s or s == "WEBVTT" or s.isdigit() or VTT_TS_RE.match(s):
            continue
        out.append(s)
    return "\n".join(out)


def _pdf(path: str) -> str:
    """PyMuPDF -> pdfplumber -> pypdf -> pdftotext, whichever is installed."""
    try:
        import fitz                                   # PyMuPDF
        with fitz.open(path) as doc:
            return "\n".join(page.get_text() for page in doc)
    except Exception:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            out = []
            for page in pdf.pages:
                out.append(page.extract_text() or "")
                for tbl in page.extract_tables() or []:
                    for row in tbl:
                        out.append(" | ".join((c or "").strip() for c in row))
            return "\n".join(out)
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        return "\n".join((p.extract_text() or "") for p in PdfReader(path).pages)
    except Exception:
        pass
    import subprocess
    try:
        return subprocess.run(["pdftotext", "-layout", path, "-"],
                              capture_output=True, text=True, timeout=120).stdout
    except Exception as exc:
        raise RuntimeError(
            f"Cannot read PDF {path!r}: install pymupdf, pdfplumber or pypdf ({exc})")
