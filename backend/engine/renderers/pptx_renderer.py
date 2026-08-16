"""
PPTX renderer.

Slides are unforgiving about overflow, so this renderer writes into the
existing text frame and run (inheriting the theme font and size) and reports
row overflow rather than silently pushing content off the slide. Visual QA
catches whatever still overflows.
"""

from __future__ import annotations

import copy as _copy
import re
from typing import Any, Dict, List, Optional

from pptx import Presentation

from ..ir import DocType
from ..logging_config import get_logger
from ..mapper import FillInstruction, FillPlan
from ..normalize import to_display
from .base import BaseRenderer, RenderResult, WriteRecord, clear_if_unresolved_placeholder, register

log = get_logger("renderers.pptx")

PLACEHOLDER_RE = re.compile(r"\{\{\s*([\w.\-]+)\s*\}\}|\$\{\s*([\w.\-]+)\s*\}|<<\s*([\w.\- ]+)\s*>>")


class PptxRenderer(BaseRenderer):
    doc_type = DocType.PPTX

    def render(self, template_path: str, output_path: str, plan: FillPlan,
              clear_unresolved: bool = True) -> RenderResult:
        out = self.prepare_output(template_path, output_path)
        prs = Presentation(out)
        result = RenderResult(output_path=out)

        for ins in plan.writable():
            if ins.kind == "field":
                result.records.append(self._write_field(prs, ins))
            else:
                rec, added = self._write_table(prs, ins)
                result.records.append(rec)
                result.rows_added += added

        if clear_unresolved:
            cleared = self._clear_unresolved_placeholders(prs)
            if cleared:
                log.info("Cleared %d unresolved placeholder(s) that had no fillable value", cleared)

        prs.save(out)
        return result

    def _clear_unresolved_placeholders(self, prs) -> int:
        n = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    n += self._clear_text_frame(shape.text_frame)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            n += self._clear_text_frame(cell.text_frame)
        return n

    def _clear_text_frame(self, tf) -> int:
        cleared_any = 0
        for para in tf.paragraphs:
            full = "".join(r.text for r in para.runs) or para.text
            cleared = clear_if_unresolved_placeholder(full)
            if cleared != full and para.runs:
                para.runs[0].text = cleared or ""
                for r in para.runs[1:]:
                    r.text = ""
                cleared_any += 1
        return cleared_any

    # ------------------------------------------------------------------
    def _shape(self, prs, slide_index: int, shape_id: int):
        slide = prs.slides[int(slide_index)]
        for shape in slide.shapes:
            if shape.shape_id == int(shape_id):
                return shape
        raise KeyError(f"shape {shape_id} not found on slide {slide_index}")

    def _write_field(self, prs, ins: FillInstruction) -> WriteRecord:
        t = ins.target
        rec = WriteRecord(ins.node_id, ins.role,
                          f"slide[{t.get('slide_index')}]/shape[{t.get('shape_id')}]", ins.value)
        try:
            shape = self._shape(prs, t["slide_index"], t["shape_id"])
            text = to_display(ins.value)
            mode = t.get("mode", "replace")
            if mode == "after_colon":
                label = t.get("label_text", "")
                self._set_text(shape, f"{label}: {text}")
            elif shape.text_frame.text and PLACEHOLDER_RE.search(shape.text_frame.text):
                self._set_text(shape, PLACEHOLDER_RE.sub(text, shape.text_frame.text, count=1))
            else:
                self._set_text(shape, text)
            rec.rows_written = 1
        except Exception as exc:
            rec.status = "failed"
            rec.message = f"{type(exc).__name__}: {exc}"
        return rec

    def _set_text(self, shape, text: str) -> None:
        """Write into run 0 so the template's font survives; extra lines clone it."""
        tf = shape.text_frame
        lines = str(text).split("\n")
        para0 = tf.paragraphs[0]
        if para0.runs:
            para0.runs[0].text = lines[0]
            for r in para0.runs[1:]:
                r.text = ""
        else:
            para0.add_run().text = lines[0]

        for extra in tf.paragraphs[1:]:
            extra._p.getparent().remove(extra._p)
        for line in lines[1:]:
            new_p = _copy.deepcopy(para0._p)
            para0._p.getparent().append(new_p)
            from pptx.text.text import _Paragraph
            p = _Paragraph(new_p, tf)
            if p.runs:
                p.runs[0].text = line
                for r in p.runs[1:]:
                    r.text = ""

    # ------------------------------------------------------------------
    def _write_table(self, prs, ins: FillInstruction):
        t = ins.target
        rec = WriteRecord(ins.node_id, ins.role,
                          f"slide[{t.get('slide_index')}]/table[{t.get('shape_id')}]")
        added = 0
        try:
            shape = self._shape(prs, t["slide_index"], t["shape_id"])
            table = shape.table
            rows = ins.rows or []
            col_map = {f: int(h.get("col")) for f, h in (t.get("columns") or {}).items()
                       if h.get("col") is not None}
            template_row_idx = int(t.get("template_row", 1))

            for i, row_data in enumerate(rows):
                idx = template_row_idx + i
                if idx >= len(table.rows):
                    self._clone_row(table, template_row_idx)
                    added += 1
                row = table.rows[idx]
                for fname, value in row_data.items():
                    ci = col_map.get(fname)
                    if ci is None or ci >= len(table.columns):
                        continue
                    cell = row.cells[ci]
                    self._set_text(cell, to_display(value))
            rec.rows_written = len(rows)
            if added:
                rec.message = f"cloned {added} row(s) from the template row"
        except Exception as exc:
            rec.status = "failed"
            rec.message = f"{type(exc).__name__}: {exc}"
        return rec, added

    def _clone_row(self, table, template_row_idx: int) -> None:
        tr = table.rows[template_row_idx]._tr
        new_tr = _copy.deepcopy(tr)
        tr.getparent().append(new_tr)
        for cell in table.rows[len(table.rows) - 1].cells:
            for para in cell.text_frame.paragraphs:
                for r in para.runs:
                    r.text = ""


register("pptx", PptxRenderer)
register("potx", PptxRenderer)
