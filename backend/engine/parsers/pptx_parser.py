"""
PPTX -> Template IR.

    Presentation > Slide(section) > {Shape/TextBox label+value, Table, Image}

Fillable shapes on a slide:
  * an empty layout placeholder (BODY/SUBTITLE/etc.) - labelled by the layout
  * a text box containing "Label:" with nothing after it
  * a text box containing {{placeholder}}
  * a table with a header row -> repeating collection
"""

from __future__ import annotations

import re
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Pt

from ..ir import (ColumnIR, DocType, Location, Node, NodeType, Section,
                  StyleSignals, TableIR, TemplateIR, ValueFormat)
from .base import (BLANK_LINE_RE, TemplateParser, detect_value_format,
                   find_placeholder, looks_like_label, register)

LABEL_VALUE_RE = re.compile(r"^\s*([^:\n]{1,60}?)\s*:\s*(.*)$", re.S)
PROMPT_TEXT_RE = re.compile(r"click to (add|edit)", re.I)


class PptxParser(TemplateParser):
    doc_type = DocType.PPTX

    def parse(self, path: str) -> TemplateIR:
        prs = Presentation(path)
        ir = TemplateIR(doc_type=DocType.PPTX, source_path=path,
                        meta={"slides": len(prs.slides),
                              "slide_width": prs.slide_width, "slide_height": prs.slide_height})

        for si, slide in enumerate(prs.slides):
            title = self._slide_title(slide) or f"Slide {si + 1}"
            sec = Section(section_id=f"slide::{si}", title=title,
                          location=Location(DocType.PPTX, {"slide_index": si}))
            ir.sections.append(sec)

            for shape in slide.shapes:
                self._parse_shape(ir, sec, si, shape)
        return ir

    # ------------------------------------------------------------------
    def _parse_shape(self, ir: TemplateIR, sec: Section, si: int, shape) -> None:
        parts = {"slide_index": si, "shape_id": shape.shape_id,
                 "shape_name": shape.name}
        if shape.is_placeholder:
            parts["placeholder_idx"] = shape.placeholder_format.idx
            parts["placeholder_type"] = str(shape.placeholder_format.type)

        if getattr(shape, "has_table", False) and shape.has_table:
            self._parse_table(ir, sec, si, shape, parts)
            return

        if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
            n = Node(node_id=f"slide{si}::shape{shape.shape_id}", type=NodeType.IMAGE,
                     location=Location(DocType.PPTX, parts), editable=False,
                     is_empty=False, section=sec.title)
            ir.nodes.append(n)
            sec.node_ids.append(n.node_id)
            return

        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            return

        text = (shape.text_frame.text or "").strip()
        style = self._shape_style(shape)
        is_prompt = bool(PROMPT_TEXT_RE.search(text))
        node_id = f"slide{si}::shape{shape.shape_id}"

        ph = find_placeholder(text)
        if ph:
            n = Node(node_id=node_id, type=NodeType.VALUE_REGION,
                     location=Location(DocType.PPTX, {**parts, "mode": "replace"}),
                     text=text, label=ph.replace("_", " "), placeholder=ph,
                     section=sec.title, style=style, is_empty=False)
            n.meta["multiline_capable"] = True
            ir.nodes.append(n)
            sec.node_ids.append(n.node_id)
            return

        # "Label:" with nothing after it
        m = LABEL_VALUE_RE.match(text)
        if m and looks_like_label(m.group(1) + ":"):
            tail = m.group(2).strip()
            if not tail or BLANK_LINE_RE.match(tail):
                n = Node(node_id=node_id, type=NodeType.VALUE_REGION,
                         location=Location(DocType.PPTX, {**parts, "mode": "after_colon",
                                                          "label_text": m.group(1).strip()}),
                         label=m.group(1).strip(), section=sec.title, style=style)
                n.meta["multiline_capable"] = True
                ir.nodes.append(n)
                sec.node_ids.append(n.node_id)
                return

        # empty placeholder inherited from the layout
        if (not text or is_prompt) and shape.is_placeholder:
            label = self._placeholder_label(shape)
            n = Node(node_id=node_id, type=NodeType.VALUE_REGION,
                     location=Location(DocType.PPTX, {**parts, "mode": "replace"}),
                     label=label, section=sec.title, style=style, is_empty=True)
            n.meta["multiline_capable"] = True
            ir.nodes.append(n)
            sec.node_ids.append(n.node_id)
            return

        if text:
            n = Node(node_id=node_id, type=NodeType.STATIC,
                     location=Location(DocType.PPTX, parts), text=text,
                     section=sec.title, style=style, editable=False, is_empty=False)
            ir.nodes.append(n)
            sec.node_ids.append(n.node_id)

    # ------------------------------------------------------------------
    def _parse_table(self, ir, sec, si, shape, parts) -> None:
        tbl = shape.table
        if len(tbl.rows) < 1 or len(tbl.columns) < 2:
            return
        header = [c.text.strip() for c in tbl.rows[0].cells]
        cols = [ColumnIR(index=i, header_text=h, location_hint={"col": i})
                for i, h in enumerate(header)]
        existing = sum(1 for r in list(tbl.rows)[1:] if any(c.text.strip() for c in r.cells))
        t = TableIR(node_id=f"slide{si}::table{shape.shape_id}",
                    location=Location(DocType.PPTX, {**parts, "kind": "table"}),
                    columns=cols, header_row=0, template_row=1,
                    existing_data_rows=existing, section=sec.title,
                    meta={"n_rows": len(tbl.rows), "n_cols": len(tbl.columns),
                          "fixed_rows": True})
        ir.tables.append(t)
        sec.table_ids.append(t.node_id)

    def _slide_title(self, slide) -> Optional[str]:
        try:
            if slide.shapes.title is not None:
                t = (slide.shapes.title.text or "").strip()
                if t and not PROMPT_TEXT_RE.search(t):
                    return t
        except Exception:
            pass
        return None

    def _placeholder_label(self, shape) -> str:
        try:
            ptype = str(shape.placeholder_format.type).split(" ")[0].replace("_", " ").title()
        except Exception:
            ptype = "Content"
        name = (shape.name or "").strip()
        return name if name and not name.lower().startswith(("textbox", "content placeholder")) else ptype

    def _shape_style(self, shape) -> StyleSignals:
        try:
            para = shape.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else None
            font = run.font if run is not None else para.font
            return StyleSignals(
                bold=bool(font.bold), italic=bool(font.italic),
                font_size=float(font.size.pt) if font.size else None,
                font_name=font.name,
                alignment=str(para.alignment) if para.alignment is not None else None,
            )
        except Exception:
            return StyleSignals()


register("pptx", PptxParser)
register("potx", PptxParser)
