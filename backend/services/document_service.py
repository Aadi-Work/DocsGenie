from __future__ import annotations

import io
import logging
import re
import tempfile
from datetime import datetime, timezone
from functools import lru_cache
from pathlib import Path
from typing import Any

from app.services.s3_service import get_s3
from app.utils.file_utils import (
    AppError,
    content_type_for,
    dedupe_answers,
    file_ext,
    require_office,
    safe_filename,
    snake,
)

log = logging.getLogger(__name__)


class DocumentService:
    def generate(
        self,
        *,
        template_key: str,
        answers: dict[str, Any],
        generated_by: str = "",
        template_id: str = "",
        template_name: str = "",
        template_version: str = "",
    ) -> dict[str, Any]:
        s3 = get_s3()
        filename, data = s3.get_object_with_name(template_key)
        require_office(filename)
        fill_answers = answers
        if template_id:
            try:
                from app.office.profiles import answers_for_fill
                from app.services.template_service import get_templates

                fill_answers = answers_for_fill(get_templates().get(template_id), answers)
            except Exception:
                fill_answers = answers
        filled = fill_office(filename, data, fill_answers)
        stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        out_name = f"{Path(filename).stem}_{stamp}{file_ext(filename)}"
        day = datetime.now(timezone.utc).strftime("%Y/%m/%d")
        key = f"{s3.documents_prefix}{day}/{out_name}"
        meta = s3.upload_object(
            key,
            filled,
            content_type=content_type_for(out_name),
            metadata={
                "template_id": template_id or template_key,
                "template_name": template_name or Path(filename).stem,
                "template_version": template_version or "",
                "generated_by": generated_by,
            },
        )
        return {
            "success": True,
            "document_name": out_name,
            "filename": out_name,
            "s3_key": meta["key"],
            "s3_uri": meta["s3_uri"],
            "size": meta["size"],
            "filled_fields": dedupe_answers(answers),
            "fill_mode": "placeholders",
            "template_id": template_id or template_key,
            "template_name": template_name or Path(filename).stem,
            "template_version": template_version or "",
            "generated_by": generated_by,
        }

    def list_generated(self, *, generated_by: str | None = None, limit: int = 50) -> list[dict[str, Any]]:
        s3 = get_s3()
        prefixes = [s3.documents_prefix]
        if s3.documents_prefix.lower() != "documents/":
            prefixes.append("Documents/")
        seen: set[str] = set()
        objects: list[dict[str, Any]] = []
        for prefix in prefixes:
            try:
                for item in s3.list_objects(prefix, limit=400):
                    if item["key"] in seen:
                        continue
                    seen.add(item["key"])
                    objects.append(item)
            except AppError:
                continue
        objects.sort(key=lambda i: i.get("last_modified") or "", reverse=True)
        wanted = (generated_by or "").strip().lower()
        out: list[dict[str, Any]] = []
        for obj in objects:
            try:
                head = s3.get_object_metadata(obj["key"])
            except AppError:
                head = {"metadata": {}, "size": obj.get("size"), "last_modified": obj.get("last_modified"), "s3_uri": f"s3://{s3.bucket}/{obj['key']}"}
            meta = {str(k).lower(): str(v) for k, v in (head.get("metadata") or {}).items()}
            owner = (meta.get("generated_by") or "").strip().lower()
            if wanted and owner and owner != wanted:
                continue
            filename = obj["name"]
            from app.utils.file_utils import download_url, preview_url

            out.append(
                {
                    "id": obj["key"],
                    "template_id": meta.get("template_id") or "",
                    "template_name": meta.get("template_name") or Path(filename).stem,
                    "template_version": meta.get("template_version") or "",
                    "filename": filename,
                    "document_name": meta.get("template_name") or filename,
                    "download_url": download_url(filename, obj["key"]),
                    "preview_url": preview_url(filename, obj["key"]),
                    "generated_by": meta.get("generated_by") or "",
                    "created_at": obj.get("last_modified") or head.get("last_modified") or "",
                    "modified_at": obj.get("last_modified") or head.get("last_modified") or "",
                    "s3_key": obj["key"],
                    "s3_uri": head.get("s3_uri") or f"s3://{s3.bucket}/{obj['key']}",
                    "status": "generated",
                    "size": head.get("size") or obj.get("size"),
                }
            )
            if len(out) >= limit:
                break
        return out

    def find_by_filename(self, filename: str) -> str:
        name = safe_filename(filename).lower()
        for doc in self.list_generated(limit=200):
            if str(doc.get("filename") or "").lower() == name:
                return str(doc["s3_key"])
        raise AppError(404, f"No generated document named {filename}")


def fill_office(filename: str, data: bytes, answers: dict[str, Any]) -> bytes:
    ext = file_ext(filename)
    mapping = _mapping(answers)
    try:
        from app.office.smart import smart_fill

        structural = smart_fill(filename, data, answers)
        if structural:
            data = structural
    except Exception as exc:
        log.exception("Structural fill failed for %s; using placeholder fill", filename)
        from app.office.smart import match_profile

        if match_profile(filename) == "sample_ppt":
            raise AppError(500, f"Could not fill Sample_ppt: {exc}") from exc
    if ext == ".docx":
        return _fill_docx(data, mapping)
    if ext == ".xlsx":
        return _fill_xlsx(data, mapping)
    if ext == ".pptx":
        return _fill_pptx(data, mapping)
    raise AppError(400, "Unsupported template type")


def _mapping(answers: dict[str, Any]) -> dict[str, str]:
    out: dict[str, str] = {}
    for key, value in (answers or {}).items():
        if value in (None, ""):
            continue
        text = value if isinstance(value, str) else str(value)
        out[key] = text
        out[snake(key)] = text
        out[key.replace(" ", "_")] = text
        out[key.lower()] = text
    return out


def _replace_text(text: str, mapping: dict[str, str]) -> str:
    if not text:
        return text
    result = text
    for key, value in mapping.items():
        result = result.replace(f"{{{{{key}}}}}", value)
        result = result.replace(f"{{{{ {key} }}}}", value)
        result = result.replace(f"[{key}]", value)
        result = result.replace(f"«{key}»", value)
        result = result.replace(f"<<{key}>>", value)
        result = result.replace(f"${{{key}}}", value)
    return result


def _fill_docx(data: bytes, mapping: dict[str, str]) -> bytes:
    from docx import Document
    from docxtpl import DocxTemplate

    if b"{{" in data:
        tmp = Path(tempfile.gettempdir()) / "ymsli-template-hub"
        tmp.mkdir(parents=True, exist_ok=True)
        src = tmp / "in.docx"
        dest = tmp / "out.docx"
        try:
            src.write_bytes(data)
            tpl = DocxTemplate(str(src))
            tpl.render(mapping)
            tpl.save(str(dest))
            return dest.read_bytes()
        except Exception:
            log.info("docxtpl render failed; using literal replacement")
        finally:
            src.unlink(missing_ok=True)
            dest.unlink(missing_ok=True)

    doc = Document(io.BytesIO(data))

    def rewrite_paragraph(paragraph) -> None:
        updated = _replace_text(paragraph.text, mapping)
        if updated != paragraph.text:
            paragraph.text = updated

    for para in doc.paragraphs:
        rewrite_paragraph(para)
    for table in doc.tables:
        for row in table.rows:
            cells = list(row.cells)
            for i, cell in enumerate(cells):
                for para in cell.paragraphs:
                    rewrite_paragraph(para)
                label = snake(cell.text.strip().rstrip(":"))
                if label in mapping and i + 1 < len(cells):
                    nxt = cells[i + 1]
                    if not (nxt.text or "").strip() or nxt.text.strip() in {"-", "—"} or nxt.text.strip().startswith("["):
                        if nxt.paragraphs:
                            nxt.paragraphs[0].text = mapping[label]
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _fill_xlsx(data: bytes, mapping: dict[str, str]) -> bytes:
    from openpyxl import load_workbook
    from openpyxl.cell.cell import MergedCell

    wb = load_workbook(io.BytesIO(data))
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for idx, cell in enumerate(row):
                if isinstance(cell, MergedCell):
                    continue
                orig = cell.value
                if isinstance(orig, str):
                    updated = _replace_text(orig, mapping)
                    if updated != orig:
                        _xlsx_set(cell, updated)
                if not isinstance(orig, str) or _xlsx_is_value_slot(orig):
                    continue
                label = snake(orig.strip().rstrip(":"))
                if label not in mapping or idx + 1 >= len(row):
                    continue
                target = _xlsx_anchor(ws, row[idx + 1])
                if target is None or target.coordinate == cell.coordinate:
                    continue
                if _xlsx_is_value_slot(target.value):
                    _xlsx_set(target, mapping[label])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _xlsx_is_value_slot(value: Any) -> bool:
    if value in (None, "", "-", "—"):
        return True
    if not isinstance(value, str):
        return False
    text = value.strip()
    return text.startswith("[") or "{{" in text or text.startswith("<<") or text.strip("_") == ""


def _xlsx_anchor(ws, cell):
    from openpyxl.cell.cell import MergedCell

    if cell is None:
        return None
    if not isinstance(cell, MergedCell):
        return cell
    for merged in ws.merged_cells.ranges:
        if cell.coordinate in merged:
            return ws.cell(merged.min_row, merged.min_col)
    return None


def _xlsx_set(cell, value) -> None:
    try:
        cell.value = value
    except AttributeError:
        return


def _fill_pptx(data: bytes, mapping: dict[str, str]) -> bytes:
    from pptx import Presentation

    from app.office.pptx_text import iter_text_frames, set_paragraph_text

    pres = Presentation(io.BytesIO(data))
    for frame in iter_text_frames(pres):
        for para in frame.paragraphs:
            full = "".join(run.text or "" for run in para.runs) or para.text
            updated = _replace_text(full, mapping)
            if updated != full:
                set_paragraph_text(para, updated)
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


@lru_cache
def get_documents() -> DocumentService:
    return DocumentService()
