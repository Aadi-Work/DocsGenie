from __future__ import annotations

import html
import io
from functools import lru_cache

from app.utils.file_utils import AppError, file_ext, pretty_template_name


class PreviewService:
    def render_html(self, filename: str, data: bytes) -> str:
        ext = file_ext(filename)
        title = html.escape(pretty_template_name(filename, filename))
        if ext == ".docx":
            body = self._docx(data)
        elif ext == ".xlsx":
            body = self._xlsx(data)
        elif ext == ".pptx":
            body = self._pptx(data)
        else:
            raise AppError(400, "Preview is only available for DOCX, PPTX, and XLSX")
        return (
            "<!DOCTYPE html><html><head><meta charset='utf-8'/>"
            f"<title>{title}</title><style>{_CSS}</style></head>"
            f"<body><header><p class='kicker'>In-app preview</p><h1>{title}</h1>"
            "<p class='note'>This is a browser preview. The original file stays in S3.</p></header>"
            f"{body}</body></html>"
        )

    def _docx(self, data: bytes) -> str:
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = html.escape(para.text)
            style = (para.style.name or "") if para.style is not None else ""
            if not text.strip():
                parts.append("<p class='empty'>&nbsp;</p>")
                continue
            if "Heading 1" in style:
                parts.append(f"<h2>{text}</h2>")
            elif "Heading" in style:
                parts.append(f"<h3>{text}</h3>")
            else:
                parts.append(f"<p>{text}</p>")
        for table in doc.tables:
            parts.append(_table([[c.text for c in row.cells] for row in table.rows]))
        return "<article class='page'>" + "".join(parts) + "</article>"

    def _xlsx(self, data: bytes) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), data_only=True)
        chunks: list[str] = []
        for ws in wb.worksheets:
            chunks.append(_xlsx_sheet(ws))
        return "".join(chunks) or "<p>Empty workbook</p>"

    def _pptx(self, data: bytes) -> str:
        from pptx import Presentation

        pres = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for i, slide in enumerate(pres.slides, start=1):
            bits: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = html.escape(shape.text_frame.text)
                    if text.strip():
                        bits.append(f"<p>{text.replace(chr(10), '<br/>')}</p>")
            slides.append(f"<section class='slide'><p class='kicker'>Slide {i}</p>{''.join(bits) or '<p>Empty slide</p>'}</section>")
        return "".join(slides) or "<p>Empty presentation</p>"


def _xlsx_sheet(ws) -> str:
    merges: dict[tuple[int, int], tuple[int, int]] = {}
    covered: set[tuple[int, int]] = set()
    for rng in ws.merged_cells.ranges:
        if rng.min_row > 200 or rng.min_col > 40:
            continue
        max_r = min(rng.max_row, 200)
        max_c = min(rng.max_col, 40)
        merges[(rng.min_row, rng.min_col)] = (max_r - rng.min_row + 1, max_c - rng.min_col + 1)
        for row in range(rng.min_row, max_r + 1):
            for col in range(rng.min_col, max_c + 1):
                if (row, col) != (rng.min_row, rng.min_col):
                    covered.add((row, col))

    raw_max_row = ws.max_row or 1
    raw_max_col = ws.max_column or 1
    max_row = min(raw_max_row, 80 if raw_max_row > 2000 else 200)
    max_col = min(raw_max_col, 20 if raw_max_col > 256 else 24)
    values: dict[tuple[int, int], str] = {}
    last_row = 0
    last_col = 0
    for row in range(1, max_row + 1):
        for col in range(1, max_col + 1):
            if (row, col) in covered:
                continue
            raw = ws.cell(row, col).value
            text = "" if raw is None else str(raw)
            if text.strip():
                values[(row, col)] = text
                span_rows, span_cols = merges.get((row, col), (1, 1))
                last_row = max(last_row, row + span_rows - 1)
                last_col = max(last_col, col + span_cols - 1)
    if last_row == 0:
        return (
            f"<section class='sheet'><h2>{html.escape(ws.title)}</h2>"
            "<p class='muted'>Empty sheet</p></section>"
        )

    skip = set(covered)
    body: list[str] = []
    for row in range(1, last_row + 1):
        cells: list[str] = []
        col = 1
        while col <= last_col:
            if (row, col) in skip:
                col += 1
                continue
            span_rows, span_cols = merges.get((row, col), (1, 1))
            text = values.get((row, col), "")
            if text and span_rows == 1 and span_cols == 1:
                nxt = col + 1
                while (
                    nxt <= last_col
                    and (row, nxt) not in skip
                    and (row, nxt) not in merges
                    and values.get((row, nxt), "") == text
                ):
                    skip.add((row, nxt))
                    span_cols += 1
                    nxt += 1
            attrs = ""
            if span_cols > 1:
                attrs += f" colspan='{span_cols}'"
            if span_rows > 1:
                attrs += f" rowspan='{span_rows}'"
            if span_cols > 1 or span_rows > 1:
                attrs += " class='merged'"
            cells.append(f"<td{attrs}>{html.escape(text).replace(chr(10), '<br/>')}</td>")
            col += span_cols
        if not cells:
            continue
        has_text = any(
            values.get((row, c), "").strip()
            for c in range(1, last_col + 1)
            if (row, c) not in covered
        )
        if has_text or any((row, c) in merges for c in range(1, last_col + 1)):
            body.append("<tr>" + "".join(cells) + "</tr>")
    return (
        f"<section class='sheet'><h2>{html.escape(ws.title)}</h2>"
        f"<div class='table-wrap'><table class='xlsx'><tbody>{''.join(body)}</tbody></table></div></section>"
    )


def _table(rows: list[list[str]], header: bool = False) -> str:
    if not rows:
        return "<p class='muted'>No table content</p>"
    def cells(row: list[str], tag: str) -> str:
        return "".join(f"<{tag}>{html.escape(c)}</{tag}>" for c in row)
    head = ""
    body_rows = rows
    if header:
        head = f"<thead><tr>{cells(rows[0], 'th')}</tr></thead>"
        body_rows = rows[1:]
    body = "".join(f"<tr>{cells(r, 'td')}</tr>" for r in body_rows)
    return f"<div class='table-wrap'><table>{head}<tbody>{body}</tbody></table></div>"


_CSS = """
:root { color-scheme: light; }
body { margin:0; font-family: Segoe UI, system-ui, sans-serif; background:#eef1f7; color:#1c2434; }
header, .page, .slide { max-width: 980px; margin: 0 auto; }
.sheet { max-width: 1100px; margin: 12px auto 20px; }
header { padding: 24px 24px 8px; }
h1 { margin: 4px 0 8px; font-size: 22px; }
h2 { font-size: 18px; }
.kicker { text-transform: uppercase; letter-spacing: .08em; font-size: 11px; color:#5b6b86; margin:0; }
.note, .muted { color:#5b6b86; font-size: 13px; }
.page, .sheet, .slide { background:#fff; padding: 28px 32px; margin: 12px auto 20px; box-shadow: 0 8px 24px rgba(28,36,52,.08); border-radius: 12px; }
.slide { min-height: 280px; }
p, td, th { overflow-wrap: anywhere; word-break: break-word; }
.table-wrap { overflow: auto; margin: 12px 0 20px; }
table { border-collapse: collapse; width: 100%; font-size: 13px; }
table.xlsx { width: max-content; min-width: 100%; table-layout: auto; }
th, td { border: 1px solid #d5dbe8; padding: 6px 8px; vertical-align: top; max-width: 28rem; }
td.merged { background: #f8fafc; }
th { background: #f3f6fb; text-align: left; }
"""


@lru_cache
def get_preview() -> PreviewService:
    return PreviewService()
