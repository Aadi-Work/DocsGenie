"""Compare two Office templates by visible text / cell values, not S3 metadata."""

from __future__ import annotations

import difflib
import hashlib
import html
import io
import re
import zipfile
from typing import Any

from app.utils.extract import (
    _CELL_RE,
    _T_RE,
    _V_RE,
    _xlsx_shared_strings,
    _xlsx_sheet_list,
)
from app.utils.file_utils import file_ext

_MAX_LINES = 4000
_MAX_DIFF = 800
_CELL_LINE = re.compile(r"^(.+![A-Z]{1,3}\d+)\s*=\s*(.*)$", re.I)


def file_digest(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()[:16]


def snapshot_office(filename: str, data: bytes) -> list[str]:
    ext = file_ext(filename)
    try:
        if ext == ".xlsx":
            return _snapshot_xlsx(data)
        if ext == ".docx":
            return _snapshot_docx(data)
        if ext == ".pptx":
            return _snapshot_pptx(data)
    except Exception:
        pass
    text = data.decode("utf-8", errors="replace")
    return [ln.rstrip() for ln in text.splitlines() if ln.strip()][:_MAX_LINES]


def diff_snapshots(old: list[str], new: list[str]) -> tuple[list[dict[str, Any]], list[dict[str, str]]]:
    """Return (unified_diff lines, structured changes)."""
    if old == new:
        return [], []
    if _mostly_cells(old) and _mostly_cells(new):
        return _diff_cells(old, new)
    return _diff_sequence(old, new)


def _snapshot_xlsx(data: bytes) -> list[str]:
    lines: list[str] = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        names = set(archive.namelist())
        sheets = _xlsx_sheet_list(archive)
        shared: list[str] = []
        if "xl/sharedStrings.xml" in names:
            shared = _xlsx_shared_strings(archive.read("xl/sharedStrings.xml"))
        for sheet_name, path in sheets:
            if path not in names:
                continue
            xml = archive.read(path).decode("utf-8", errors="ignore")
            cells = _xlsx_cells(xml, shared)
            if not cells and not lines:
                continue
            lines.append(f"# {sheet_name}")
            for ref, value in cells:
                lines.append(f"{sheet_name}!{ref} = {value}")
                if len(lines) >= _MAX_LINES:
                    return lines
    return lines


def _xlsx_cells(xml: str, shared: list[str]) -> list[tuple[str, str]]:
    out: list[tuple[str, str]] = []
    for attrs, body in _CELL_RE.findall(xml):
        ref_m = re.search(r'\br="([^"]+)"', attrs)
        ref = ref_m.group(1) if ref_m else ""
        if ref and not _ref_in_bounds(ref):
            continue
        text = _cell_text(attrs, body, shared)
        if not text:
            continue
        out.append((ref or "?", text))
        if len(out) >= 2500:
            break
    return out


def _cell_text(attrs: str, body: str, shared: list[str]) -> str:
    runs = _T_RE.findall(body)
    if runs:
        text = html.unescape("".join(runs)).strip()
        return _one_line(text)
    kind_m = re.search(r'\bt="([^"]+)"', attrs)
    kind = kind_m.group(1) if kind_m else ""
    value = _V_RE.search(body)
    raw = html.unescape(value.group(1)).strip() if value else ""
    if kind == "s":
        try:
            return _one_line(shared[int(raw)])
        except (ValueError, IndexError):
            return ""
    if kind == "b":
        return "TRUE" if raw in {"1", "true", "TRUE"} else "FALSE" if raw else ""
    return _one_line(raw)


def _ref_in_bounds(ref: str, max_row: int = 200, max_col: int = 40) -> bool:
    match = re.match(r"([A-Z]+)(\d+)$", ref, re.I)
    if not match:
        return True
    col = 0
    for ch in match.group(1).upper():
        col = col * 26 + (ord(ch) - 64)
    return int(match.group(2)) <= max_row and col <= max_col


def _one_line(text: str) -> str:
    return re.sub(r"[\r\n]+", " / ", (text or "")).strip()


def _snapshot_docx(data: bytes) -> list[str]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    lines: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            lines.append(text)
        if len(lines) >= _MAX_LINES:
            return lines
    for t_i, table in enumerate(doc.tables, start=1):
        lines.append(f"# Table {t_i}")
        for row in table.rows:
            cells = [_one_line(c.text) for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
            if len(lines) >= _MAX_LINES:
                return lines
    return lines


def _snapshot_pptx(data: bytes) -> list[str]:
    from pptx import Presentation

    pres = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = (para.text or "").strip()
                if text:
                    lines.append(text)
                if len(lines) >= _MAX_LINES:
                    return lines
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [_one_line(c.text_frame.text if c.text_frame else "") for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
    return lines


def _mostly_cells(lines: list[str]) -> bool:
    content = [ln for ln in lines if ln and not ln.startswith("#")]
    if len(content) < 3:
        return False
    hits = sum(1 for ln in content if _CELL_LINE.match(ln))
    return hits / len(content) >= 0.5


def _parse_cells(lines: list[str]) -> dict[str, str]:
    cells: dict[str, str] = {}
    for line in lines:
        match = _CELL_LINE.match(line)
        if match:
            cells[match.group(1)] = match.group(2)
    return cells


def _addr_key(addr: str) -> tuple:
    sheet, _, cell = addr.partition("!")
    match = re.match(r"([A-Z]+)(\d+)$", cell, re.I)
    if not match:
        return (sheet.lower(), 0, cell)
    col = 0
    for ch in match.group(1).upper():
        col = col * 26 + (ord(ch) - 64)
    return (sheet.lower(), int(match.group(2)), col)


def _diff_cells(old: list[str], new: list[str]) -> tuple[list[dict[str, Any]], list[dict[str, str]]]:
    left = _parse_cells(old)
    right = _parse_cells(new)
    keys = sorted(set(left) | set(right), key=_addr_key)
    unified: list[dict[str, Any]] = []
    changes: list[dict[str, str]] = []
    last_sheet = ""
    for key in keys:
        before = left.get(key)
        after = right.get(key)
        if before == after:
            continue
        sheet = key.split("!", 1)[0]
        if sheet != last_sheet:
            unified.append({"type": "context", "text": f"# {sheet}"})
            last_sheet = sheet
        if before is None:
            unified.append({"type": "added", "text": f"{key} = {after}"})
            changes.append({"field": key, "change": "added", "before": "", "after": after or ""})
        elif after is None:
            unified.append({"type": "removed", "text": f"{key} = {before}"})
            changes.append({"field": key, "change": "removed", "before": before, "after": ""})
        else:
            unified.append({"type": "removed", "text": f"{key} = {before}"})
            unified.append({"type": "added", "text": f"{key} = {after}"})
            changes.append({"field": key, "change": "updated", "before": before, "after": after})
        if len(unified) >= _MAX_DIFF:
            unified.append({"type": "context", "text": f"… truncated after {_MAX_DIFF} diff lines"})
            break
    return unified, changes


def _diff_sequence(old: list[str], new: list[str]) -> tuple[list[dict[str, Any]], list[dict[str, str]]]:
    matcher = difflib.SequenceMatcher(a=old, b=new, autojunk=False)
    unified: list[dict[str, Any]] = []
    changes: list[dict[str, str]] = []
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "equal":
            if i2 - i1 <= 4:
                block = old[i1:i2]
            else:
                block = [*old[i1 : i1 + 2], "…", *old[i2 - 2 : i2]]
            for line in block:
                unified.append({"type": "context", "text": line})
            continue
        removed = old[i1:i2]
        added = new[j1:j2]
        for line in removed:
            unified.append({"type": "removed", "text": line})
        for line in added:
            unified.append({"type": "added", "text": line})
        n = max(len(removed), len(added), 1)
        for k in range(n):
            before = removed[k] if k < len(removed) else ""
            after = added[k] if k < len(added) else ""
            if before == after:
                continue
            if not before:
                change = "added"
            elif not after:
                change = "removed"
            else:
                change = "updated"
            changes.append(
                {
                    "field": f"line {i1 + k + 1}",
                    "change": change,
                    "before": before,
                    "after": after,
                }
            )
        if len(unified) >= _MAX_DIFF:
            unified.append({"type": "context", "text": f"… truncated after {_MAX_DIFF} diff lines"})
            break
    return unified, changes
