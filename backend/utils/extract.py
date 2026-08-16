from __future__ import annotations

import html
import io
import json
import logging
import re
import zipfile
from typing import Any
from xml.etree import ElementTree as ET

log = logging.getLogger(__name__)

# {{token}} / {{#SECTION}} / {{/SECTION}} as used in BFL_Template.xlsx
MUSTACHE_RE = re.compile(r"\{\{\s*([#^/]?)([^}]+?)\s*\}\}")
PLACEHOLDER_RE = re.compile(
    r"\{\{\s*([A-Za-z_][\w.\- ]{0,80})\s*\}\}|"
    r"\{%\s*([A-Za-z][A-Za-z0-9_ ]{0,60})\s*%\}|"
    r"«\s*([A-Za-z][A-Za-z0-9_ ]{0,60})\s*»|"
    r"<<\s*([A-Za-z][A-Za-z0-9_ ]{0,60})\s*>>|"
    r"\[([A-Za-z][A-Za-z0-9_ /-]{1,60})\]"
)
LABEL_RE = re.compile(r"^(.{2,48}?)\s*[:：]\s*(\{\{[^}]+\}\}|_+|\[.*\]|\s*)$")
SKIP_LABELS = {
    "if",
    "else",
    "each",
    "with",
    "unless",
    "this",
    "for",
    "endif",
    "endfor",
    "none",
    "n_a",
    "sheet1",
    "sheet2",
    "total",
    "sum",
    "s_no",
    "sr_no",
    "sno",
}
_XLSX_NS = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
_CELL_RE = re.compile(r"<c\b([^>]*)>(.*?)</c>", re.S)
_T_RE = re.compile(r"<t\b[^>]*>(.*?)</t>", re.S)
_V_RE = re.compile(r"<v\b[^>]*>(.*?)</v>", re.S)

from app.utils.file_utils import snake, title_from_ident


def _add_field(
    fields: dict[str, dict[str, Any]],
    ident: str,
    label: str | None = None,
    *,
    source: str = "detected",
) -> None:
    key = snake(ident)
    if not key or key in SKIP_LABELS:
        return
    if key not in fields:
        nice = label or title_from_ident(ident)
        fields[key] = {
            "id": key,
            "label": nice,
            "question": f"What is the {nice}?",
            "required": True,
            "field_type": "string",
            "source": source,
        }


def _mustache_ident(raw: str) -> str:
    token = (raw or "").strip().split("|", 1)[0].strip()
    token = token.split()[0] if token else ""
    return re.sub(r"[^\w.\- ]+", "", token).strip("._-")


def harvest_mustache(
    text: str,
    fields: dict[str, dict[str, Any]],
    sections: list[str] | None = None,
) -> int:
    """Collect {{token}} values. {{#block}} / {{/block}} become section names, not fields."""
    added = 0
    for match in MUSTACHE_RE.finditer(text or ""):
        prefix = match.group(1) or ""
        ident = _mustache_ident(match.group(2) or "")
        if not ident:
            continue
        if prefix in {"#", "/", "^"}:
            if sections is not None and ident not in sections and snake(ident) not in SKIP_LABELS:
                sections.append(ident)
            continue
        before = len(fields)
        _add_field(fields, ident, source="placeholder")
        added += len(fields) - before
    return added


def harvest(text: str, fields: dict[str, dict[str, Any]]) -> None:
    harvest_mustache(text, fields)
    for match in PLACEHOLDER_RE.finditer(text or ""):
        ident = next((g for g in match.groups() if g), "")
        if ident:
            _add_field(fields, ident)
    for line in (text or "").splitlines():
        matched = LABEL_RE.match(line.strip())
        if matched:
            _add_field(fields, matched.group(1), matched.group(1).strip())


def extract_office(filename: str, data: bytes) -> dict[str, Any]:
    name = (filename or "").lower()
    fields: dict[str, dict[str, Any]] = {}
    sections: list[str] = []
    tables = 0
    text = ""
    try:
        if name.endswith(".docx"):
            text, tables, sections = _docx(data, fields)
        elif name.endswith(".xlsx"):
            text, tables, sections = _xlsx(data, fields)
        elif name.endswith(".pptx"):
            text, tables, sections = _pptx(data, fields)
        else:
            text = data.decode("utf-8", errors="ignore")[:20000]
            harvest(text, fields)
    except Exception as exc:
        return {
            "ok": False,
            "error": str(exc),
            "placeholders": [],
            "field_config": [],
            "context_questions": [],
            "content_outline": [],
            "tables": 0,
            "preview_text": "",
            "literal_placeholders": False,
        }
    literal = any(f.get("source") == "placeholder" for f in fields.values())
    if not literal:
        _merge_engine_fields(filename, data, fields)
    config = list(fields.values())
    return {
        "ok": True,
        "placeholders": [f["label"] for f in config],
        "field_config": config,
        "context_questions": [f["question"] for f in config],
        "content_outline": sections or [f["label"] for f in config[:8]],
        "tables": tables,
        "preview_text": text[:4000],
        "text": text[:24000],
        "literal_placeholders": literal,
    }


def _docx(data: bytes, fields: dict[str, dict[str, Any]]) -> tuple[str, int, list[str]]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    chunks: list[str] = []
    sections: list[str] = []
    for para in doc.paragraphs:
        chunks.append(para.text)
        harvest(para.text, fields)
        style = (para.style.name or "") if para.style is not None else ""
        if "Heading" in style and para.text.strip():
            sections.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            labels = [c.text.strip() for c in row.cells]
            chunks.append(" | ".join(labels))
            for cell in row.cells:
                harvest(cell.text, fields)
            if len(labels) >= 2:
                if not labels[1] or labels[1] in {"", "-", "—"} or labels[1].startswith("[") or "{{" in labels[1]:
                    _add_field(fields, labels[0], labels[0])
                elif all(len(x) <= 40 for x in labels if x):
                    for label in labels:
                        if label and not label.isdigit():
                            _add_field(fields, label, label)
    return "\n".join(chunks), len(doc.tables), sections


def _blank(value: Any) -> bool:
    if value is None:
        return True
    if isinstance(value, str):
        text = value.strip()
        return text in {"", "-", "—", "...", "n/a", "N/A"} or text.startswith("[") or "{{" in text or text.strip("_") == ""
    return False


def _xlsx_rel_path(target: str) -> str:
    path = (target or "").replace("\\", "/").lstrip("/")
    return path if path.startswith("xl/") else f"xl/{path}"


def _xlsx_sheet_list(archive: zipfile.ZipFile) -> list[tuple[str, str]]:
    workbook = ET.fromstring(archive.read("xl/workbook.xml"))
    rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
    rid_to_path = {
        rel.attrib.get("Id"): _xlsx_rel_path(rel.attrib.get("Target") or "")
        for rel in rels
        if rel.attrib.get("Id")
    }
    sheets: list[tuple[str, str]] = []
    for sheet in workbook.findall("m:sheets/m:sheet", _XLSX_NS):
        name = sheet.attrib.get("name") or "Sheet"
        path = rid_to_path.get(sheet.attrib.get(_R_ID) or "")
        if path:
            sheets.append((name, path))
    return sheets


def _xlsx_shared_strings(xml: bytes) -> list[str]:
    root = ET.fromstring(xml)
    out: list[str] = []
    for item in root.findall("m:si", _XLSX_NS):
        parts = [html.unescape(node.text or "") for node in item.iter() if node.tag.endswith("}t") or node.tag == "t"]
        out.append("".join(parts))
    return out


def _xlsx_cell_texts(xml: str, shared: list[str]) -> list[str]:
    texts: list[str] = []
    for attrs, body in _CELL_RE.findall(xml):
        runs = _T_RE.findall(body)
        if runs:
            texts.append(html.unescape("".join(runs)))
            continue
        kind = ""
        match = re.search(r'\bt="([^"]+)"', attrs)
        if match:
            kind = match.group(1)
        if kind != "s":
            continue
        value = _V_RE.search(body)
        if not value:
            continue
        try:
            texts.append(shared[int(html.unescape(value.group(1)).strip())])
        except (ValueError, IndexError):
            continue
    return texts


def _xlsx_scan_zip(data: bytes) -> tuple[list[str], list[str], list[str], list[str]]:
    """Read every worksheet tab and collect {{placeholder}} tokens without loading styles."""
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        names = set(archive.namelist())
        sheets = _xlsx_sheet_list(archive)
        shared: list[str] = []
        if "xl/sharedStrings.xml" in names:
            shared = _xlsx_shared_strings(archive.read("xl/sharedStrings.xml"))
        sheet_names = [name for name, _ in sheets]
        texts: list[str] = []
        tokens: list[str] = []
        blocks: list[str] = []
        seen_tokens: set[str] = set()
        seen_blocks: set[str] = set()
        for _name, path in sheets:
            if path not in names:
                continue
            xml = archive.read(path).decode("utf-8", errors="ignore")
            for cell_text in _xlsx_cell_texts(xml, shared):
                if cell_text.strip():
                    texts.append(cell_text)
                for match in MUSTACHE_RE.finditer(cell_text):
                    prefix = match.group(1) or ""
                    ident = _mustache_ident(match.group(2) or "")
                    if not ident:
                        continue
                    key = ident.lower()
                    if prefix in {"#", "/", "^"}:
                        if key not in seen_blocks and snake(ident) not in SKIP_LABELS:
                            seen_blocks.add(key)
                            blocks.append(ident)
                    elif key not in seen_tokens:
                        seen_tokens.add(key)
                        tokens.append(ident)
        return sheet_names, texts, tokens, blocks


def _xlsx(data: bytes, fields: dict[str, dict[str, Any]]) -> tuple[str, int, list[str]]:
    sheet_names, texts, tokens, blocks = _xlsx_scan_zip(data)
    chunks = [f"# {name}" for name in sheet_names]
    chunks.extend(texts[:120])
    if tokens or blocks:
        for ident in tokens:
            _add_field(fields, ident, source="placeholder")
        sections = list(sheet_names)
        for block in blocks:
            if block not in sections:
                sections.append(block)
        return "\n".join(chunks), len(sheet_names), sections
    return _xlsx_labels(data, fields, sheet_names)


def _xlsx_labels(
    data: bytes,
    fields: dict[str, dict[str, Any]],
    sheet_names: list[str] | None = None,
) -> tuple[str, int, list[str]]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=False, read_only=True)
    chunks: list[str] = []
    sections = list(sheet_names or wb.sheetnames)
    try:
        for ws in wb.worksheets:
            chunks.append(f"# {ws.title}")
            header_done = False
            for row in ws.iter_rows(min_row=1, max_row=80, max_col=24):
                values = ["" if c.value is None else str(c.value) for c in row]
                if any(values):
                    chunks.append(" | ".join(values))
                nonempty = [str(c.value).strip() for c in row if isinstance(c.value, str) and str(c.value).strip()]
                if not header_done and len(nonempty) >= 2 and all(len(v) <= 40 for v in nonempty):
                    for label in nonempty:
                        _add_field(fields, label, label)
                    header_done = True
                    continue
                for idx, cell in enumerate(row):
                    raw = cell.value
                    if not isinstance(raw, str):
                        continue
                    harvest(raw, fields)
                    label = raw.strip().rstrip(":")
                    nxt = row[idx + 1].value if idx + 1 < len(row) else None
                    if label and 2 <= len(label) <= 48 and _blank(nxt) and not label.isdigit():
                        _add_field(fields, label, label)
                header_done = True
    finally:
        wb.close()
    return "\n".join(chunks), len(sections), sections


def _pptx(data: bytes, fields: dict[str, dict[str, Any]]) -> tuple[str, int, list[str]]:
    from pptx import Presentation

    pres = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    sections: list[str] = []
    tables = 0
    for i, slide in enumerate(pres.slides, start=1):
        title = f"Slide {i}"
        bits: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                tables += 1
                for row in shape.table.rows:
                    labels = [(cell.text_frame.text if cell.text_frame else "").strip() for cell in row.cells]
                    chunks.append(" | ".join(labels))
                    for cell in row.cells:
                        text = cell.text_frame.text if cell.text_frame else ""
                        harvest(text, fields)
                    if labels and all(len(x) <= 40 for x in labels if x):
                        for label in labels:
                            if label and not label.isdigit():
                                _add_field(fields, label, label)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text if shape.has_text_frame else ""
            if text.strip():
                bits.append(text)
                harvest(text, fields)
        if bits:
            title = bits[0].splitlines()[0][:80]
            sections.append(title)
        chunks.append(f"# {title}\n" + "\n".join(bits))
    return "\n".join(chunks), tables or len(list(pres.slides)), sections


def _merge_engine_fields(filename: str, data: bytes, fields: dict[str, dict[str, Any]]) -> None:
    try:
        from app.office.smart import detect_engine_fields

        for item in detect_engine_fields(filename, data):
            _add_field(fields, str(item.get("id") or ""), str(item.get("label") or "") or None)
            key = snake(str(item.get("id") or ""))
            if key in fields and item.get("field_type"):
                fields[key]["field_type"] = item["field_type"]
            if key in fields and item.get("source") and fields[key].get("source") == "detected":
                fields[key]["source"] = "engine"
    except Exception:
        log.debug("Engine field merge skipped for %s", filename, exc_info=True)


def extract_plain_upload(filename: str, data: bytes) -> str:
    name = (filename or "").lower()
    if name.endswith((".docx", ".xlsx", ".pptx")):
        return str(extract_office(filename, data).get("text") or "")
    try:
        return data.decode("utf-8", errors="ignore")[:24000]
    except Exception:
        return ""


def dump_fields(fields: list[dict[str, Any]]) -> str:
    return json.dumps(fields, ensure_ascii=False, indent=2)
