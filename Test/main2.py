import argparse
import io
import json
import os
import random
import sys
import time
import zipfile
import math
from copy import copy
from datetime import datetime

import pandas as pd
import requests

# Document parsers
from docx import Document
from docxtpl import DocxTemplate
from pptx import Presentation

# Charting
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# Excel generation
from openpyxl import load_workbook
from openpyxl.drawing.image import Image
from openpyxl.drawing.spreadsheet_drawing import AnchorMarker, OneCellAnchor
from openpyxl.drawing.xdr import XDRPositiveSize2D
from openpyxl.utils import range_boundaries
from openpyxl.utils.units import pixels_to_EMU

# Vector Database & Embeddings
import chromadb
from chromadb.utils import embedding_functions

# Google GenAI SDK
from google import genai
from google.genai import types
from google.genai.errors import APIError

# Environment detection for Google Colab vs Local
try:
    from IPython.display import HTML, display
    from google.colab import files, userdata
    COLAB_ENV = True
except ImportError:
    COLAB_ENV = False

try:
    import pypdf
except ImportError:
    pypdf = None

# ==============================================================================
# 1. CONFIGURATION & CLIENT INITIALIZATION
# ==============================================================================

GEMINI_API_KEY = ""

if not GEMINI_API_KEY and COLAB_ENV:
    try:
        GEMINI_API_KEY = userdata.get("GEMINI_API_KEY")
    except Exception:
        GEMINI_API_KEY = None

CLIENT_ID = "2194a676-0580-4882-a21c-d7d4d0a8966c"
TENANT_ID = "common"
SCOPES = ["Files.Read.All", "Files.ReadWrite.All", "User.Read"]

MODEL_NAME = "gemini-2.5-flash"
CHROMA_PATH = "./onedrive_vector_db"
DELTA_TOKEN_FILE = "./onedrive_delta_token.json"
DEFAULT_MOM_WORD_TEMPLATE = "mom_template.docx"
DEFAULT_MOM_EXCEL_TEMPLATE = "MOM_Template.xlsx"

CURRENT_ACCESS_TOKEN = None
TOKEN_EXPIRES_AT = 0.0
USER_EMAIL = None

FORCE_RESYNC_ON_START = False
MAX_API_CALLS_PER_TURN = 4
MIN_SECONDS_BETWEEN_CALLS = 23.0
_last_gemini_call_time = 0.0

if not GEMINI_API_KEY:
    raise RuntimeError("GEMINI_API_KEY environment variable is not set.")
client = genai.Client(api_key=GEMINI_API_KEY)

def _enforce_global_pacing():
    global _last_gemini_call_time
    elapsed = time.perf_counter() - _last_gemini_call_time
    remaining = MIN_SECONDS_BETWEEN_CALLS - elapsed
    if remaining > 0:
        time.sleep(remaining)
    _last_gemini_call_time = time.perf_counter()

# ==============================================================================
# 2. MICROSOFT GRAPH AUTHENTICATION & SECURITY
# ==============================================================================

def acquire_microsoft_graph_token() -> str:
    global CURRENT_ACCESS_TOKEN, TOKEN_EXPIRES_AT
    
    # Return cached token if valid with a 5-minute buffer
    if CURRENT_ACCESS_TOKEN and time.time() < (TOKEN_EXPIRES_AT - 300):
        return CURRENT_ACCESS_TOKEN

    print("\n--- Initiating Microsoft Graph Authentication ---")
    device_code_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/devicecode"
    payload = {"client_id": CLIENT_ID, "scope": " ".join(SCOPES)}

    res = requests.post(device_code_url, data=payload)
    if res.status_code != 200:
        raise RuntimeError(f"Device Code Auth failed: {res.status_code} - {res.text}")

    data = res.json()
    verification_uri = data.get("verification_uri", "https://microsoft.com/devicelogin")
    user_code = data["user_code"]

    if COLAB_ENV:
        html_banner = f"""
        <div style="background-color: #f0f4f9; border-left: 6px solid #1a73e8; padding: 15px; margin: 10px 0;">
          <h3 style="margin-top: 0; color: #1a73e8;">🔑 Microsoft Authentication Required</h3>
          <p>1. Open Link: <a href="{verification_uri}" target="_blank">{verification_uri}</a></p>
          <p>2. Enter Code: <b>{user_code}</b></p>
        </div>
        """
        display(HTML(html_banner))
    else:
        print(f"\n🔑 Open {verification_uri} and enter code: {user_code}\n")

    token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    token_payload = {
        "grant_type": "urn:ietf:params:oauth:grant-type:device_code",
        "client_id": CLIENT_ID,
        "device_code": data["device_code"],
    }

    start_time = time.time()
    interval = data.get("interval", 5)
    print("⏳ Waiting for browser sign-in completion...")

    while time.time() - start_time < data.get("expires_in", 900):
        time.sleep(interval)
        token_res = requests.post(token_url, data=token_payload)
        token_data = token_res.json()

        if "access_token" in token_data:
            print("✅ Access token acquired successfully!\n")
            CURRENT_ACCESS_TOKEN = token_data["access_token"]
            TOKEN_EXPIRES_AT = time.time() + token_data.get("expires_in", 3600)
            return CURRENT_ACCESS_TOKEN

        error_code = token_data.get("error")
        if error_code == "authorization_pending":
            continue
        elif error_code in ["authorization_declined", "expired_token"]:
            raise RuntimeError(f"❌ OAuth Error: {error_code}")
            
    raise TimeoutError("Authentication device code expired.")


def get_user_permissions_for_file(file_id: str) -> list:
    if not USER_EMAIL:
        return []
        
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    permissions_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/permissions"

    res = requests.get(permissions_url, headers=headers)
    if res.status_code != 200:
        return []

    permissions = res.json().get("value", [])
    email_lower = USER_EMAIL.strip().lower()
    matched_roles = []

    for perm in permissions:
        roles = perm.get("roles", [])
        is_matched = False
        
        for field in ["grantedTo", "grantedToV2"]:
            if field in perm and perm[field]:
                entity = perm[field].get("user") or perm[field].get("siteUser", {})
                if entity.get("email", "").lower() == email_lower:
                    is_matched = True

        for arr_field in ["grantedToIdentities", "grantedToIdentitiesV2"]:
            if arr_field in perm:
                for identity in perm[arr_field]:
                    entity = identity.get("user") or identity.get("siteUser", {})
                    if entity.get("email", "").lower() == email_lower:
                        is_matched = True

        if is_matched:
            matched_roles.extend(roles)

    return list(set(matched_roles))


def save_modifications_to_onedrive(file_id: str, updated_bytes: bytes) -> str:
    user_roles = get_user_permissions_for_file(file_id)
    can_write = any(role in ["write", "owner", "edit"] for role in user_roles)

    if not can_write:
        return "❌ [WRITE ACTION BLOCKED] Security Policy: Your account holds 'read' access only."

    token = acquire_microsoft_graph_token()
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/octet-stream"
    }
    upload_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
    res = requests.put(upload_url, headers=headers, data=updated_bytes)

    if res.status_code in [200, 201]:
        return "✅ Success: Modifications saved back to master storage."
    return f"⚠️ Failed to save changes (HTTP {res.status_code}): {res.text}"

# ==============================================================================
# 3. DOCUMENT PARSING & UTILS
# ==============================================================================

def semantic_chunk_text(text: str, chunk_words: int = 250, overlap_words: int = 50) -> list:
    """Chunks text by word boundaries rather than arbitrary character slicing."""
    words = text.split()
    if not words:
        return [text]
    
    chunks = []
    step = max(1, chunk_words - overlap_words)
    for i in range(0, len(words), step):
        chunks.append(" ".join(words[i:i + chunk_words]))
    return chunks

class ParserRegistry:
    @staticmethod
    def truncate_text(text: str, max_chars: int = 4000) -> str:
        return text[:max_chars] + "\n\n...[Content truncated]..." if len(text) > max_chars else text

    @classmethod
    def parse(cls, content_bytes: bytes, file_name: str) -> str:
        ext = os.path.splitext(file_name)[1].lower()
        try:
            if ext in [".xlsx", ".xls", ".csv"]:
                xls = pd.ExcelFile(io.BytesIO(content_bytes))
                return "\n\n".join([f"**Sheet: {s}**\n" + pd.read_excel(xls, sheet_name=s, nrows=20).to_markdown(index=False) for s in xls.sheet_names[:3]])
            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(io.BytesIO(content_bytes))
                slides = list(prs.slides)[:15]
                return cls.truncate_text("\n".join([f"Slide {i+1}: " + " | ".join([s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]) for i, slide in enumerate(slides)]))
            elif ext in [".docx", ".doc"]:
                doc = Document(io.BytesIO(content_bytes))
                return cls.truncate_text("\n".join([p.text for p in doc.paragraphs if p.text.strip()]))
            elif ext == ".pdf" and pypdf:
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                pages = list(reader.pages)[:10]
                return cls.truncate_text("\n".join([p.extract_text() for p in pages if p.extract_text()]))
            else:
                return cls.truncate_text(content_bytes.decode("utf-8", errors="ignore"))
        except Exception as e:
            return f"Parsing error for {file_name}: {str(e)}"

# ==============================================================================
# 4. VECTOR DB & ONEDRIVE SYNC
# ==============================================================================

embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)
folder_collection = chroma_client.get_or_create_collection(name="folder_paths", embedding_function=embedding_fn)
document_collection = chroma_client.get_or_create_collection(name="document_chunks", embedding_function=embedding_fn)

def fetch_template(template_name: str) -> str:
    if os.path.exists(template_name):
        return template_name

    try:
        results = document_collection.get(where={"file_name": template_name})
        if results and results.get("metadatas") and len(results["metadatas"]) > 0:
            file_id = results["metadatas"][0]["file_id"]
            token = acquire_microsoft_graph_token()
            res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content", headers={"Authorization": f"Bearer {token}"})
            if res.status_code == 200:
                with open(template_name, "wb") as f:
                    f.write(res.content)
                return template_name
    except Exception as e:
        print(f"Template DB fetch error: {e}")

    try:
        token = acquire_microsoft_graph_token()
        search_url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{template_name}')"
        res = requests.get(search_url, headers={"Authorization": f"Bearer {token}"})
        if res.status_code == 200:
            for item in res.json().get("value", []):
                if item.get("name", "").lower() == template_name.lower():
                    file_res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{item['id']}/content", headers={"Authorization": f"Bearer {token}"})
                    if file_res.status_code == 200:
                        with open(template_name, "wb") as f:
                            f.write(file_res.content)
                        return template_name
    except Exception as e:
        print(f"Template Graph fetch error: {e}")
        
    raise FileNotFoundError(f"Template '{template_name}' not found locally or in OneDrive.")


def sync_onedrive_to_vector_db(force_resync: bool = False):
    global folder_collection, document_collection
    if force_resync:
        try:
            chroma_client.delete_collection("folder_paths")
            chroma_client.delete_collection("document_chunks")
            if os.path.exists(DELTA_TOKEN_FILE): os.remove(DELTA_TOKEN_FILE)
        except Exception:
            pass
        folder_collection = chroma_client.create_collection(name="folder_paths", embedding_function=embedding_fn)
        document_collection = chroma_client.create_collection(name="document_chunks", embedding_function=embedding_fn)

    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}

    if os.path.exists(DELTA_TOKEN_FILE) and not force_resync:
        with open(DELTA_TOKEN_FILE, "r") as f: next_url = json.load(f).get("delta_link")
    else:
        next_url = "https://graph.microsoft.com/v1.0/me/drive/root/delta?$select=id,name,createdDateTime,lastModifiedDateTime,folder,file,parentReference,deleted"

    delta_link = None
    processed_count = 0

    while next_url:
        res = requests.get(next_url, headers=headers)
        if res.status_code == 410 or "resyncRequired" in res.text:
            if os.path.exists(DELTA_TOKEN_FILE): os.remove(DELTA_TOKEN_FILE)
            return sync_onedrive_to_vector_db(force_resync=True)

        data = res.json()
        for item in data.get("value", []):
            item_id, item_name = item["id"], item.get("name", "Unknown")
            if "deleted" in item:
                try:
                    document_collection.delete(where={"file_id": item_id})
                    folder_collection.delete(ids=[item_id])
                except: pass
                continue

            if "folder" in item:
                folder_collection.upsert(ids=[item_id], documents=[f"Folder: {item_name}"], metadatas=[{"folder_name": item_name, "full_path": item_name}])
            elif "file" in item:
                file_res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content", headers=headers)
                if file_res.status_code == 200:
                    text = ParserRegistry.parse(file_res.content, item_name)
                    chunks = semantic_chunk_text(text)
                    ids = [f"{item_id}_chunk_{i}" for i in range(len(chunks))]
                    metas = [{"file_id": item_id, "file_name": item_name, "created_date": item.get("createdDateTime", ""), "chunk_idx": i} for i in range(len(chunks))]
                    document_collection.upsert(ids=ids, documents=chunks, metadatas=metas)
                    processed_count += 1
                    
        next_url = data.get("@odata.nextLink")
        if "@odata.deltaLink" in data: delta_link = data["@odata.deltaLink"]

    if delta_link:
        with open(DELTA_TOKEN_FILE, "w") as f: json.dump({"delta_link": delta_link}, f)
    print(f"✅ OneDrive sync finished. Indexed {processed_count} files.")

# ==============================================================================
# 5. EXCEL MOM GENERATION UTILITIES
# ==============================================================================

def preserve_template_logo(ws, template_path):
    if len(ws._images) == 0:
        try:
            with zipfile.ZipFile(template_path, 'r') as z:
                media_files = [f for f in z.namelist() if 'xl/media/' in f]
                if media_files:
                    img_data = z.read(media_files[0])
                    img = Image(io.BytesIO(img_data))
                    ws['A1'].value = None
                    marker = AnchorMarker(col=6, colOff=pixels_to_EMU(180), row=0, rowOff=pixels_to_EMU(2))
                    size = XDRPositiveSize2D(pixels_to_EMU(img.width or 100), pixels_to_EMU(img.height or 24))
                    img.anchor = OneCellAnchor(_from=marker, ext=size)
                    ws.add_image(img)
        except Exception as e:
            print(f"Warning: Could not extract logo: {e}")

def clone_cell_style(source_cell, target_cell):
    if source_cell.has_style:
        target_cell.font = copy(source_cell.font)
        target_cell.border = copy(source_cell.border)
        target_cell.fill = copy(source_cell.fill)
        target_cell.number_format = copy(source_cell.number_format)
        target_cell.protection = copy(source_cell.protection)
        target_cell.alignment = copy(source_cell.alignment)

def clone_row_style(ws, source_row_idx, target_row_idx):
    for col in range(1, ws.max_column + 1):
        clone_cell_style(ws.cell(row=source_row_idx, column=col), ws.cell(row=target_row_idx, column=col))
    for merged_range in list(ws.merged_cells.ranges):
        min_col, min_row, max_col, max_row = range_boundaries(str(merged_range))
        if min_row == source_row_idx and max_row == source_row_idx:
            ws.merge_cells(start_row=target_row_idx, start_column=min_col, end_row=target_row_idx, end_column=max_col)

def shift_merged_ranges_down(ws, start_row, amount):
    if amount <= 0: return
    ranges_to_transform = []
    for merged_range in list(ws.merged_cells.ranges):
        min_col, min_row, max_col, max_row = range_boundaries(str(merged_range))
        if min_row > start_row:
            ws.merged_cells.remove(merged_range)
            ranges_to_transform.append((min_row + amount, min_col, max_row + amount, max_col))
        elif min_row <= start_row < max_row:
            ws.merged_cells.remove(merged_range)
            ranges_to_transform.append((min_row, min_col, max_row + amount, max_col))
    for min_r, min_c, max_r, max_c in ranges_to_transform:
        ws.merge_cells(start_row=min_r, start_column=min_c, end_row=max_r, end_column=max_c)

def replace_scalar_placeholders(ws, metadata):
    for row in ws.iter_rows():
        for cell in row:
            if cell.value and isinstance(cell.value, str):
                for key, val in metadata.items():
                    tag = f'{{{{{key}}}}}'
                    if tag in cell.value: cell.value = cell.value.replace(tag, str(val))

def populate_table_section(ws, anchor_tag, items, field_keys):
    anchor_row, anchor_col = None, None
    for r in range(1, ws.max_row + 1):
        for c in range(1, ws.max_column + 1):
            val = ws.cell(row=r, column=c).value
            if val and isinstance(val, str) and anchor_tag in val:
                anchor_row, anchor_col = r, c
                ws.cell(row=r, column=c).value = ''
                break
        if anchor_row: break

    if not anchor_row or not items: return
    num_items = len(items)
    if num_items > 1:
        shift_merged_ranges_down(ws, start_row=anchor_row, amount=num_items - 1)
        ws.insert_rows(anchor_row + 1, amount=num_items - 1)

    for i, item in enumerate(items):
        curr_row = anchor_row + i
        if i > 0: clone_row_style(ws, source_row_idx=anchor_row, target_row_idx=curr_row)
        for col_offset, key in enumerate(field_keys):
            cell = ws.cell(row=curr_row, column=anchor_col + col_offset)
            cell.value = (i + 1) if key == '_index' else item.get(key, '')

def generate_mom_excel(summary_data, template_path, output_path):
    wb = load_workbook(template_path)
    ws = wb.active
    preserve_template_logo(ws, template_path)
    if 'metadata' in summary_data: replace_scalar_placeholders(ws, summary_data['metadata'])
    if 'discussions' in summary_data:
        populate_table_section(ws, '{{#DISCUSSIONS}}', summary_data['discussions'], ['_index', 'details'])
    if 'action_items' in summary_data:
        populate_table_section(ws, '{{#ACTION_ITEMS}}', summary_data['action_items'], ['_index', 'task', 'owner', 'target_date', 'closure_date', 'status', 'remarks'])
    wb.save(output_path)
    return output_path

# ==============================================================================
# 6. TOOL IMPLEMENTATIONS
# ==============================================================================

def process_notes_with_gemini(raw_text: str) -> dict:
    prompt = f"""
    Read the following meeting notes. Extract info into this EXACT JSON structure suitable for both Word and Excel templates:
    {{
      "metadata": {{
         "date": "YYYY-MM-DD",
         "purpose": "Brief summary of meeting purpose",
         "prepared_by": "Author name or Team",
         "venue": "Location or online platform",
         "attendees_ymsli": "Names of YMSLI attendees",
         "attendees_ymesg": "Names of YMESG attendees"
      }},
      "discussions": [
         {{"details": "Discussion point 1"}},
         {{"details": "Discussion point 2"}}
      ],
      "action_items": [
         {{"task": "Task name", "owner": "Owner", "target_date": "YYYY-MM-DD", "closure_date": "", "status": "Pending", "remarks": ""}}
      ]
    }}
    Notes: {raw_text}
    """
    _enforce_global_pacing()
    response = client.models.generate_content(
        model=MODEL_NAME, contents=prompt,
        config=types.GenerateContentConfig(response_mime_type="application/json"),
    )
    return json.loads(response.text)


def summarize_and_generate_mom(file_id: str, file_name: str, project_name: str = "Project Summary", doc_format: str = "excel") -> dict:
    token = acquire_microsoft_graph_token()
    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content", headers={"Authorization": f"Bearer {token}"})
    raw_text = ParserRegistry.parse(res.content, file_name) if res.status_code == 200 else f"Download Error"
    meeting_data = process_notes_with_gemini(raw_text)
    
    if doc_format.lower() == "excel":
        template = fetch_template(DEFAULT_MOM_EXCEL_TEMPLATE)
        out_name = "Final_MOM_Output.xlsx"
        generate_mom_excel(meeting_data, template, out_name)
    else:
        template = fetch_template(DEFAULT_MOM_WORD_TEMPLATE)
        out_name = "final_minutes.docx"
        doc = DocxTemplate(template)
        context = {
            "project_title": project_name,
            "meeting_date": meeting_data.get("metadata", {}).get("date", "Date Not Specified"),
            "discussion_points": [d.get("details") for d in meeting_data.get("discussions", [])],
            "action_items": meeting_data.get("action_items", [])
        }
        doc.render(context)
        doc.save(out_name)
        
    return {"status": f"MOM rendered successfully in {doc_format} format!", "file_path": out_name}


def search_filename_and_content(query: str, find_latest_only: bool = False) -> dict:
    results = document_collection.query(query_texts=[query], n_results=5)
    if not results.get("metadatas") or not results["metadatas"][0]:
        return {"result": f"No documents found matching '{query}'."}

    files_map = {}
    for i, meta in enumerate(results["metadatas"][0]):
        f_id = meta["file_id"]
        if f_id not in files_map:
            files_map[f_id] = meta
            files_map[f_id]["score"] = 1
        else:
            files_map[f_id]["score"] += 1

    sorted_files = sorted(files_map.values(), key=lambda x: (x["score"], x["created_date"]), reverse=True)
    return {"status": "Found", "latest_file": sorted_files[0]} if find_latest_only else {"matching_documents": sorted_files[:5]}


def generate_version_history_chart(file_id: str, file_name: str = None) -> dict:
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    if not file_name:
        res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}", headers=headers)
        file_name = res.json().get("name", file_id)

    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/versions", headers=headers)
    versions = [{"version_id": v.get("id"), "modified_datetime": v.get("lastModifiedDateTime"), "size_bytes": v.get("size", 0), "modified_by": (v.get("lastModifiedBy", {}).get("user", {})).get("displayName", "Unknown")} for v in res.json().get("value", [])]
    versions.sort(key=lambda x: x["modified_datetime"])
    
    dates = [datetime.fromisoformat(v["modified_datetime"].replace("Z", "+00:00")) for v in versions]
    sizes = [v["size_bytes"] / 1024 for v in versions]
    
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.plot(dates, sizes, marker="o", color="#1a73e8")
    ax.set_title(f"Version History — {file_name}")
    fig.autofmt_xdate(rotation=30)
    out_path = f"version_history_{file_name.replace(' ', '_')}.png"
    fig.savefig(out_path)
    plt.close(fig)
    return {"chart_path": out_path, "versions": versions}


def update_file_content(file_id: str, new_content: str) -> dict:
    """Updates a text-based file on OneDrive (e.g., .txt, .csv) with new content, constrained by permissions."""
    result = save_modifications_to_onedrive(file_id, new_content.encode('utf-8'))
    return {"status": result}


def delete_file(file_id: str) -> dict:
    """Deletes a file from OneDrive if the user has owner or write access."""
    user_roles = get_user_permissions_for_file(file_id)
    if not any(r in ["write", "owner"] for r in user_roles):
        return {"error": "❌ Permission denied. You do not have sufficient rights to delete this file."}
        
    token = acquire_microsoft_graph_token()
    res = requests.delete(
        f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}", 
        headers={"Authorization": f"Bearer {token}"}
    )
    if res.status_code == 204:
        # Also clean up the vector DB manually if we successfully delete it in the backend
        try:
            document_collection.delete(where={"file_id": file_id})
            folder_collection.delete(ids=[file_id])
        except Exception:
            pass
        return {"status": "✅ File successfully deleted from OneDrive."}
    return {"error": f"⚠️ Failed to delete. HTTP {res.status_code}: {res.text}"}


TOOLS_MAP = {
    "search_filename_and_content": search_filename_and_content,
    "summarize_and_generate_mom": summarize_and_generate_mom,
    "generate_version_history_chart": generate_version_history_chart,
    "update_file_content": update_file_content,
    "delete_file": delete_file,
}

SYSTEM_PROMPT = """
You are a knowledgeable and action-oriented OneDrive assistant. 
1. If the user asks to generate MOM or summarize a document, use `search_filename_and_content` to find it, then `summarize_and_generate_mom`. You can pass doc_format="excel" or "word".
2. If asked for a file's history, use `generate_version_history_chart`.
3. If the user asks to UPDATE or EDIT a text file, find the file ID, then use `update_file_content`.
4. If the user asks to DELETE a file, find the file ID, then use `delete_file`.
Call one tool at a time. Do not invoke multiple search functions in a single batch.
"""

def send_message_and_handle_tools(chat, user_input: str):
    _enforce_global_pacing()
    response = chat.send_message(user_input)
    api_call_count = 1

    while response.function_calls and api_call_count < MAX_API_CALLS_PER_TURN:
        for fn_call in response.function_calls:
            fn_name, fn_args = fn_call.name, dict(fn_call.args)
            print(f"\n⚙️ [Executing Tool: {fn_name}]")
            try:
                tool_result = TOOLS_MAP[fn_name](**fn_args) if fn_name in TOOLS_MAP else {"error": "Not found"}
            except Exception as e:
                tool_result = {"error": str(e)}

            tool_part = types.Part.from_function_response(name=fn_name, response={"result": tool_result})
            _enforce_global_pacing()
            response = chat.send_message(tool_part)
            api_call_count += 1

    return response.text if hasattr(response, "text") else "Completed task."

# ==============================================================================
# 7. MAIN APPLICATION
# ==============================================================================

def run_assistant():
    print("==================================================")
    print(" 🔐 Application User Portal Login")
    print("==================================================")
    
    global USER_EMAIL
    USER_EMAIL = input("Enter your email address to establish permissions context: ").strip()
    
    if not USER_EMAIL:
        print("❌ Login failed: Email cannot be empty.")
        return

    print(f"\n🔍 Authenticating master session for: {USER_EMAIL}...")
    acquire_microsoft_graph_token()
    
    sync_onedrive_to_vector_db(force_resync=FORCE_RESYNC_ON_START)

    config = types.GenerateContentConfig(system_instruction=SYSTEM_PROMPT, tools=[
        search_filename_and_content, 
        summarize_and_generate_mom, 
        generate_version_history_chart,
        update_file_content,
        delete_file
    ])
    chat = client.chats.create(model=MODEL_NAME, config=config)

    print("\n==================================================")
    print(f"  OneDrive Knowledge, Edit & MOM Assistant ({MODEL_NAME})")
    print("==================================================")
    print("Type 'exit' to end session.\n")

    while True:
        try:
            user_input = input("\nUSER > ").strip()
            if not user_input: continue
            if user_input.lower() in ["exit", "quit"]: break
            print(f"\nAGENT > {send_message_and_handle_tools(chat, user_input)}\n{'-'*50}")
        except KeyboardInterrupt:
            break
        except Exception as e:
            print(f"\n[ERROR] {str(e)}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="OneDrive Knowledge Assistant")
    parser.add_argument("--version-history", nargs="?", const="", default=None)
    args, _ = parser.parse_known_args()

    if args.version_history is not None:
        file_id = args.version_history or input("Enter file ID: ")
        print(generate_version_history_chart(file_id))
    else:
        run_assistant()