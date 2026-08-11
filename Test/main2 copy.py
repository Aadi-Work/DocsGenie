import io
import json
import os
import random
import time
from datetime import datetime

# Prefer vanilla HF downloads instead of the Xet reconstruction path,
# which can fail in locked-down/offline environments with a 403 response.
os.environ.setdefault("HF_HUB_DISABLE_XET", "1")
os.environ.setdefault("HF_HUB_DISABLE_PROGRESS_BARS", "1")

import pandas as pd
import requests

# Document parsers
try:
    import docx  # type: ignore
    from docx import Document  # type: ignore
except Exception:
    docx = None
    Document = None

try:
    from docxtpl import DocxTemplate  # type: ignore
except Exception:
    DocxTemplate = None

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

# Environment detection for Google Colab
try:
    from IPython.display import HTML, display  # type: ignore
    from google.colab import files  # type: ignore
    COLAB_ENV = True
except Exception:
    HTML = None
    display = None
    files = None
    COLAB_ENV = False

try:
    import pypdf  # type: ignore
except Exception:
    pypdf = None

# Vector Database & Embeddings
try:
    import chromadb
    from chromadb.utils import embedding_functions
except Exception:
    chromadb = None
    embedding_functions = None

# Google GenAI SDK (v1 / google-genai)
try:
    from google import genai
    from google.genai import types
    from google.genai.errors import APIError
    GENAI_AVAILABLE = True
except Exception:
    genai = None
    types = None
    APIError = Exception
    GENAI_AVAILABLE = False

# ==============================================================================
# 1. CONFIGURATION
# ==============================================================================

CLIENT_ID = "2194a676-0580-4882-a21c-d7d4d0a8966c"

# Set to "common" for standard personal/multi-tenant accounts
TENANT_ID = "common"

GEMINI_API_KEY = ""

MODEL_NAME = "gemini-3.6-flash"
CHROMA_PATH = "./onedive_vector_db"
DELTA_TOKEN_FILE = "./onedrive_delta_token.json"
SCOPES = ["Files.Read.All", "User.Read"]

DEFAULT_MOM_TEMPLATE = "mom_template.docx"
CURRENT_ACCESS_TOKEN = None

# Set to True only if you want to manually wipe local vector cache and do a full resync
FORCE_RESYNC_ON_START = False

# Initialize Google GenAI Client when the SDK is available and a key exists
client = None


def configure_gemini_client(api_key: str | None = None):
    global client, GEMINI_API_KEY
    if api_key:
        GEMINI_API_KEY = api_key

    if GENAI_AVAILABLE and genai is not None and GEMINI_API_KEY and GEMINI_API_KEY != "AQ.Ab8RN6IqblEL1lXP034ndepILeQnz5jQlgjHr-1AuRrsaKWCiQ":
        try:
            client = genai.Client(api_key=GEMINI_API_KEY)
        except Exception as e:
            print(f"⚠️ Failed to initialize GenAI Client: {e}")
            client = None
    else:
        client = None
    return client


configure_gemini_client()

# ==============================================================================
# 2. HELPER FUNCTIONS
# ==============================================================================

def log_execution_time(action_name: str, start_time: float):
    """Logs real-time execution duration in seconds."""
    elapsed = time.perf_counter() - start_time
    print(f"⏱️ [{action_name}] completed in {elapsed:.3f} seconds.")


def resolve_available_model_name() -> str:
    """Resolves a Gemini model name against the SDK's live catalog when possible."""
    configured = MODEL_NAME or "gemini-2.5-flash"
    if client is None:
        return configured

    try:
        model_names = []
        model_list = list(client.models.list())
        for model in model_list:
            model_name = getattr(model, "name", "")
            if model_name:
                model_names.append(model_name)
        normalized = {name.replace("models/", "") for name in model_names}
        if configured.replace("models/", "") in normalized:
            return configured

        candidate_order = [
            "gemini-2.5-flash",
            "gemini-2.5-pro",
            "gemini-2.0-flash",
            "gemini-2.0-flash-lite",
            "gemini-1.5-flash",
        ]
        for candidate in candidate_order:
            if candidate in normalized:
                print(f"⚠️ Falling back from '{configured}' to supported Gemini model '{candidate}'.")
                return candidate

        if model_names:
            first_model = model_names[0].replace("models/", "")
            print(f"⚠️ Falling back from '{configured}' to first available Gemini model '{first_model}'.")
            return first_model
    except Exception as exc:
        print(f"⚠️ Gemini model catalog lookup failed: {exc}")

    return configured


def acquire_microsoft_graph_token() -> str:
    """Acquires a Microsoft Graph access token using OAuth Device Code flow."""
    global CURRENT_ACCESS_TOKEN
    if CURRENT_ACCESS_TOKEN:
        return CURRENT_ACCESS_TOKEN

    print("\n--- Initiating Microsoft Graph Authentication ---")
    device_code_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/devicecode"
    payload = {"client_id": CLIENT_ID, "scope": " ".join(SCOPES)}

    res = requests.post(device_code_url, data=payload)
    if res.status_code != 200:
        raise RuntimeError(f"Device Code Auth failed: {res.status_code} - {res.text}")

    data = res.json()
    verification_uri = data.get("verification_uri", "https://microsoft.com/devicelogin")
    user_code = data["user_code"]

    if COLAB_ENV:
        html_banner = f"""
        <div style="background-color: #f0f4f9; border-left: 6px solid #1a73e8; padding: 15px; margin: 10px 0; font-family: sans-serif;">
          <h3 style="margin-top: 0; color: #1a73e8;">🔑 Microsoft Authentication Required</h3>
          <p>1. Open Link: <a href="{verification_uri}" target="_blank" style="font-weight: bold; font-size: 16px;">{verification_uri}</a></p>
          <p>2. Enter Code: <span style="background-color: #ffffff; border: 1px solid #ccc; padding: 4px 8px; font-family: monospace; font-size: 18px; font-weight: bold;">{user_code}</span></p>
        </div>
        """
        display(HTML(html_banner))
    else:
        print(f"\n🔑 Open {verification_uri} and enter code: {user_code}\n")

    token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    token_payload = {
        "grant_type": "urn:ietf:params:oauth:grant-type:device_code",
        "client_id": CLIENT_ID,
        "device_code": data["device_code"],
    }

    start_time = time.time()
    interval = data.get("interval", 5)
    print("⏳ Waiting for browser sign-in completion...")

    while time.time() - start_time < data.get("expires_in", 900):
        time.sleep(interval)
        token_res = requests.post(token_url, data=token_payload)
        token_data = token_res.json()

        if "access_token" in token_data:
            print("✅ Access token acquired successfully!\n")
            CURRENT_ACCESS_TOKEN = token_data["access_token"]
            return CURRENT_ACCESS_TOKEN

        error_code = token_data.get("error")
        if error_code == "authorization_pending":
            continue
        elif error_code == "authorization_declined":
            raise RuntimeError("❌ User declined authentication in browser.")
        elif error_code == "expired_token":
            raise RuntimeError("❌ Device code expired. Please re-run the script.")
        else:
            raise RuntimeError(f"❌ OAuth Error ({error_code}): {token_data.get('error_description')}")

    raise TimeoutError("Authentication device code expired.")


class ParserRegistry:
    @staticmethod
    def truncate_text(text: str, max_chars: int = 4000) -> str:
        return text[:max_chars] + "\n\n...[Content truncated]..." if len(text) > max_chars else text

    @classmethod
    def parse(cls, content_bytes: bytes, file_name: str) -> str:
        ext = os.path.splitext(file_name)[1].lower()
        try:
            if ext in [".xlsx", ".xls", ".csv"]:
                xls = pd.ExcelFile(io.BytesIO(content_bytes))
                return "\n\n".join(
                    [
                        f"**Sheet: {s}**\n" + pd.read_excel(xls, sheet_name=s, nrows=20).to_markdown(index=False)
                        for s in xls.sheet_names[:3]
                    ]
                )
            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(io.BytesIO(content_bytes))
                slides = list(prs.slides)[:15]
                return cls.truncate_text(
                    "\n".join(
                        [
                            f"Slide {i+1}: " + " | ".join(
                                [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
                            )
                            for i, slide in enumerate(slides)
                        ]
                    )
                )
            elif ext in [".docx", ".doc"]:
                doc = Document(io.BytesIO(content_bytes))
                return cls.truncate_text("\n".join([p.text for p in doc.paragraphs if p.text.strip()]))
            elif ext == ".pdf" and pypdf:
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                return cls.truncate_text(
                    "\n".join([p.extract_text() for p in reader.pages[:10] if p.extract_text()])
                )
            else:
                return cls.truncate_text(content_bytes.decode("utf-8", errors="ignore"))
        except Exception as e:
            return f"Parsing error for {file_name}: {str(e)}"


# ==============================================================================
# 3. TEMPLATE RESOLUTION & ONEDRIVE AUTO-FETCH
# ==============================================================================

def ensure_mom_template_exists(template_path: str = DEFAULT_MOM_TEMPLATE) -> str:
    """Finds template or auto-creates a valid default template if missing completely."""
    if os.path.exists(template_path):
        return template_path

    cwd = os.getcwd()
    for f in os.listdir(cwd):
        if "mom_template" in f.lower() and f.endswith(".docx"):
            resolved_path = os.path.join(cwd, f)
            print(f"🔍 Located local template file: '{f}'")
            return resolved_path

    print(f"⚠️ Creating default local template '{template_path}'...")
    doc = Document()
    doc.add_heading("{{ project_title }} - Meeting Minutes", level=1)
    doc.add_paragraph("Date: {{ meeting_date }}")

    doc.add_heading("Discussion Points", level=2)
    doc.add_paragraph("{% for point in discussion_points %}")
    doc.add_paragraph("• {{ point }}")
    doc.add_paragraph("{% endfor %}")

    doc.add_heading("Action Items", level=2)
    doc.add_paragraph("{% for item in action_items %}")
    doc.add_paragraph("• Task: {{ item.task }} | Owner: {{ item.owner }} | Status: {{ item.status }}")
    doc.add_paragraph("{% endfor %}")

    doc.save(template_path)
    print(f"✅ Created baseline MOM template at '{template_path}'")
    return template_path


def fetch_template_from_onedrive(template_name: str = DEFAULT_MOM_TEMPLATE) -> str:
    """Checks if template exists locally. If not, automatically searches & downloads it from OneDrive."""
    if os.path.exists(template_name):
        return template_name

    print(f"📡 '{template_name}' not found on local disk. Searching OneDrive cloud...")
    try:
        token = acquire_microsoft_graph_token()
        headers = {"Authorization": f"Bearer {token}"}
        search_url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{template_name}')"
        res = requests.get(search_url, headers=headers)

        if res.status_code == 200:
            items = res.json().get("value", [])
            for item in items:
                if item.get("name", "").lower() == template_name.lower():
                    file_id = item["id"]
                    download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
                    file_res = requests.get(download_url, headers=headers)
                    if file_res.status_code == 200:
                        with open(template_name, "wb") as f:
                            f.write(file_res.content)
                        print(f"✅ Successfully downloaded '{template_name}' from OneDrive into local workspace!")
                        return template_name
    except Exception as e:
        print(f"⚠️ Cloud search error: {str(e)}")

    print(f"⚠️ Could not fetch '{template_name}' from OneDrive. Falling back to default generation...")
    return ensure_mom_template_exists(template_name)


# ==============================================================================
# 4. VECTOR DB & INCREMENTAL ONEDRIVE SYNC
# ==============================================================================

embedding_fn = None
chroma_client = None
folder_collection = None
document_collection = None

if embedding_functions is not None and chromadb is not None:
    try:
        embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
        chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)
        folder_collection = chroma_client.get_or_create_collection(
            name="folder_paths", embedding_function=embedding_fn
        )
        document_collection = chroma_client.get_or_create_collection(
            name="document_chunks", embedding_function=embedding_fn
        )
    except Exception as exc:
        print(f"⚠️ Chroma/SentenceTransformer initialization failed: {exc}. Vector search disabled.")
        embedding_fn = None
        chroma_client = None
        folder_collection = None
        document_collection = None


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks if chunks else [text]


def sync_onedrive_to_vector_db(force_resync: bool = False):
    """Syncs OneDrive to ChromaDB with 410 auto-recovery and verbose real-time logging."""
    global folder_collection, document_collection
    sync_start = time.perf_counter()

    if chroma_client is None or folder_collection is None or document_collection is None:
        print("⚠️ ChromaDB collections are unavailable. OneDrive sync skipped.")
        return

    if force_resync:
        print("🧹 Clearing local vector database and delta tokens for full resync...")
        try:
            chroma_client.delete_collection("folder_paths")
            chroma_client.delete_collection("document_chunks")
            if os.path.exists(DELTA_TOKEN_FILE):
                os.remove(DELTA_TOKEN_FILE)
        except Exception:
            pass

        folder_collection = chroma_client.create_collection(name="folder_paths", embedding_function=embedding_fn)
        document_collection = chroma_client.create_collection(name="document_chunks", embedding_function=embedding_fn)

    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}

    if os.path.exists(DELTA_TOKEN_FILE) and not force_resync:
        with open(DELTA_TOKEN_FILE, "r") as f:
            next_url = json.load(f).get("delta_link")
        print("⚡ Checking OneDrive for NEW, MODIFIED, or DELETED files/folders...\n")
    else:
        next_url = "https://graph.microsoft.com/v1.0/me/drive/root/delta?$select=id,name,createdDateTime,lastModifiedDateTime,folder,file,parentReference,deleted"
        print("🔄 Performing initial OneDrive index...\n")

    delta_link = None
    processed_count = 0
    deleted_count = 0

    while next_url:
        res = requests.get(next_url, headers=headers)

        if res.status_code == 410 or "resyncRequired" in res.text:
            print("\n⚠️ Stale delta token detected (410 resyncRequired). Auto-recovering via full resync...")
            if os.path.exists(DELTA_TOKEN_FILE):
                os.remove(DELTA_TOKEN_FILE)
            return sync_onedrive_to_vector_db(force_resync=True)

        if res.status_code != 200:
            print(f"⚠️ Delta API Error ({res.status_code}): {res.text}")
            break

        data = res.json()
        items = data.get("value", [])

        for item in items:
            item_id = item["id"]
            item_name = item.get("name", "Unknown")

            if "deleted" in item:
                try:
                    document_collection.delete(where={"file_id": item_id})
                    folder_collection.delete(ids=[item_id])
                    deleted_count += 1
                    print(f"  🗑️ Removed Deleted Item: {item_name}")
                except Exception:
                    pass
                continue

            if "folder" in item:
                folder_collection.upsert(
                    ids=[item_id],
                    documents=[f"Folder Name: {item_name}"],
                    metadatas=[{"folder_name": item_name, "full_path": item_name}],
                )
                print(f"  📁 Indexed Folder: {item_name}")

            elif "file" in item:
                download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content"
                file_res = requests.get(download_url, headers=headers)

                if file_res.status_code == 200:
                    parsed_text = ParserRegistry.parse(file_res.content, item_name)
                    chunks = chunk_text(parsed_text)
                    chunk_ids = [f"{item_id}_chunk_{i}" for i in range(len(chunks))]
                    metadatas = [
                        {
                            "file_id": item_id,
                            "file_name": item_name,
                            "created_date": item.get("createdDateTime", "1970-01-01T00:00:00Z"),
                            "chunk_idx": i,
                        }
                        for i in range(len(chunks))
                    ]
                    document_collection.upsert(ids=chunk_ids, documents=chunks, metadatas=metadatas)
                    processed_count += 1
                    print(f"  📄 Indexed/Updated File: {item_name} ({len(chunks)} text chunks)")

        if "@odata.nextLink" in data:
            next_url = data["@odata.nextLink"]
        elif "@odata.deltaLink" in data:
            delta_link = data["@odata.deltaLink"]
            next_url = None

    if delta_link:
        with open(DELTA_TOKEN_FILE, "w") as f:
            json.dump({"delta_link": delta_link}, f)

    print(f"\n✅ Sync finished! (Indexed/Updated: {processed_count} files | Removed: {deleted_count} items)")
    log_execution_time("OneDrive Delta Sync Total", sync_start)


# ==============================================================================
# 5. STRUCTURED GEMINI JSON EXTRACTION & AGENT TOOLS
# ==============================================================================

def process_notes_with_gemini(raw_text: str) -> dict:
    """Uses Gemini JSON mode to extract date, discussion points and action items."""
    print("🧠 Extracting structured meeting summary with Gemini JSON mode...")
    start_time = time.perf_counter()

    prompt = f"""
    Read the following meeting notes. Do three things:
    1. Extract the date of the meeting (or 'Date Not Specified' if absent).
    2. Provide a brief summary of the discussion points as a list of strings.
    3. Extract the action items.

    Return the result EXCLUSIVELY as a valid JSON object with this exact structure:
    {{
      "meeting_date": "YYYY-MM-DD or formatted date string",
      "discussion_points": ["Point 1", "Point 2"],
      "action_items": [
        {{"task": "Task name", "owner": "Name or Team", "status": "Pending/In Progress/Done"}}
      ]
    }}

    Meeting Notes:
    {raw_text}
    """

    if client is None or types is None:
        summary = {
            "meeting_date": "Date Not Specified",
            "discussion_points": ["Meeting notes processed locally (GenAI disabled)."],
            "action_items": [],
        }
        log_execution_time("Local Fallback Extraction", start_time)
        return summary

    response = client.models.generate_content(
        model=MODEL_NAME,
        contents=prompt,
        config=types.GenerateContentConfig(response_mime_type="application/json"),
    )

    log_execution_time("Gemini JSON Extraction", start_time)
    return json.loads(response.text)


def generate_meeting_minutes(
    raw_notes_text: str,
    project_name: str,
    date: str = None,
    template_path: str = DEFAULT_MOM_TEMPLATE,
    output_path: str = "final_minutes.docx",
) -> dict:
    """Generates a structured Word document using docxtpl and Gemini JSON extraction."""
    start_time = time.perf_counter()
    resolved_template_path = fetch_template_from_onedrive(template_path)
    meeting_data = process_notes_with_gemini(raw_notes_text)
    
    final_date = date if date else meeting_data.get("meeting_date", "Date Not Specified")

    print(f"📝 Rendering Word document with docxtpl using '{resolved_template_path}'...")
    if DocxTemplate is None:
        raise RuntimeError("docxtpl is not installed. Install it before generating the MOM.")

    doc = DocxTemplate(resolved_template_path)
    context = {
        "project_title": project_name,
        "meeting_date": final_date,
        "discussion_points": meeting_data.get("discussion_points", []),
        "action_items": meeting_data.get("action_items", []),
    }
    doc.render(context)
    doc.save(output_path)

    if COLAB_ENV:
        try:
            files.download(output_path)
            download_status = "Triggered browser file download."
        except Exception as e:
            download_status = f"Saved locally. Download note: {str(e)}"
    else:
        download_status = f"Saved locally at {output_path}"

    log_execution_time("Full MOM Generation Pipeline", start_time)
    return {
        "status": "MOM rendered successfully!",
        "file_path": output_path,
        "download_status": download_status,
        "extracted_summary": meeting_data,
    }


def find_folder_path(query: str) -> dict:
    if folder_collection is None:
        return {"result": "Vector DB is not initialized. Re-run the assistant after dependency setup."}
    start_time = time.perf_counter()
    results = folder_collection.query(query_texts=[query], n_results=3)
    log_execution_time("Folder Vector Search", start_time)

    metadata = results.get("metadatas") or []
    if not metadata or not metadata[0]:
        return {"result": "No matching folder paths found."}

    matches = []
    for m in metadata[0]:
        if isinstance(m, dict):
            matches.append({"folder_name": m.get("folder_name"), "full_path": m.get("full_path")})
    return {"matching_folders": matches} if matches else {"result": "No matching folder paths found."}


def search_content_only(query: str) -> dict:
    if document_collection is None:
        return {"result": "Vector DB is not initialized. Re-run the assistant after dependency setup."}
    start_time = time.perf_counter()
    results = document_collection.query(query_texts=[query], n_results=5)
    log_execution_time("Content-Only Vector Search", start_time)

    metadata = results.get("metadatas") or []
    documents = results.get("documents") or []
    if not metadata or not metadata[0] or not documents or not documents[0]:
        return {"result": f"No content matches found for '{query}'."}

    matches = [
        {"file_name": meta.get("file_name"), "file_id": meta.get("file_id"), "text_snippet": documents[0][i]}
        for i, meta in enumerate(metadata[0]) if isinstance(meta, dict)
    ]
    return {"status": "Content-only search results", "matching_documents": matches}


def search_filename_and_content(query: str, find_latest_only: bool = False) -> dict:
    if document_collection is None:
        return {"result": "Vector DB is not initialized. Re-run the assistant after dependency setup."}
    start_time = time.perf_counter()
    results = document_collection.query(query_texts=[query], n_results=5)
    log_execution_time("ChromaDB HNSW Vector Search", start_time)

    metadata = results.get("metadatas") or []
    documents = results.get("documents") or []
    if not metadata or not metadata[0] or not documents or not documents[0]:
        return {"result": f"No documents found matching '{query}'."}

    query_keywords = [k.lower() for k in query.split() if len(k) > 2]
    files_map = {}

    for i, meta in enumerate(metadata[0]):
        if not isinstance(meta, dict):
            continue
        f_id = meta.get("file_id")
        if not f_id:
            continue
        file_name = str(meta.get("file_name", "")).lower()
        filename_score = sum(3 for kw in query_keywords if kw in file_name)

        if f_id not in files_map:
            files_map[f_id] = {
                "file_id": f_id,
                "file_name": meta.get("file_name", "Unknown"),
                "created_date": meta.get("created_date", "1970-01-01T00:00:00Z"),
                "score": filename_score,
                "snippet": documents[0][i],
            }

    if find_latest_only:
        sorted_files = sorted(files_map.values(), key=lambda x: (x["score"], x["created_date"]), reverse=True)
        return {"status": "Found latest document", "latest_file": sorted_files[0] if sorted_files else None}

    sorted_files = sorted(files_map.values(), key=lambda x: x["score"], reverse=True)
    return {"status": "Hybrid search results", "matching_documents": sorted_files[:5]}


def read_file_content(file_id: str, file_name: str) -> str:
    start_time = time.perf_counter()
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content", headers=headers)
    parsed_content = ParserRegistry.parse(res.content, file_name) if res.status_code == 200 else f"Download Error: {res.status_code}"
    log_execution_time("OneDrive File Fetch & Parse", start_time)
    return parsed_content


def summarize_and_generate_mom(file_id: str, file_name: str, project_name: str, date: str = None) -> dict:
    raw_text = read_file_content(file_id, file_name)
    return generate_meeting_minutes(raw_notes_text=raw_text, project_name=project_name, date=date)


TOOLS_MAP = {
    "find_folder_path": find_folder_path,
    "search_content_only": search_content_only,
    "search_filename_and_content": search_filename_and_content,
    "read_file_content": read_file_content,
    "summarize_and_generate_mom": summarize_and_generate_mom,
}

# ==============================================================================
# 6. AGENT SYSTEM PROMPT & AUTOMATED CHAT LOOP
# ==============================================================================

SYSTEM_PROMPT = """
You are an onboarding and knowledge management assistant for Microsoft OneDrive.

SEARCH ROUTING RULES:
1. If user asks specifically to search inside file text or content:
   --> Use `search_content_only(query)`.
2. If user asks for a document by name or topic:
   --> Use `search_filename_and_content(query, find_latest_only=False)`.
3. If user asks for the "latest" or "newest" MOM/document on a topic:
   --> Use `search_filename_and_content(query, find_latest_only=True)`.
   --> Present the identified file details and ask: "Would you like me to summarize this document and generate a Word document?"

WHEN USER CONFIRMS "YES":
1. Call `summarize_and_generate_mom(file_id, file_name, project_name, date)`.
2. Present the extracted data and confirm output path.
"""

def send_with_429_retry(chat, payload, max_retries: int = 5, pacing_seconds: float = 1.5):
    for attempt in range(max_retries):
        try:
            time.sleep(pacing_seconds)
            return chat.send_message(payload)
        except APIError as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                wait_time = (2 ** (attempt + 1) * 2) + random.uniform(0.5, 1.5)
                print(f"\n⚠️ [RATE LIMIT 429] Retrying in {wait_time:.1f}s (Attempt {attempt+1}/{max_retries})...")
                time.sleep(wait_time)
            else:
                raise e
    raise RuntimeError("Failed after max retries due to persistent rate limits.")


def send_message_and_handle_tools(chat, user_input: str):
    if chat is None:
        return "Gemini chat is unavailable. Check your API key and dependencies."

    response = send_with_429_retry(chat, user_input)

    while response.function_calls:
        for fn_call in response.function_calls:
            fn_name = fn_call.name
            fn_args = dict(fn_call.args)
            print(f"\n⚙️ [Executing Tool: {fn_name}]")

            if fn_name in TOOLS_MAP:
                try:
                    tool_result = TOOLS_MAP[fn_name](**fn_args)
                except Exception as e:
                    tool_result = {"error": str(e)}
            else:
                tool_result = {"error": f"Tool '{fn_name}' not found."}

            tool_part = types.Part.from_function_response(name=fn_name, response={"result": tool_result})
            response = send_with_429_retry(chat, tool_part)

    if hasattr(response, "text") and response.text:
        return response.text
    try:
        parts = response.candidates[0].content.parts
        texts = [p.text for p in parts if hasattr(p, "text") and p.text]
        return "\n".join(texts) if texts else "Completed task."
    except Exception:
        return "Completed task."


def run_assistant():
    global MODEL_NAME
    
    # Validation updated: only check if the API key wasn't changed from the default string.
    if GEMINI_API_KEY == "AQ.Ab8RN6IqblEL1lXP034ndepILeQnz5jQlgjHr-1AuRrsaKWCiQ":
        print("\n❌ CRITICAL ERROR: API Key is missing. Please paste your key into the code.")
        return

    MODEL_NAME = resolve_available_model_name()

    if chromadb is not None and embedding_functions is not None:
        sync_onedrive_to_vector_db(force_resync=FORCE_RESYNC_ON_START)
    else:
        print("⚠️ ChromaDB stack is unavailable. Local vector sync is skipped.")

    if client is None or types is None:
        print("⚠️ Google GenAI client is unavailable. The assistant will operate in degraded mode.")

    chat = None
    if client is not None and types is not None:
        config = types.GenerateContentConfig(
            system_instruction=SYSTEM_PROMPT,
            tools=[
                find_folder_path,
                search_content_only,
                search_filename_and_content,
                read_file_content,
                summarize_and_generate_mom
            ],
        )
        try:
            chat = client.chats.create(model=MODEL_NAME, config=config)
        except Exception as exc:
            print(f"⚠️ Chat creation with model '{MODEL_NAME}' failed: {exc}")
            try:
                available_models = list(client.models.list())
                if available_models:
                    fallback = available_models[0].name.replace("models/", "")
                    print(f"⚠️ Falling back to API-reported model: {fallback}")
                    MODEL_NAME = fallback
                    chat = client.chats.create(model=MODEL_NAME, config=config)
            except Exception as low_exc:
                print(f"⚠️ Final fallback initialization failed: {low_exc}")

    print("\n==================================================")
    print(f"  OneDrive Knowledge & MOM Assistant ({MODEL_NAME})")
    print("==================================================")
    print("Type 'exit' or 'quit' to end session.\n")

    while True:
        try:
            user_input = input("\nUSER > ").strip()
            if not user_input:
                continue
            if user_input.lower() in ["exit", "quit"]:
                print("Session ended. Goodbye!")
                break

            response_text = send_message_and_handle_tools(chat, user_input)
            print(f"\nAGENT > {response_text}")
            print("-" * 50)

        except KeyboardInterrupt:
            print("\nSession cancelled.")
            break
        except Exception as e:
            print(f"\n[ERROR] {str(e)}")


if __name__ == "__main__":
    run_assistant()