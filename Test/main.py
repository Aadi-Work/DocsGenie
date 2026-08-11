import argparse
import io
import json
import os
import random
import sys
import time
from datetime import datetime

import pandas as pd
import requests

# Document parsers
from docx import Document
from docxtpl import DocxTemplate
from pptx import Presentation

# Charting (non-interactive backend — must be set before importing pyplot)
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# Vector Database & Embeddings
import chromadb
from chromadb.utils import embedding_functions

# Google GenAI SDK (google-genai)
from google import genai
from google.genai import types
from google.genai.errors import APIError

# Environment detection for Google Colab vs Local
try:
    from IPython.display import HTML, display
    from google.colab import files, userdata

    COLAB_ENV = True
except ImportError:
    COLAB_ENV = False

try:
    import pypdf
except ImportError:
    pypdf = None

# ==============================================================================
# 1. CONFIGURATION & CLIENT INITIALIZATION
# ==============================================================================

# SECURITY: the API key must NOT be hardcoded in source. It is read from the
# GEMINI_API_KEY environment variable. Set it before running, e.g.:
#   export GEMINI_API_KEY="your-key-here"        (macOS/Linux)
#   $env:GEMINI_API_KEY="your-key-here"           (Windows PowerShell)
# If a real key was ever committed to source control or shared in plaintext,
# treat it as compromised and rotate it in the Google AI Studio / Cloud console.


# GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
GEMINI_API_KEY = ""

# In Colab, fall back to the userdata secrets store if the env var isn't set.
if not GEMINI_API_KEY and COLAB_ENV:
    try:
        GEMINI_API_KEY = userdata.get("GEMINI_API_KEY")
    except Exception:
        GEMINI_API_KEY = None

CLIENT_ID = "2194a676-0580-4882-a21c-d7d4d0a8966c"
TENANT_ID = "common"
SCOPES = ["Files.Read.All", "User.Read"]

# Standard valid Gemini Model string
MODEL_NAME = "gemini-3.6-flash"

CHROMA_PATH = "./onedrive_vector_db"
DELTA_TOKEN_FILE = "./onedrive_delta_token.json"

DEFAULT_MOM_TEMPLATE = "mom_template.docx"
CURRENT_ACCESS_TOKEN = None

# Set to True if you want to manually wipe local vector cache and perform a full resync
FORCE_RESYNC_ON_START = False

# Hard cap on Gemini API calls (chat.send_message invocations) per single user turn.
# Prevents a runaway tool-calling loop from generating unbounded API spend.
MAX_API_CALLS_PER_TURN = 3

# Minimum seconds enforced between ANY two Gemini API calls, session-wide (not just
# within one turn). Set to match a 3 requests/minute quota (60s / 3 = 20s), with a
# small safety margin so we don't ride the exact edge of the quota window.
MIN_SECONDS_BETWEEN_CALLS = 23.0
_last_gemini_call_time = 0.0

# Initialize Google GenAI Client only once we know we have a key.
if not GEMINI_API_KEY:
    raise RuntimeError(
        "GEMINI_API_KEY environment variable is not set. "
        "Set it before running this script (see comment above)."
    )
client = genai.Client(api_key=GEMINI_API_KEY)


def _enforce_global_pacing():
    """Blocks until at least MIN_SECONDS_BETWEEN_CALLS has passed since the last
    Gemini API call, regardless of which function is making it."""
    global _last_gemini_call_time
    elapsed = time.perf_counter() - _last_gemini_call_time
    remaining = MIN_SECONDS_BETWEEN_CALLS - elapsed
    if remaining > 0:
        time.sleep(remaining)
    _last_gemini_call_time = time.perf_counter()


def log_execution_time(action_name: str, start_time: float):
    """Logs real-time execution duration in seconds."""
    elapsed = time.perf_counter() - start_time
    print(f"⏱️ [{action_name}] completed in {elapsed:.3f} seconds.")


# ==============================================================================
# 2. MICROSOFT GRAPH AUTHENTICATION (device code flow — single shared copy)
# ==============================================================================


def acquire_microsoft_graph_token() -> str:
    """Acquires a Microsoft Graph access token using OAuth Device Code flow."""
    global CURRENT_ACCESS_TOKEN
    if CURRENT_ACCESS_TOKEN:
        return CURRENT_ACCESS_TOKEN

    print("\n--- Initiating Microsoft Graph Authentication ---")
    device_code_url = (
        f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/devicecode"
    )
    payload = {"client_id": CLIENT_ID, "scope": " ".join(SCOPES)}

    res = requests.post(device_code_url, data=payload)
    if res.status_code != 200:
        raise RuntimeError(
            f"Device Code Auth failed: {res.status_code} - {res.text}"
        )

    data = res.json()
    verification_uri = data.get(
        "verification_uri", "https://microsoft.com/devicelogin"
    )
    user_code = data["user_code"]

    if COLAB_ENV:
        html_banner = f"""
        <div style="background-color: #f0f4f9; border-left: 6px solid #1a73e8; padding: 15px; margin: 10px 0; font-family: sans-serif;">
          <h3 style="margin-top: 0; color: #1a73e8;">🔑 Microsoft Authentication Required</h3>
          <p>1. Open Link: <a href="{verification_uri}" target="_blank" style="font-weight: bold; font-size: 16px;">{verification_uri}</a></p>
          <p>2. Enter Code: <span style="background-color: #ffffff; border: 1px solid #ccc; padding: 4px 8px; font-family: monospace; font-size: 18px; font-weight: bold;">{user_code}</span></p>
        </div>
        """
        display(HTML(html_banner))
    else:
        print(f"\n🔑 Open {verification_uri} and enter code: {user_code}\n")

    token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    token_payload = {
        "grant_type": "urn:ietf:params:oauth:grant-type:device_code",
        "client_id": CLIENT_ID,
        "device_code": data["device_code"],
    }

    start_time = time.time()
    interval = data.get("interval", 5)
    print("⏳ Waiting for browser sign-in completion...")

    while time.time() - start_time < data.get("expires_in", 900):
        time.sleep(interval)
        token_res = requests.post(token_url, data=token_payload)
        token_data = token_res.json()

        if "access_token" in token_data:
            print("✅ Access token acquired successfully!\n")
            CURRENT_ACCESS_TOKEN = token_data["access_token"]
            return CURRENT_ACCESS_TOKEN

        error_code = token_data.get("error")
        if error_code == "authorization_pending":
            continue
        elif error_code == "authorization_declined":
            raise RuntimeError("❌ User declined authentication in browser.")
        elif error_code == "expired_token":
            raise RuntimeError("❌ Device code expired. Please re-run the script.")
        else:
            raise RuntimeError(
                f"❌ OAuth Error ({error_code}): {token_data.get('error_description')}"
            )

    raise TimeoutError("Authentication device code expired.")


# ==============================================================================
# 3. DOCUMENT PARSING
# ==============================================================================


class ParserRegistry:
    """Extracts plain text from Office docs, PDFs, spreadsheets, and raw text files."""

    @staticmethod
    def truncate_text(text: str, max_chars: int = 4000) -> str:
        return (
            text[:max_chars] + "\n\n...[Content truncated]..."
            if len(text) > max_chars
            else text
        )

    @classmethod
    def parse(cls, content_bytes: bytes, file_name: str) -> str:
        ext = os.path.splitext(file_name)[1].lower()
        try:
            if ext in [".xlsx", ".xls", ".csv"]:
                xls = pd.ExcelFile(io.BytesIO(content_bytes))
                return "\n\n".join(
                    [
                        f"**Sheet: {s}**\n"
                        + pd.read_excel(xls, sheet_name=s, nrows=20).to_markdown(
                            index=False
                        )
                        for s in xls.sheet_names[:3]
                    ]
                )
            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(io.BytesIO(content_bytes))
                # NOTE: python-pptx's `Slides` collection does not support slicing
                # (only integer __getitem__), so `prs.slides[:15]` raises TypeError.
                # Materialize to a list first, then slice.
                slides = list(prs.slides)[:15]
                return cls.truncate_text(
                    "\n".join(
                        [
                            f"Slide {i+1}: "
                            + " | ".join(
                                [
                                    s.text_frame.text.strip()
                                    for s in slide.shapes
                                    if s.has_text_frame
                                ]
                            )
                            for i, slide in enumerate(slides)
                        ]
                    )
                )
            elif ext in [".docx", ".doc"]:
                doc = Document(io.BytesIO(content_bytes))
                return cls.truncate_text(
                    "\n".join(
                        [p.text for p in doc.paragraphs if p.text.strip()]
                    )
                )
            elif ext == ".pdf" and pypdf:
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                pages = list(reader.pages)[:10]
                return cls.truncate_text(
                    "\n".join(
                        [p.extract_text() for p in pages if p.extract_text()]
                    )
                )
            else:
                return cls.truncate_text(
                    content_bytes.decode("utf-8", errors="ignore")
                )
        except Exception as e:
            return f"Parsing error for {file_name}: {str(e)}"


# ==============================================================================
# 4. VECTOR DB INITIALIZATION
# ==============================================================================

embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="all-MiniLM-L6-v2"
)
chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)

folder_collection = chroma_client.get_or_create_collection(
    name="folder_paths", embedding_function=embedding_fn
)
document_collection = chroma_client.get_or_create_collection(
    name="document_chunks", embedding_function=embedding_fn
)

# ==============================================================================
# 5. THREE-TIER TEMPLATE RESOLUTION ARCHITECTURE
# ==============================================================================


def fetch_template(template_name: str = DEFAULT_MOM_TEMPLATE) -> str:
    """3-TIER ARCHITECTURE FOR TEMPLATE RESOLUTION:
    Tier 1: Check Local Disk Workspace.
    Tier 2: Query ChromaDB Vector Index for metadata/file ID.
    Tier 3: Query OneDrive Cloud Graph API search endpoint directly.
    """

    # Tier 1: Local Workspace
    if os.path.exists(template_name):
        print(f"✅ [Tier 1] Template '{template_name}' found on local disk.")
        return template_name

    # Tier 2: ChromaDB Vector Index
    print(f"🔍 [Tier 2] Searching ChromaDB index for template '{template_name}'...")
    try:
        results = document_collection.get(where={"file_name": template_name})
        if results and results.get("metadatas") and len(results["metadatas"]) > 0:
            file_id = results["metadatas"][0]["file_id"]
            print(
                f"⚡ Found template in ChromaDB Index (File ID: {file_id}). Fetching from OneDrive..."
            )
            token = acquire_microsoft_graph_token()
            headers = {"Authorization": f"Bearer {token}"}
            download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
            res = requests.get(download_url, headers=headers)
            if res.status_code == 200:
                with open(template_name, "wb") as f:
                    f.write(res.content)
                print(f"✅ Downloaded '{template_name}' via ChromaDB metadata lookup!")
                return template_name
    except Exception as e:
        print(f"⚠️ Tier 2 lookup notice: {str(e)}")

    # Tier 3: OneDrive Cloud Graph API Direct Search
    print(f"📡 [Tier 3] Searching OneDrive Graph API directly for '{template_name}'...")
    try:
        token = acquire_microsoft_graph_token()
        headers = {"Authorization": f"Bearer {token}"}
        search_url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{template_name}')"
        res = requests.get(search_url, headers=headers)

        if res.status_code == 200:
            items = res.json().get("value", [])
            for item in items:
                if item.get("name", "").lower() == template_name.lower():
                    file_id = item["id"]
                    download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
                    file_res = requests.get(download_url, headers=headers)
                    if file_res.status_code == 200:
                        with open(template_name, "wb") as f:
                            f.write(file_res.content)
                        print(f"✅ Downloaded '{template_name}' from OneDrive Cloud search!")
                        return template_name
    except Exception as e:
        print(f"⚠️ Tier 3 cloud search notice: {str(e)}")

    raise FileNotFoundError(
        f"❌ Template '{template_name}' could not be resolved in Tier 1 (Local), Tier 2 (ChromaDB), or Tier 3 (OneDrive)."
    )


# ==============================================================================
# 6. INCREMENTAL ONEDRIVE SYNC
# ==============================================================================


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks if chunks else [text]


def sync_onedrive_to_vector_db(force_resync: bool = False):
    """Syncs OneDrive to ChromaDB using Graph Delta API with auto 410 recovery."""
    global folder_collection, document_collection
    sync_start = time.perf_counter()

    if force_resync:
        print("🧹 Clearing local vector database and delta tokens for full resync...")
        try:
            chroma_client.delete_collection("folder_paths")
            chroma_client.delete_collection("document_chunks")
            if os.path.exists(DELTA_TOKEN_FILE):
                os.remove(DELTA_TOKEN_FILE)
        except Exception:
            pass

        folder_collection = chroma_client.create_collection(
            name="folder_paths", embedding_function=embedding_fn
        )
        document_collection = chroma_client.create_collection(
            name="document_chunks", embedding_function=embedding_fn
        )

    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}

    if os.path.exists(DELTA_TOKEN_FILE) and not force_resync:
        with open(DELTA_TOKEN_FILE, "r") as f:
            next_url = json.load(f).get("delta_link")
        print("⚡ Checking OneDrive for NEW, MODIFIED, or DELETED files/folders...\n")
    else:
        next_url = (
            "https://graph.microsoft.com/v1.0/me/drive/root/delta"
            "?$select=id,name,createdDateTime,lastModifiedDateTime,folder,file,parentReference,deleted"
        )
        print("🔄 Performing initial OneDrive index...\n")

    delta_link = None
    processed_count = 0
    deleted_count = 0

    while next_url:
        res = requests.get(next_url, headers=headers)

        if res.status_code == 410 or "resyncRequired" in res.text:
            print(
                "\n⚠️ Stale delta token detected (410 resyncRequired). Auto-recovering via full resync..."
            )
            if os.path.exists(DELTA_TOKEN_FILE):
                os.remove(DELTA_TOKEN_FILE)
            return sync_onedrive_to_vector_db(force_resync=True)

        if res.status_code != 200:
            raise RuntimeError(f"❌ Graph Delta API Error ({res.status_code}): {res.text}")

        data = res.json()
        items = data.get("value", [])

        for item in items:
            item_id = item["id"]
            item_name = item.get("name", "Unknown")

            if "deleted" in item:
                try:
                    document_collection.delete(where={"file_id": item_id})
                    folder_collection.delete(ids=[item_id])
                    deleted_count += 1
                    print(f"  🗑️ Removed Deleted Item: {item_name}")
                except Exception:
                    pass
                continue

            if "folder" in item:
                folder_collection.upsert(
                    ids=[item_id],
                    documents=[f"Folder Name: {item_name}"],
                    metadatas=[{"folder_name": item_name, "full_path": item_name}],
                )
                print(f"  📁 Indexed Folder: {item_name}")

            elif "file" in item:
                download_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content"
                file_res = requests.get(download_url, headers=headers)

                if file_res.status_code == 200:
                    parsed_text = ParserRegistry.parse(file_res.content, item_name)
                    chunks = chunk_text(parsed_text)

                    chunk_ids = [f"{item_id}_chunk_{i}" for i in range(len(chunks))]
                    metadatas = [
                        {
                            "file_id": item_id,
                            "file_name": item_name,
                            "created_date": item.get(
                                "createdDateTime", "1970-01-01T00:00:00Z"
                            ),
                            "chunk_idx": i,
                        }
                        for i in range(len(chunks))
                    ]

                    document_collection.upsert(
                        ids=chunk_ids, documents=chunks, metadatas=metadatas
                    )
                    processed_count += 1
                    print(
                        f"  📄 Indexed/Updated File: {item_name} ({len(chunks)} text chunks)"
                    )
                else:
                    print(
                        f"  ⚠️ Failed to download '{item_name}' "
                        f"({file_res.status_code}): {file_res.text[:200]}"
                    )

        if "@odata.nextLink" in data:
            next_url = data["@odata.nextLink"]
        elif "@odata.deltaLink" in data:
            delta_link = data["@odata.deltaLink"]
            next_url = None

    if delta_link:
        with open(DELTA_TOKEN_FILE, "w") as f:
            json.dump({"delta_link": delta_link}, f)

    print(
        f"\n✅ Sync finished! (Indexed/Updated: {processed_count} files | Removed: {deleted_count} items)"
    )

    if processed_count == 0 and deleted_count == 0 and force_resync:
        print(
            "⚠️ WARNING: Full resync completed but indexed 0 files and 0 folders. "
            "This usually means the Graph token lacks consent for Files.Read.All, "
            "or the delta query is still malformed. Check the raw response above."
        )

    log_execution_time("OneDrive Delta Sync Total", sync_start)


# ==============================================================================
# 7. VERSION HISTORY: METADATA FETCH + CHART RENDERING (from script 1)
# ==============================================================================


def fetch_file_metadata(file_id: str, headers: dict) -> dict:
    """Confirms the file ID is valid/accessible and gets its current name."""
    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}"
    res = requests.get(url, headers=headers)

    if res.status_code != 200:
        raise RuntimeError(
            f"Could not access file with ID '{file_id}': {res.status_code} - {res.text}\n"
            "  → Double check the file ID, and that the signed-in account has access to it."
        )

    item = res.json()
    return {"id": item["id"], "name": item.get("name", file_id)}


def fetch_document_versions(file_id: str, headers: dict) -> list[dict]:
    """Retrieves the complete version history metadata for a OneDrive file via the Graph API."""
    versions_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/versions"
    res = requests.get(versions_url, headers=headers)

    if res.status_code != 200:
        raise RuntimeError(
            f"Graph API version history fetch failed: {res.status_code} - {res.text}"
        )

    raw_versions = res.json().get("value", [])
    parsed_versions = []
    for v in raw_versions:
        modified_by = (v.get("lastModifiedBy") or {}).get("user") or {}
        parsed_versions.append(
            {
                "version_id": v.get("id", "Unknown"),
                "modified_datetime": v.get("lastModifiedDateTime", "1970-01-01T00:00:00Z"),
                "modified_by": modified_by.get("displayName", "Unknown User"),
                "size_bytes": v.get("size", 0),
            }
        )

    # Graph API returns newest-first; sort oldest -> newest for chronological plotting
    parsed_versions.sort(key=lambda x: x["modified_datetime"])
    return parsed_versions


def render_version_history_graph(
    versions: list[dict], file_name: str, output_path: str = None
) -> str:
    """Renders the version-history size-over-time chart and saves it as a PNG."""
    if not versions:
        raise ValueError(
            f"'{file_name}' has no recorded version history (it may only "
            "have a single saved state, which OneDrive doesn't always version)."
        )

    if output_path is None:
        safe_name = file_name.replace(" ", "_").rsplit(".", 1)[0]
        output_path = f"version_history_{safe_name}.png"

    dates = [
        datetime.fromisoformat(v["modified_datetime"].replace("Z", "+00:00"))
        for v in versions
    ]
    sizes_kb = [v["size_bytes"] / 1024 for v in versions]
    labels = [f"v{v['version_id']}\n{v['modified_by']}" for v in versions]

    fig, ax = plt.subplots(figsize=(11, 6))
    ax.plot(
        dates,
        sizes_kb,
        marker="o",
        markersize=7,
        linewidth=2,
        color="#1a73e8",
        zorder=3,
    )
    ax.margins(x=0.08, y=0.25)

    for d, s, label in zip(dates, sizes_kb, labels):
        ax.annotate(
            label,
            (d, s),
            textcoords="offset points",
            xytext=(0, 12),
            ha="center",
            fontsize=8,
        )

    ax.set_title(f"Version History — {file_name}", fontsize=14, fontweight="bold")
    ax.set_xlabel("Modified Date")
    ax.set_ylabel("File Size (KB)")
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d %H:%M"))
    fig.autofmt_xdate(rotation=30)
    ax.grid(True, linestyle="--", alpha=0.4)
    fig.tight_layout()

    fig.savefig(output_path, dpi=150)
    plt.close(fig)

    return output_path


# ==============================================================================
# 8. TOOL IMPLEMENTATIONS
# ==============================================================================


def process_notes_with_gemini(raw_text: str) -> dict:
    """Uses Gemini JSON mode to extract structured meeting information."""
    print("🧠 Extracting structured meeting summary with Gemini JSON mode...")
    start_time = time.perf_counter()

    prompt = f"""
    Read the following meeting notes. Extract three things:
    1. The date of the meeting (or 'Date Not Specified' if absent).
    2. A brief summary of discussion points as a list of strings.
    3. Action items.

    Return the result EXCLUSIVELY as a valid JSON object with this exact structure:
    {{
      "meeting_date": "YYYY-MM-DD or formatted date string",
      "discussion_points": ["Point 1", "Point 2"],
      "action_items": [
        {{"task": "Task name", "owner": "Name or Team", "status": "Pending/In Progress/Done"}}
      ]
    }}

    Meeting Notes:
    {raw_text}
    """

    response = generate_content_with_retry(prompt)

    log_execution_time("Gemini JSON Extraction", start_time)
    return json.loads(response.text)


def generate_meeting_minutes(
    raw_notes_text: str,
    project_name: str,
    date: str = None,
    template_path: str = DEFAULT_MOM_TEMPLATE,
    output_path: str = "final_minutes.docx",
) -> dict:
    """Renders Word document using docxtpl and the 3-tier resolved MOM template."""
    start_time = time.perf_counter()

    resolved_template_path = fetch_template(template_path)

    meeting_data = process_notes_with_gemini(raw_notes_text)
    final_date = date if date else meeting_data.get("meeting_date", "Date Not Specified")

    print(f"📝 Rendering Word document with docxtpl using '{resolved_template_path}'...")
    doc = DocxTemplate(resolved_template_path)
    context = {
        "project_title": project_name,
        "meeting_date": final_date,
        "discussion_points": meeting_data.get("discussion_points", []),
        "action_items": meeting_data.get("action_items", []),
    }
    doc.render(context)
    doc.save(output_path)

    if COLAB_ENV:
        try:
            files.download(output_path)
            download_status = "Triggered browser file download."
        except Exception as e:
            download_status = f"Saved locally. Download note: {str(e)}"
    else:
        download_status = f"Saved locally at {output_path}"

    log_execution_time("Full MOM Generation Pipeline", start_time)
    return {
        "status": "MOM rendered successfully!",
        "file_path": output_path,
        "download_status": download_status,
        "extracted_summary": meeting_data,
    }


def find_folder_path(query: str) -> dict:
    """Finds folder path locations based on folder name or topic description."""
    start_time = time.perf_counter()
    results = folder_collection.query(query_texts=[query], n_results=3)

    log_execution_time("Folder Vector Search", start_time)

    matches = (
        [
            {"folder_name": m["folder_name"], "full_path": m["full_path"]}
            for m in results["metadatas"][0]
        ]
        if results.get("metadatas") and results["metadatas"][0]
        else []
    )
    return (
        {"matching_folders": matches}
        if matches
        else {"result": "No matching folder paths found."}
    )


def search_content_only(query: str) -> dict:
    """Searches strictly inside document body text using semantic similarity."""
    start_time = time.perf_counter()
    results = document_collection.query(query_texts=[query], n_results=5)

    log_execution_time("Content-Only Vector Search", start_time)

    if not results.get("metadatas") or not results["metadatas"][0]:
        return {"result": f"No content matches found for '{query}'."}

    matches = [
        {
            "file_name": meta["file_name"],
            "file_id": meta["file_id"],
            "text_snippet": results["documents"][0][i],
        }
        for i, meta in enumerate(results["metadatas"][0])
    ]
    return {"status": "Content-only search results", "matching_documents": matches}


def search_filename_and_content(query: str, find_latest_only: bool = False) -> dict:
    """Searches documents by evaluating BOTH filename relevance and text content matches."""
    start_time = time.perf_counter()
    results = document_collection.query(query_texts=[query], n_results=5)

    log_execution_time("ChromaDB Vector Search", start_time)

    if not results.get("metadatas") or not results["metadatas"][0]:
        return {"result": f"No documents found matching '{query}'."}

    query_keywords = [k.lower() for k in query.split() if len(k) > 2]
    files_map = {}

    for i, meta in enumerate(results["metadatas"][0]):
        f_id = meta["file_id"]
        file_name = meta.get("file_name", "").lower()

        filename_score = sum(3 for kw in query_keywords if kw in file_name)

        if f_id not in files_map:
            files_map[f_id] = {
                "file_id": f_id,
                "file_name": meta.get("file_name", "Unknown"),
                "created_date": meta.get("created_date", "1970-01-01T00:00:00Z"),
                "score": filename_score,
                "snippet": results["documents"][0][i],
            }

    if find_latest_only:
        sorted_files = sorted(
            files_map.values(), key=lambda x: (x["score"], x["created_date"]), reverse=True
        )
        return {"status": "Found target document", "latest_file": sorted_files[0]}
    else:
        sorted_files = sorted(files_map.values(), key=lambda x: x["score"], reverse=True)
        return {"status": "Search results", "matching_documents": sorted_files[:5]}


def read_file_content(file_id: str, file_name: str) -> str:
    """Downloads and reads complete file text from OneDrive."""
    start_time = time.perf_counter()
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    res = requests.get(
        f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content",
        headers=headers,
    )
    parsed_content = (
        ParserRegistry.parse(res.content, file_name)
        if res.status_code == 200
        else f"Download Error: {res.status_code}"
    )

    log_execution_time("OneDrive File Fetch & Parse", start_time)
    return parsed_content


def summarize_and_generate_mom(
    file_id: str, file_name: str, project_name: str = "Project Summary", date: str = None
) -> dict:
    """Reads OneDrive file content, extracts structured MOM via Gemini, and generates Word document."""
    raw_text = read_file_content(file_id, file_name)
    return generate_meeting_minutes(raw_notes_text=raw_text, project_name=project_name, date=date)


def generate_version_history_chart(file_id: str, file_name: str = None) -> dict:
    """Fetches a OneDrive file's full version history and renders a size-over-time
    line chart as a PNG. If file_name is not supplied, it is looked up from the
    file's current metadata. Returns the chart path plus the raw version list."""
    start_time = time.perf_counter()
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}

    if not file_name:
        meta = fetch_file_metadata(file_id, headers)
        file_name = meta["name"]

    versions = fetch_document_versions(file_id, headers)
    chart_path = render_version_history_graph(versions, file_name)

    log_execution_time("Version History Chart Generation", start_time)
    return {
        "status": "Version history chart rendered",
        "file_name": file_name,
        "version_count": len(versions),
        "chart_path": chart_path,
        "versions": versions,
    }


TOOLS_MAP = {
    "find_folder_path": find_folder_path,
    "search_content_only": search_content_only,
    "search_filename_and_content": search_filename_and_content,
    "read_file_content": read_file_content,
    "summarize_and_generate_mom": summarize_and_generate_mom,
    "generate_version_history_chart": generate_version_history_chart,
}

# ==============================================================================
# 9. REASONING & PROMPT ROUTING ARCHITECTURE
# ==============================================================================

SYSTEM_PROMPT = """
You are an onboarding and knowledge management assistant for Microsoft OneDrive.

STRICT REASONING & ROUTING ARCHITECTURE:

1. INTENT: USER REQUESTS TO "DOWNLOAD SUMMARY", "GENERATE MOM", "MAKE MOM", OR "SUMMARIZE AND DOWNLOAD" A FILE
   --> Step A: Execute `search_filename_and_content(query="<file_keyword>", find_latest_only=True)` ONCE to locate the target file_id and file_name.
   --> Step B: Immediately call `summarize_and_generate_mom(file_id, file_name, project_name)` using the retrieved metadata.
   --> DO NOT execute `search_content_only` or `find_folder_path`.
   --> DO NOT shotgun or trigger multiple search functions simultaneously.

2. INTENT: USER SEARCHES FOR A SPECIFIC FILE OR DOCUMENT BY NAME/TOPIC
   --> Execute ONLY `search_filename_and_content(query, find_latest_only=False)`.

3. INTENT: USER SEARCHES STRICTLY FOR IN-BODY TEXT/QUOTES WITHOUT KNOWING THE FILE NAME
   --> Execute ONLY `search_content_only(query)`.

4. INTENT: USER ASKS FOR FOLDER PATH LOCATIONS
   --> Execute ONLY `find_folder_path(query)`.

5. INTENT: USER ASKS FOR A FILE'S "VERSION HISTORY", "CHANGE LOG", "SIZE OVER TIME", OR TO "CHART HOW A FILE GREW/CHANGED"
   --> Step A: If the file_id is not already known, execute `search_filename_and_content(query="<file_keyword>", find_latest_only=True)` ONCE to locate it.
   --> Step B: Immediately call `generate_version_history_chart(file_id, file_name)` using the retrieved metadata.
   --> DO NOT call `summarize_and_generate_mom` for this intent.

CRITICAL RULE:
Always execute ONE tool per reasoning turn. Never invoke multiple search functions in a single batch.
"""


def send_with_429_retry(chat, payload, max_retries: int = 5, pacing_seconds: float = 1.0):
    """Sends requests with dynamic spacing and exponential backoff on 429 rate limit errors."""
    for attempt in range(max_retries):
        try:
            _enforce_global_pacing()
            time.sleep(pacing_seconds)
            return chat.send_message(payload)
        except APIError as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                wait_time = MIN_SECONDS_BETWEEN_CALLS * (attempt + 1) + random.uniform(0.5, 1.5)
                print(
                    f"\n⚠️ [RATE LIMIT 429] API throttled. Retrying in {wait_time:.1f}s (Attempt {attempt+1}/{max_retries})..."
                )
                time.sleep(wait_time)
            else:
                raise e
    raise RuntimeError("Failed after max retries due to persistent rate limits.")


def generate_content_with_retry(prompt: str, max_retries: int = 5, pacing_seconds: float = 1.0):
    """Same 429 retry/backoff + global pacing as send_with_429_retry, but for the
    direct client.models.generate_content() call used by process_notes_with_gemini."""
    for attempt in range(max_retries):
        try:
            _enforce_global_pacing()
            time.sleep(pacing_seconds)
            return client.models.generate_content(
                model=MODEL_NAME,
                contents=prompt,
                config=types.GenerateContentConfig(response_mime_type="application/json"),
            )
        except APIError as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                wait_time = MIN_SECONDS_BETWEEN_CALLS * (attempt + 1) + random.uniform(0.5, 1.5)
                print(
                    f"\n⚠️ [RATE LIMIT 429] API throttled. Retrying in {wait_time:.1f}s (Attempt {attempt+1}/{max_retries})..."
                )
                time.sleep(wait_time)
            else:
                raise e
    raise RuntimeError("Failed after max retries due to persistent rate limits.")


def send_message_and_handle_tools(chat, user_input: str):
    """Sends user message and handles function tool invocations sequentially.

    Capped at MAX_API_CALLS_PER_TURN total calls to the Gemini API (the initial
    message plus every subsequent tool-result round-trip) so a runaway tool loop
    can't generate unbounded requests in a single turn.
    """
    api_call_count = 1  # counts the initial send below
    response = send_with_429_retry(chat, user_input)

    while response.function_calls:
        if api_call_count >= MAX_API_CALLS_PER_TURN:
            print(
                f"\n⚠️ [LIMIT REACHED] Stopped after {api_call_count} API calls "
                f"(MAX_API_CALLS_PER_TURN={MAX_API_CALLS_PER_TURN}). "
                "Returning partial result instead of continuing the tool loop."
            )
            break

        for fn_call in response.function_calls:
            fn_name = fn_call.name
            fn_args = dict(fn_call.args)

            print(f"\n⚙️ [Executing Tool: {fn_name}]")

            if fn_name in TOOLS_MAP:
                try:
                    tool_result = TOOLS_MAP[fn_name](**fn_args)
                except Exception as e:
                    tool_result = {"error": str(e)}
            else:
                tool_result = {"error": f"Tool '{fn_name}' not found."}

            tool_part = types.Part.from_function_response(
                name=fn_name, response={"result": tool_result}
            )

            if api_call_count >= MAX_API_CALLS_PER_TURN:
                print(
                    f"\n⚠️ [LIMIT REACHED] Stopped after {api_call_count} API calls "
                    f"(MAX_API_CALLS_PER_TURN={MAX_API_CALLS_PER_TURN}) mid-batch."
                )
                break

            response = send_with_429_retry(chat, tool_part)
            api_call_count += 1

    if hasattr(response, "text") and response.text:
        return response.text
    else:
        try:
            parts = response.candidates[0].content.parts
            texts = [p.text for p in parts if hasattr(p, "text") and p.text]
            return "\n".join(texts) if texts else "Completed task."
        except Exception:
            return "Completed task."


def run_assistant():
    sync_onedrive_to_vector_db(force_resync=FORCE_RESYNC_ON_START)

    config = types.GenerateContentConfig(
        system_instruction=SYSTEM_PROMPT,
        tools=[
            find_folder_path,
            search_content_only,
            search_filename_and_content,
            read_file_content,
            summarize_and_generate_mom,
            generate_version_history_chart,
        ],
    )
    chat = client.chats.create(model=MODEL_NAME, config=config)

    print("\n==================================================")
    print(f"  OneDrive Knowledge, MOM & Version-History Assistant ({MODEL_NAME})")
    print("==================================================")
    print("Type 'exit' or 'quit' to end session.\n")

    while True:
        try:
            user_input = input("\nUSER > ").strip()
            if not user_input:
                continue
            if user_input.lower() in ["exit", "quit"]:
                print("Session ended. Goodbye!")
                break

            response_text = send_message_and_handle_tools(chat, user_input)
            print(f"\nAGENT > {response_text}")
            print("-" * 50)

        except KeyboardInterrupt:
            print("\nSession cancelled.")
            break
        except Exception as e:
            print(f"\n[ERROR] {str(e)}")


# ==============================================================================
# 10. STANDALONE VERSION-HISTORY MODE (lightweight, skips vector DB entirely)
# ==============================================================================


def run_version_history_standalone(file_id: str = None):
    """Runs just the version-history chart flow for a known OneDrive file ID,
    without touching the vector DB or Gemini at all. Mirrors the original
    standalone script's CLI behavior."""
    file_id = file_id or input("Enter the OneDrive file ID: ").strip()

    if not file_id:
        print("❌ No file ID provided. Exiting.")
        raise SystemExit(1)

    if "http" in file_id or "onedrive" in file_id.lower():
        print("\n❌ Error: It looks like you provided a URL. This script requires a Graph API Item ID.")
        print("An Item ID usually looks like '01ABCD...' or '0123456789ABCDEF!123'.")
        raise SystemExit(1)

    try:
        result = generate_version_history_chart(file_id)

        print(f"\n📋 Version history ({result['version_count']} versions) for '{result['file_name']}':")
        for v in result["versions"]:
            print(
                f"  v{v['version_id']:<5} | {v['modified_datetime']} | "
                f"{v['modified_by']:<20} | {v['size_bytes'] / 1024:.1f} KB"
            )

        print(f"\n✅ Done! Chart saved to: {result['chart_path']}")

    except (RuntimeError, ValueError, TimeoutError) as e:
        print(f"\n❌ {e}")


# ==============================================================================
# 11. ENTRY POINT
# ==============================================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="OneDrive Knowledge, MOM & Version-History Assistant"
    )
    parser.add_argument(
        "--version-history",
        nargs="?",
        const="",
        default=None,
        metavar="FILE_ID",
        help=(
            "Skip the interactive Gemini assistant and just render a version-history "
            "chart for the given OneDrive file ID. If FILE_ID is omitted you'll be "
            "prompted for it."
        ),
    )

    args, _unknown_args = parser.parse_known_args()

    if args.version_history is not None:
        run_version_history_standalone(args.version_history or None)
    else:
        run_assistant()