import io
import json
import os
import random
import time
from datetime import datetime

import pandas as pd
import requests
import streamlit as st

# Document parsers
from docx import Document
from docxtpl import DocxTemplate
from pptx import Presentation

# Charting
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# Vector Database & Embeddings
import chromadb
from chromadb.utils import embedding_functions

# Google GenAI SDK
from google import genai
from google.genai import types
from google.genai.errors import APIError

try:
    import pypdf
except ImportError:
    pypdf = None

# ==============================================================================
# 1. CONFIGURATION & STATE INITIALIZATION
# ==============================================================================
GEMINI_API_KEY = ""
CLIENT_ID = "2194a676-0580-4882-a21c-d7d4d0a8966c"
TENANT_ID = "common"
SCOPES = ["Files.Read.All", "User.Read"]
MODEL_NAME = "gemini-3.6-flash"

CHROMA_PATH = "./onedrive_vector_db"
DELTA_TOKEN_FILE = "./onedrive_delta_token.json"
DEFAULT_MOM_TEMPLATE = "mom_template.docx"

MAX_API_CALLS_PER_TURN = 3
MIN_SECONDS_BETWEEN_CALLS = 20.0
_last_gemini_call_time = 0.0

# Streamlit Page Config
st.set_page_config(page_title="OneDrive AI Assistant", page_icon="☁️", layout="centered")

# Persistent Session State
if "access_token" not in st.session_state:
    st.session_state.access_token = None
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "Hello! I am your OneDrive AI Assistant. How can I help you today?"}]

# Initialize GenAI Client
client = genai.Client(api_key=GEMINI_API_KEY)

# ==============================================================================
# 2. UTILITIES & RATE LIMITING
# ==============================================================================
def _enforce_global_pacing():
    global _last_gemini_call_time
    elapsed = time.perf_counter() - _last_gemini_call_time
    remaining = MIN_SECONDS_BETWEEN_CALLS - elapsed
    if remaining > 0:
        time.sleep(remaining)
    _last_gemini_call_time = time.perf_counter()

def log_execution_time(action_name: str, start_time: float):
    elapsed = time.perf_counter() - start_time
    print(f"⏱️ [{action_name}] completed in {elapsed:.3f} seconds.")

def acquire_microsoft_graph_token() -> str:
    if not st.session_state.access_token:
        raise RuntimeError("Not authenticated with Microsoft Graph.")
    return st.session_state.access_token

# ==============================================================================
# 3. DOCUMENT PARSING
# ==============================================================================
class ParserRegistry:
    @staticmethod
    def truncate_text(text: str, max_chars: int = 4000) -> str:
        return text[:max_chars] + "\n\n...[Content truncated]..." if len(text) > max_chars else text

    @classmethod
    def parse(cls, content_bytes: bytes, file_name: str) -> str:
        ext = os.path.splitext(file_name)[1].lower()
        try:
            if ext in [".xlsx", ".xls", ".csv"]:
                xls = pd.ExcelFile(io.BytesIO(content_bytes))
                return "\n\n".join([f"**Sheet: {s}**\n" + pd.read_excel(xls, sheet_name=s, nrows=20).to_markdown(index=False) for s in xls.sheet_names[:3]])
            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(io.BytesIO(content_bytes))
                slides = list(prs.slides)[:15]
                return cls.truncate_text("\n".join([f"Slide {i+1}: " + " | ".join([s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]) for i, slide in enumerate(slides)]))
            elif ext in [".docx", ".doc"]:
                doc = Document(io.BytesIO(content_bytes))
                return cls.truncate_text("\n".join([p.text for p in doc.paragraphs if p.text.strip()]))
            elif ext == ".pdf" and pypdf:
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                pages = list(reader.pages)[:10]
                return cls.truncate_text("\n".join([p.extract_text() for p in pages if p.extract_text()]))
            else:
                return cls.truncate_text(content_bytes.decode("utf-8", errors="ignore"))
        except Exception as e:
            return f"Parsing error for {file_name}: {str(e)}"

# ==============================================================================
# 4. VECTOR DB INITIALIZATION
# ==============================================================================
@st.cache_resource
def get_chroma_client():
    return chromadb.PersistentClient(path=CHROMA_PATH)

chroma_client = get_chroma_client()
embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
folder_collection = chroma_client.get_or_create_collection(name="folder_paths", embedding_function=embedding_fn)
document_collection = chroma_client.get_or_create_collection(name="document_chunks", embedding_function=embedding_fn)

# ==============================================================================
# 5. ONEDRIVE SYNC & TEMPLATE RESOLUTION
# ==============================================================================
def fetch_template(template_name: str = DEFAULT_MOM_TEMPLATE) -> str:
    if os.path.exists(template_name):
        return template_name
    try:
        results = document_collection.get(where={"file_name": template_name})
        if results and results.get("metadatas") and len(results["metadatas"]) > 0:
            file_id = results["metadatas"][0]["file_id"]
            token = acquire_microsoft_graph_token()
            res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content", headers={"Authorization": f"Bearer {token}"})
            if res.status_code == 200:
                with open(template_name, "wb") as f:
                    f.write(res.content)
                return template_name
    except Exception: pass
    
    try:
        token = acquire_microsoft_graph_token()
        res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{template_name}')", headers={"Authorization": f"Bearer {token}"})
        if res.status_code == 200:
            for item in res.json().get("value", []):
                if item.get("name", "").lower() == template_name.lower():
                    file_res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{item['id']}/content", headers={"Authorization": f"Bearer {token}"})
                    if file_res.status_code == 200:
                        with open(template_name, "wb") as f:
                            f.write(file_res.content)
                        return template_name
    except Exception: pass
    raise FileNotFoundError(f"❌ Template '{template_name}' could not be resolved.")

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks if chunks else [text]

def sync_onedrive_to_vector_db(force_resync: bool = False):
    global folder_collection, document_collection
    sync_start = time.perf_counter()
    
    if force_resync:
        try:
            chroma_client.delete_collection("folder_paths")
            chroma_client.delete_collection("document_chunks")
            if os.path.exists(DELTA_TOKEN_FILE): os.remove(DELTA_TOKEN_FILE)
        except Exception: pass
        folder_collection = chroma_client.get_or_create_collection(name="folder_paths", embedding_function=embedding_fn)
        document_collection = chroma_client.get_or_create_collection(name="document_chunks", embedding_function=embedding_fn)

    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    next_url = "https://graph.microsoft.com/v1.0/me/drive/root/delta?$select=id,name,createdDateTime,lastModifiedDateTime,folder,file,parentReference,deleted"
    
    if os.path.exists(DELTA_TOKEN_FILE) and not force_resync:
        with open(DELTA_TOKEN_FILE, "r") as f:
            next_url = json.load(f).get("delta_link")

    delta_link = None
    processed_count = 0
    deleted_count = 0

    while next_url:
        res = requests.get(next_url, headers=headers)
        if res.status_code == 410 or "resyncRequired" in res.text:
            if os.path.exists(DELTA_TOKEN_FILE): os.remove(DELTA_TOKEN_FILE)
            return sync_onedrive_to_vector_db(force_resync=True)
        if res.status_code != 200: raise RuntimeError(f"Graph Delta API Error: {res.text}")

        data = res.json()
        for item in data.get("value", []):
            item_id = item["id"]
            item_name = item.get("name", "Unknown")
            if "deleted" in item:
                try:
                    document_collection.delete(where={"file_id": item_id})
                    folder_collection.delete(ids=[item_id])
                    deleted_count += 1
                except Exception: pass
            elif "folder" in item:
                folder_collection.upsert(ids=[item_id], documents=[f"Folder Name: {item_name}"], metadatas=[{"folder_name": item_name, "full_path": item_name}])
            elif "file" in item:
                file_res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content", headers=headers)
                if file_res.status_code == 200:
                    chunks = chunk_text(ParserRegistry.parse(file_res.content, item_name))
                    chunk_ids = [f"{item_id}_chunk_{i}" for i in range(len(chunks))]
                    metadatas = [{"file_id": item_id, "file_name": item_name, "created_date": item.get("createdDateTime", "1970-01-01T00:00:00Z"), "chunk_idx": i} for i in range(len(chunks))]
                    document_collection.upsert(ids=chunk_ids, documents=chunks, metadatas=metadatas)
                    processed_count += 1

        if "@odata.nextLink" in data: next_url = data["@odata.nextLink"]
        elif "@odata.deltaLink" in data:
            delta_link = data["@odata.deltaLink"]
            next_url = None

    if delta_link:
        with open(DELTA_TOKEN_FILE, "w") as f: json.dump({"delta_link": delta_link}, f)
    log_execution_time("OneDrive Delta Sync", sync_start)
    return processed_count, deleted_count

# ==============================================================================
# 6. VERSION HISTORY & TOOL FUNCTIONS (With Docstrings for Gemini)
# ==============================================================================
def fetch_file_metadata(file_id: str, headers: dict) -> dict:
    """Confirms the file ID is valid/accessible and gets its current name."""
    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}", headers=headers)
    if res.status_code != 200: raise RuntimeError(f"Could not access file: {res.text}")
    return {"id": res.json()["id"], "name": res.json().get("name", file_id)}

def fetch_document_versions(file_id: str, headers: dict) -> list[dict]:
    """Retrieves the complete version history metadata for a OneDrive file."""
    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/versions", headers=headers)
    if res.status_code != 200: raise RuntimeError(f"Version history fetch failed: {res.text}")
    parsed = [{"version_id": v.get("id"), "modified_datetime": v.get("lastModifiedDateTime", "1970-01-01T00:00:00Z"), "modified_by": (v.get("lastModifiedBy") or {}).get("user", {}).get("displayName", "Unknown"), "size_bytes": v.get("size", 0)} for v in res.json().get("value", [])]
    parsed.sort(key=lambda x: x["modified_datetime"])
    return parsed

def render_version_history_graph(versions: list[dict], file_name: str, output_path: str = None) -> str:
    """Renders the version-history size-over-time chart and saves it as a PNG."""
    if not versions: raise ValueError("No version history.")
    if output_path is None: output_path = f"version_history_{file_name.replace(' ', '_').rsplit('.', 1)[0]}.png"
    dates = [datetime.fromisoformat(v["modified_datetime"].replace("Z", "+00:00")) for v in versions]
    sizes_kb = [v["size_bytes"] / 1024 for v in versions]
    labels = [f"v{v['version_id']}\n{v['modified_by']}" for v in versions]

    fig, ax = plt.subplots(figsize=(11, 6))
    ax.plot(dates, sizes_kb, marker="o", markersize=7, linewidth=2, color="#1a73e8", zorder=3)
    ax.margins(x=0.08, y=0.25)
    for d, s, label in zip(dates, sizes_kb, labels):
        ax.annotate(label, (d, s), textcoords="offset points", xytext=(0, 12), ha="center", fontsize=8)
    ax.set_title(f"Version History — {file_name}", fontsize=14, fontweight="bold")
    ax.set_xlabel("Modified Date"); ax.set_ylabel("File Size (KB)")
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d %H:%M"))
    fig.autofmt_xdate(rotation=30)
    ax.grid(True, linestyle="--", alpha=0.4)
    fig.tight_layout()
    fig.savefig(output_path, dpi=150)
    plt.close(fig)
    return output_path

def process_notes_with_gemini(raw_text: str) -> dict:
    """Uses Gemini JSON mode to extract structured meeting information."""
    prompt = f"Read the meeting notes. Extract: 1. Date 2. Summary 3. Action items. Return EXCLUSIVELY as JSON.\nNotes: {raw_text}"
    response = generate_content_with_retry(prompt)
    return json.loads(response.text)

def generate_meeting_minutes(raw_notes_text: str, project_name: str, date: str = None, template_path: str = DEFAULT_MOM_TEMPLATE, output_path: str = "final_minutes.docx") -> dict:
    """Renders Word document using docxtpl and the resolved MOM template."""
    resolved_template_path = fetch_template(template_path)
    meeting_data = process_notes_with_gemini(raw_notes_text)
    doc = DocxTemplate(resolved_template_path)
    doc.render({"project_title": project_name, "meeting_date": date or meeting_data.get("meeting_date", "Date Not Specified"), "discussion_points": meeting_data.get("discussion_points", []), "action_items": meeting_data.get("action_items", [])})
    doc.save(output_path)
    return {"status": "MOM rendered successfully!", "file_path": output_path, "extracted_summary": meeting_data}

def find_folder_path(query: str) -> dict:
    """Finds folder path locations based on folder name or topic description."""
    results = folder_collection.query(query_texts=[query], n_results=3)
    matches = [{"folder_name": m["folder_name"], "full_path": m["full_path"]} for m in results["metadatas"][0]] if results.get("metadatas") and results["metadatas"][0] else []
    return {"matching_folders": matches} if matches else {"result": "No matching folder paths found."}

def search_content_only(query: str) -> dict:
    """Searches strictly inside document body text using semantic similarity."""
    results = document_collection.query(query_texts=[query], n_results=5)
    if not results.get("metadatas") or not results["metadatas"][0]: return {"result": f"No content matches found."}
    return {"status": "Content search results", "matching_documents": [{"file_name": meta["file_name"], "file_id": meta["file_id"], "text_snippet": results["documents"][0][i]} for i, meta in enumerate(results["metadatas"][0])]}

def search_filename_and_content(query: str, find_latest_only: bool = False) -> dict:
    """Searches documents by evaluating BOTH filename relevance and text content matches."""
    results = document_collection.query(query_texts=[query], n_results=5)
    if not results.get("metadatas") or not results["metadatas"][0]: return {"result": f"No documents found."}
    query_keywords = [k.lower() for k in query.split() if len(k) > 2]
    files_map = {}
    for i, meta in enumerate(results["metadatas"][0]):
        f_id = meta["file_id"]
        if f_id not in files_map:
            files_map[f_id] = {"file_id": f_id, "file_name": meta.get("file_name", "Unknown"), "created_date": meta.get("created_date", "1970-01-01T00:00:00Z"), "score": sum(3 for kw in query_keywords if kw in meta.get("file_name", "").lower()), "snippet": results["documents"][0][i]}
    if find_latest_only:
        return {"status": "Found target document", "latest_file": sorted(files_map.values(), key=lambda x: (x["score"], x["created_date"]), reverse=True)[0]}
    return {"status": "Search results", "matching_documents": sorted(files_map.values(), key=lambda x: x["score"], reverse=True)[:5]}

def read_file_content(file_id: str, file_name: str) -> str:
    """Downloads and reads complete file text from OneDrive."""
    token = acquire_microsoft_graph_token()
    res = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content", headers={"Authorization": f"Bearer {token}"})
    return ParserRegistry.parse(res.content, file_name) if res.status_code == 200 else f"Download Error: {res.status_code}"

def summarize_and_generate_mom(file_id: str, file_name: str, project_name: str = "Project Summary", date: str = None) -> dict:
    """Reads OneDrive file content, extracts structured MOM via Gemini, and generates Word document."""
    raw_text = read_file_content(file_id, file_name)
    return generate_meeting_minutes(raw_notes_text=raw_text, project_name=project_name, date=date)

def generate_version_history_chart(file_id: str, file_name: str = None) -> dict:
    """Fetches a OneDrive file's full version history and renders a size-over-time line chart."""
    token = acquire_microsoft_graph_token()
    headers = {"Authorization": f"Bearer {token}"}
    if not file_name: file_name = fetch_file_metadata(file_id, headers)["name"]
    versions = fetch_document_versions(file_id, headers)
    chart_path = render_version_history_graph(versions, file_name)
    return {"status": "Version history chart rendered", "file_name": file_name, "chart_path": chart_path}

TOOLS_MAP = {
    "find_folder_path": find_folder_path,
    "search_content_only": search_content_only,
    "search_filename_and_content": search_filename_and_content,
    "read_file_content": read_file_content,
    "summarize_and_generate_mom": summarize_and_generate_mom,
    "generate_version_history_chart": generate_version_history_chart,
}

# ==============================================================================
# 7. AGENT ROUTING & MESSAGE HANDLING
# ==============================================================================
SYSTEM_PROMPT = """
You are an onboarding and knowledge management assistant for Microsoft OneDrive.

You have access to the following tools:

1. find_folder_path
   - Use when the user asks where a folder is located.

2. search_content_only
   - Use when the user asks about information contained inside documents.

3. search_filename_and_content
   - Use when the user asks to find a specific document/file.

4. read_file_content
   - Use when the complete contents of a known file are required.

5. summarize_and_generate_mom
   - Use when the user asks to generate Meeting Minutes / MOM.

6. generate_version_history_chart
   - Use when the user asks for file version history or a version history chart.

ROUTING RULES:

- GENERATE MOM:
  First use search_filename_and_content with find_latest_only=True.
  Then use summarize_and_generate_mom with the file_id and file_name returned.

- SEARCH SPECIFIC FILE:
  Use search_filename_and_content.

- SEARCH DOCUMENT CONTENT:
  Use search_content_only.

- FIND FOLDER:
  Use find_folder_path.

- VERSION HISTORY:
  First use search_filename_and_content with find_latest_only=True.
  Then use generate_version_history_chart.

IMPORTANT:

- Execute only one tool at a time.
- After receiving a tool result, analyze it and decide whether another tool is required.
- Do not stop after receiving the first tool result if the user's request requires another operation.
- After all required tools have completed, provide a clear final response to the user.
"""

def send_with_429_retry(
    chat,
    payload,
    max_retries: int = 5,
    pacing_seconds: float = 0.5
):
    for attempt in range(max_retries):

        try:
            _enforce_global_pacing()

            if pacing_seconds > 0:
                time.sleep(pacing_seconds)

            return chat.send_message(payload)

        except APIError as e:

            error_text = str(e)

            if "429" in error_text or "RESOURCE_EXHAUSTED" in error_text:

                wait_time = (
                    MIN_SECONDS_BETWEEN_CALLS * (attempt + 1)
                    + random.uniform(0.5, 1.5)
                )

                print(
                    f"⚠️ Rate limited. "
                    f"Retrying in {wait_time:.1f} seconds..."
                )

                time.sleep(wait_time)

            else:
                raise

    raise RuntimeError(
        "Gemini request failed after maximum retry attempts."
    )

def generate_content_with_retry(prompt: str, max_retries: int = 5, pacing_seconds: float = 1.0):
    for attempt in range(max_retries):
        try:
            _enforce_global_pacing()
            time.sleep(pacing_seconds)
            return client.models.generate_content(model=MODEL_NAME, contents=prompt, config=types.GenerateContentConfig(response_mime_type="application/json"))
        except APIError as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                time.sleep(MIN_SECONDS_BETWEEN_CALLS * (attempt + 1) + random.uniform(0.5, 1.5))
            else: raise e
    raise RuntimeError("Failed after max retries due to rate limits.")

def send_message_and_handle_tools(chat, user_input: str):
    """
    Send a user message to Gemini and execute function calls sequentially.

    Important:
    - One tool is executed at a time.
    - Every tool result is immediately sent back to Gemini.
    - Gemini gets a chance to decide the next step.
    - The loop continues until Gemini produces a final text response.
    """

    try:
        # First user message
        response = send_with_429_retry(chat, user_input)

        tool_call_count = 0

        while True:

            # ---------------------------------------------------------
            # 1. Check whether Gemini wants to call a tool
            # ---------------------------------------------------------
            function_calls = response.function_calls

            if not function_calls:
                break

            # ---------------------------------------------------------
            # 2. Execute ONLY ONE function call
            # ---------------------------------------------------------
            fn_call = function_calls[0]

            fn_name = fn_call.name
            fn_args = dict(fn_call.args or {})

            print(f"🔧 Gemini requested tool: {fn_name}")
            print(f"📦 Arguments: {fn_args}")

            if fn_name not in TOOLS_MAP:
                tool_result = {
                    "error": f"Tool '{fn_name}' not found."
                }
            else:
                try:
                    tool_result = TOOLS_MAP[fn_name](**fn_args)

                    print(f"✅ Tool completed: {fn_name}")
                    print(f"📤 Result: {tool_result}")

                except Exception as e:
                    tool_result = {
                        "error": str(e)
                    }

                    print(f"❌ Tool failed: {fn_name}")
                    print(f"Error: {e}")

            # ---------------------------------------------------------
            # 3. Send the tool result back to Gemini
            # ---------------------------------------------------------
            tool_response = types.Part.from_function_response(
                name=fn_name,
                response={
                    "result": tool_result
                }
            )

            response = send_with_429_retry(
                chat,
                [tool_response]
            )

            tool_call_count += 1

            # Safety limit
            if tool_call_count >= 10:
                return (
                    "I reached the maximum number of tool operations "
                    "allowed for this request."
                )

        # -------------------------------------------------------------
        # 4. Gemini has finished processing
        # -------------------------------------------------------------

        if hasattr(response, "text") and response.text:
            return response.text

        # Fallback extraction
        try:
            texts = []

            for candidate in response.candidates:
                if not candidate.content:
                    continue

                for part in candidate.content.parts:
                    if hasattr(part, "text") and part.text:
                        texts.append(part.text)

            if texts:
                return "\n".join(texts)

        except Exception as e:
            print(f"⚠️ Could not extract response text: {e}")

        return "I completed the request, but Gemini did not return a text response."

    except Exception as e:
        print(f"❌ Gemini processing error: {e}")
        raise
# Setup GenAI Chat Session in Streamlit State
if "chat_session" not in st.session_state:
    config = types.GenerateContentConfig(
        system_instruction=SYSTEM_PROMPT,
        tools=[find_folder_path, search_content_only, search_filename_and_content, read_file_content, summarize_and_generate_mom, generate_version_history_chart]
    )
    st.session_state.chat_session = client.chats.create(model=MODEL_NAME, config=config)

# ==============================================================================
# 8. STREAMLIT UI & AUTHENTICATION
# ==============================================================================
def authenticate_microsoft():
    if st.session_state.access_token: return True
    st.warning("🔑 Microsoft Authentication Required")
    
    device_code_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/devicecode"
    res = requests.post(device_code_url, data={"client_id": CLIENT_ID, "scope": " ".join(SCOPES)})
    data = res.json()

    st.markdown(f"**1. Open this link:** [{data.get('verification_uri')}]({data.get('verification_uri')})")
    st.markdown(f"**2. Enter Code:** `{data['user_code']}`")
    
    with st.spinner("Waiting for you to complete sign-in on your browser..."):
        token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
        token_payload = {"grant_type": "urn:ietf:params:oauth:grant-type:device_code", "client_id": CLIENT_ID, "device_code": data["device_code"]}
        start_time = time.time()
        
        while time.time() - start_time < data.get("expires_in", 900):
            time.sleep(data.get("interval", 5))
            token_res = requests.post(token_url, data=token_payload).json()
            if "access_token" in token_res:
                st.session_state.access_token = token_res["access_token"]
                st.rerun()
            if token_res.get("error") not in ["authorization_pending"]: break
    return False

# Build UI
st.title("☁️ OneDrive AI Assistant")

if authenticate_microsoft():
    with st.sidebar:
        st.header("Settings")
        
        # 🔴 THE FIX: THIS BUTTON PURGES THE BROKEN CACHE 🔴
        if st.button("Reset Chat Session & Clear Cache"):
            st.session_state.clear()
            st.rerun()
            
        st.divider()
        if st.button("Sync OneDrive Database"):
            with st.spinner("Syncing your files to local Vector DB..."):
                processed, deleted = sync_onedrive_to_vector_db(force_resync=False)
                st.success(f"Sync complete! Processed: {processed}, Deleted: {deleted}")
                
        if st.button("Force Full Resync"):
            with st.spinner("Clearing DB and running full resync..."):
                processed, deleted = sync_onedrive_to_vector_db(force_resync=True)
                st.success(f"Full resync complete! Processed: {processed}, Deleted: {deleted}")

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if prompt := st.chat_input("Ask about your OneDrive files..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            with st.spinner("Executing request (this may take a moment due to API pacing)..."):
                try:
                    response_text = send_message_and_handle_tools(st.session_state.chat_session, prompt)
                    st.markdown(response_text)
                    st.session_state.messages.append({"role": "assistant", "content": response_text})
                except Exception as e:
                    error_msg = f"An error occurred: {str(e)}"
                    st.error(error_msg)
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})